"""
gen_pptx_full.py
================
Generate NAS_PINN_Thermal_Quenching.pptx
Professional 16-slide presentation with architecture diagram,
heat maps, full k-sweep tables, and subframe results.
"""

import sys
from pathlib import Path
from thermal_pinn.runtime_paths import REPORT_DIR, RESULT_DIR

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: pip install python-pptx"); sys.exit(1)

ROOT   = Path(__file__).resolve().parent
RES    = RESULT_DIR
FIG    = ROOT / "reports" / "thesis" / "figures"
KF     = RES / "01_all_k_fields" / "bayesian"
SFDIR  = RES / "subframe_result_application"
OUT    = REPORT_DIR / "NAS_PINN_Thermal_Quenching.pptx"

NAVY  = RGBColor(0x1C, 0x35, 0x57)
BLUE  = RGBColor(0x1A, 0x5C, 0x96)
CYAN  = RGBColor(0x00, 0xBC, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF0, 0xF4, 0xF8)
MGRAY = RGBColor(0xCC, 0xD1, 0xD9)
DGRAY = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED   = RGBColor(0xC6, 0x28, 0x28)
GOLD  = RGBColor(0xF0, 0xA5, 0x00)
ROSE  = RGBColor(0xBE, 0x18, 0x5D)

SW = Inches(13.33)
SH = Inches(7.5)


# ── primitives ────────────────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill):
    from pptx.util import Inches as I
    sh = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    return sh


def txt(slide, text, x, y, w, h, fs=13, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    color = color or DGRAY
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(fs)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tf


def img(slide, path, x, y, w=None, h=None):
    p = Path(path)
    if not p.exists():
        print(f"  [WARN] missing: {p.name}")
        return
    kw = {}
    if w: kw['width']  = Inches(w)
    if h: kw['height'] = Inches(h)
    slide.shapes.add_picture(str(p), Inches(x), Inches(y), **kw)


def stripe(slide, color=NAVY, h=0.65):
    rect(slide, 0, 0, 13.33, h, color)


def header(slide, title, sub=None):
    rect(slide, 0, 0, 13.33, 7.5, LGRAY)
    stripe(slide, NAVY, 0.65)
    txt(slide, title, 0.4, 0.07, 12.5, 0.55, fs=22, bold=True,
        color=WHITE, align=PP_ALIGN.LEFT)
    if sub:
        txt(slide, sub, 0.4, 0.65, 12.5, 0.35, fs=12,
            color=MGRAY, align=PP_ALIGN.LEFT)


def footer(slide, text):
    rect(slide, 0, 7.15, 13.33, 0.35, NAVY)
    txt(slide, text, 0.4, 7.18, 12.5, 0.3, fs=11,
        color=WHITE, align=PP_ALIGN.CENTER)


def bullets(slide, items, x, y, w, h, fs=12, color=None):
    color = color or DGRAY
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = "  \u2022  " + item
        p.font.size = Pt(fs)
        p.font.color.rgb = color
        p.space_before = Pt(4)
        first = False


def draw_table(slide, headers, rows, x, y, col_widths,
               row_h=0.38, hdr_bg=None, alt_bg=None, fs=10):
    hdr_bg = hdr_bg or NAVY
    alt_bg = alt_bg or LGRAY
    cx = [x]
    for w in col_widths[:-1]:
        cx.append(cx[-1] + w)
    # header
    for j, (h, cxj, wj) in enumerate(zip(headers, cx, col_widths)):
        rect(slide, cxj, y, wj, row_h, hdr_bg)
        txt(slide, h, cxj + 0.03, y + 0.03, wj - 0.06, row_h - 0.06,
            fs=fs, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # rows
    for i, row in enumerate(rows):
        ry = y + row_h * (i + 1)
        bg = alt_bg if i % 2 == 0 else WHITE
        for j, (val, cxj, wj) in enumerate(zip(row, cx, col_widths)):
            rect(slide, cxj, ry, wj, row_h, bg)
            fc = GREEN if (str(val).startswith('\u2605') or
                           (j == len(row)-1 and '%' in str(val) and str(val) != '0%')) else DGRAY
            txt(slide, str(val), cxj + 0.03, ry + 0.03, wj - 0.06, row_h - 0.06,
                fs=fs, color=fc, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════

def s01_title(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, NAVY)
    rect(sl, 0, 5.5, 13.33, 0.07, CYAN)
    txt(sl, "NAS-PINN Framework for Thermal Quenching",
        0.6, 1.0, 12.1, 1.0, fs=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "of A356 Aluminium Castings",
        0.6, 2.0, 12.1, 0.75, fs=30, color=CYAN, align=PP_ALIGN.CENTER)
    txt(sl, "Multi-Step Window Prediction  |  Neural Architecture Search  |  FEM Cost Reduction",
        0.6, 3.1, 12.1, 0.45, fs=15, color=MGRAY, align=PP_ALIGN.CENTER)
    txt(sl, "Reference: Mortensen et al. (2026) Int. J. Adv. Manuf. Technol.  doi:10.1007/s00170-026-17515-w",
        0.6, 4.2, 12.1, 0.35, fs=11, color=MGRAY, align=PP_ALIGN.CENTER, italic=True)
    txt(sl, "April 2026", 0.6, 5.65, 12.1, 0.35, fs=12, color=MGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MOTIVATION
# ══════════════════════════════════════════════════════════════════════════

def s02_motivation(prs):
    sl = blank(prs)
    header(sl, "Motivation", "Why reduce FEM calls in aluminium quenching?")

    rect(sl, 0.3, 0.85, 6.0, 5.9, WHITE)
    txt(sl, "The Problem", 0.5, 0.9, 5.6, 0.38, fs=14, bold=True, color=RED)
    bullets(sl, [
        "Water quenching of A356: 540 C to 20 C in 30 s",
        "FEM needs 60 Crank-Nicolson steps per simulation",
        "Each FEM step: full mesh assembly + linear solve",
        "Repeated evaluations needed for design optimisation",
        "Industrial subframe geometry: prohibitively expensive",
        "Reference: Mortensen et al. (2026)",
    ], 0.5, 1.35, 5.6, 5.0, fs=12)

    rect(sl, 6.6, 0.85, 6.4, 5.9, WHITE)
    txt(sl, "Our Solution", 6.8, 0.9, 6.0, 0.38, fs=14, bold=True, color=GREEN)
    bullets(sl, [
        "PINN surrogate covers k FEM intervals per window",
        "k=5 reduces 20 FEM anchor calls down to 4",
        "NAS finds the best architecture automatically",
        "IC-consistent output: exact IC at each window start",
        "Two protocols: FEM-anchored hybrid + PINN-only",
        "Tested on 7 canonical + 1 subframe domain",
    ], 6.8, 1.35, 6.0, 5.0, fs=12)

    rect(sl, 6.4, 0.85, 0.05, 5.9, MGRAY)
    footer(sl, "Goal: reduce FEM anchor calls by up to 80% while keeping mean MAE below 10 C")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PHYSICAL PROBLEM + DOMAINS
# ══════════════════════════════════════════════════════════════════════════

def s03_problem(prs):
    sl = blank(prs)
    header(sl, "Physical Problem and Domains Tested",
           "Transient heat equation — A356 aluminium, 7 canonical + 1 subframe geometry")

    rect(sl, 0.3, 0.85, 6.0, 3.5, WHITE)
    txt(sl, "Governing Equation", 0.5, 0.9, 5.6, 0.38, fs=13, bold=True, color=NAVY)
    txt(sl, "\u03c1c\u209a \u2202T/\u2202t = \u2207\u00b7(\u03bb\u2207T)",
        0.5, 1.35, 5.6, 0.55, fs=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(sl, "BC:  -\u03bb \u2207T\u00b7n = h(T - T\u221e)    (Robin cooling)",
        0.5, 2.0, 5.6, 0.38, fs=12, color=DGRAY, align=PP_ALIGN.CENTER)
    txt(sl, "IC:  T(x,0) = 540 \u00b0C    |    T\u221e = 20 \u00b0C",
        0.5, 2.45, 5.6, 0.38, fs=12, color=DGRAY, align=PP_ALIGN.CENTER)
    bullets(sl, [
        "\u03bb = 150 W/mK (2D)  |  160 W/mK (3D)",
        "h = 5000 W/m\u00b2K (2D)  |  4000 W/m\u00b2K (3D)",
        "\u0394t_int = 0.5 s (60 steps)  |  \u0394t_anc = 1.5 s (20 anchors)",
    ], 0.5, 2.9, 5.6, 1.3, fs=11)

    rect(sl, 0.3, 4.55, 6.0, 2.2, WHITE)
    txt(sl, "Three NAS Architectures", 0.5, 4.6, 5.6, 0.35, fs=13, bold=True, color=NAVY)
    draw_table(sl,
        ["Strategy", "Layers", "Neurons", "Activation"],
        [("Bayesian/TPE", "5", "151", "ReLU"),
         ("NSGA-II",      "3", "153", "tanh"),
         ("NSGA-III",     "3", "75",  "tanh")],
        x=0.35, y=5.0, col_widths=[1.6, 0.8, 0.9, 1.5], row_h=0.36, fs=10)

    rect(sl, 6.6, 0.85, 6.4, 5.9, WHITE)
    txt(sl, "8 Domains Tested", 6.8, 0.9, 6.0, 0.38, fs=13, bold=True, color=NAVY)
    draw_table(sl,
        ["Domain", "Dim", "Challenge"],
        [("Rectangle",         "2D", "Regular baseline"),
         ("Circle",            "2D", "Curved boundary"),
         ("L-shape",           "2D", "Re-entrant corner"),
         ("Rect. prism",       "3D", "Regular 3D"),
         ("Cylinder",          "3D", "Curved 3D surface"),
         ("Stacked cubes",     "3D", "Stepped gradients"),
         ("L-shape prism",     "3D", "3D re-entrant edge"),
         ("Subframe surrogate","3D", "Multi-region, hetero.")],
        x=6.65, y=1.3, col_widths=[1.9, 0.55, 3.7], row_h=0.36, fs=10)

    footer(sl, "Heat transfer coefficient h differs between 2D (5000) and 3D (4000) W/m\u00b2K benchmark groups")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — NAS-PINN ARCHITECTURE DIAGRAM
# ══════════════════════════════════════════════════════════════════════════

def s04_architecture(prs):
    sl = blank(prs)
    header(sl, "NAS-PINN Architecture",
           "IC-consistent output, Fourier embedding, self-adaptive loss weights")

    rect(sl, 0.3, 0.85, 12.73, 5.75, WHITE)
    arch_fig = FIG / "pinn_framework_fig_pinn_framework.png"
    img(sl, arch_fig, 0.35, 0.9, w=12.63)

    footer(sl, "Output formula: T\u0302(x,\u03c4) = \u03b8_ic(x) + \u03c4\u00b7N\u03b8(x,\u03c4,\u03b8_ic)  "
           "\u2014  exact IC satisfaction at \u03c4=0 by construction, no penalty weight needed")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MSWP TIME STRUCTURE
# ══════════════════════════════════════════════════════════════════════════

def s05_timestep(prs):
    sl = blank(prs)
    header(sl, "Multi-Step Window Prediction (MSWP)",
           "FEM-anchored hybrid sweep + PINN-only adaptive rollout")

    rect(sl, 0.3, 0.85, 12.73, 5.75, WHITE)
    img(sl, FIG / "fem_pinn_timestep_fig_fem_pinn_timestep.png", 0.35, 0.9, w=12.63)

    footer(sl, "Bottom lane (pink): PINN-only rollout \u2014 FEM used only at t=0, "
           "all subsequent ICs from previous PINN prediction")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — k-SKIP TABLE
# ══════════════════════════════════════════════════════════════════════════

def s06_kskip(prs):
    sl = blank(prs)
    header(sl, "k-Skip Window Parameters",
           "How the skip factor controls FEM anchor cost")

    rect(sl, 0.3, 0.85, 5.8, 5.85, WHITE)
    txt(sl, "Window parameters per k", 0.5, 0.9, 5.4, 0.38, fs=13, bold=True, color=NAVY)
    draw_table(sl,
        ["k", "\u0394t window", "PINN windows", "FEM anchors", "FEM saved"],
        [("1 (base)", "1.5 s", "20", "20", "0%"),
         ("2",        "3.0 s", "10", "10", "50%"),
         ("3",        "4.5 s",  "7",  "7", "67%"),
         ("4",        "6.0 s",  "5",  "5", "75%"),
         ("5",        "7.5 s",  "4",  "4", "80%")],
        x=0.4, y=1.35, col_widths=[0.9, 1.1, 1.3, 1.2, 1.1], row_h=0.42, fs=11)

    rect(sl, 6.4, 0.85, 6.6, 5.85, WHITE)
    txt(sl, "IC-consistent output layer", 6.6, 0.9, 6.2, 0.38, fs=13, bold=True, color=NAVY)
    txt(sl, "T\u0302(x,\u03c4) = \u03b8_ic(x) + \u03c4 \u00b7 N\u03b8(x, \u03c4, \u03b8_ic)",
        6.6, 1.4, 6.2, 0.55, fs=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(sl, "\u03c4 \u2208 [0,1]  \u2014  normalised time inside window",
        6.6, 2.05, 6.2, 0.35, fs=11, color=DGRAY, align=PP_ALIGN.CENTER)

    txt(sl, "Loss function", 6.6, 2.65, 6.2, 0.38, fs=13, bold=True, color=NAVY)
    txt(sl, "L = W_pde\u00b7L_pde + W_ic\u00b7L_ic + W_bc\u00b7L_bc + W_end\u00b7L_end",
        6.6, 3.1, 6.2, 0.55, fs=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(sl, "PINN-only: W_end\u00b7L_end replaced by surface cooling prior",
        6.6, 3.72, 6.2, 0.35, fs=11, color=ROSE, align=PP_ALIGN.CENTER)

    txt(sl, "Fourier embedding", 6.6, 4.35, 6.2, 0.38, fs=13, bold=True, color=NAVY)
    txt(sl, "\u03b3(v) = [sin(2\u03c0Bv), cos(2\u03c0Bv)]   F=64   \u03c3=1.0/1.5",
        6.6, 4.8, 6.2, 0.38, fs=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(sl, "Mitigates spectral bias near steep boundary gradients",
        6.6, 5.25, 6.2, 0.35, fs=11, color=DGRAY, align=PP_ALIGN.CENTER)

    footer(sl, "Larger k = fewer FEM anchors required; IC-consistent formula guarantees "
           "exact initial condition at every window start")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 2D RESULTS TABLE (best per k)
# ══════════════════════════════════════════════════════════════════════════

def s07_2d_table(prs):
    sl = blank(prs)
    header(sl, "2D Domain Results \u2014 Best Result per k",
           "Lowest MAE across all domains and architectures for each k value")

    rect(sl, 0.3, 0.85, 12.73, 2.75, WHITE)
    txt(sl, "Best result for each k  (all 2D domains and architectures)",
        0.5, 0.9, 12.3, 0.38, fs=13, bold=True, color=NAVY)
    draw_table(sl,
        ["k", "Best domain", "Best arch", "MAE (C)", "FEM saved", "FEM calls"],
        [("1 (baseline)", "L-shape",  "Bayesian/TPE", "1.93",         "0% (baseline)", "20"),
         ("2",            "L-shape",  "Bayesian/TPE", "\u2605 2.48",  "50%",           "10"),
         ("3",            "L-shape",  "Bayesian/TPE", "\u2605 3.25",  "67%",           "7"),
         ("4",            "L-shape",  "Bayesian/TPE", "\u2605 4.10",  "75%",           "5"),
         ("5",            "L-shape",  "Bayesian/TPE", "\u2605 5.23",  "80%",           "4")],
        x=0.4, y=1.35, col_widths=[1.3, 1.6, 1.7, 1.2, 1.5, 1.2], row_h=0.38, fs=11)

    rect(sl, 0.3, 3.75, 12.73, 2.9, WHITE)
    txt(sl, "Full breakdown by domain (Bayesian/TPE, mean MAE in C)",
        0.5, 3.8, 12.3, 0.38, fs=12, bold=True, color=NAVY)
    draw_table(sl,
        ["Domain / Arch", "k=1 (0%)", "k=2 (50%)", "k=3 (67%)", "k=4 (75%)", "k=5 (80%)"],
        [("Rectangle / Bayesian", "2.49", "3.49", "4.34", "5.35", "5.80"),
         ("Circle / Bayesian",    "2.70", "2.70", "3.63", "4.30", "5.70"),
         ("Circle / NSGA-III",    "1.64", "3.70", "5.66", "9.68", "12.50"),
         ("L-shape / Bayesian",   "1.93", "2.48", "3.25", "4.10", "5.23"),
         ("L-shape / NSGA-II",    "2.96", "4.96", "7.59","11.74","17.24")],
        x=0.4, y=4.25, col_widths=[2.4, 1.4, 1.4, 1.4, 1.4, 1.4], row_h=0.38, fs=10)

    footer(sl, "Bayesian/TPE dominates 2D across all k  \u2014  L-shape best domain due to IC-consistent layer handling the re-entrant corner")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 2D HEAT MAPS (L-shape k=1 and k=5)
# ══════════════════════════════════════════════════════════════════════════

def s08_2d_heatmaps(prs):
    sl = blank(prs)
    header(sl, "2D L-shape Heat Maps \u2014 Bayesian/TPE",
           "FEM reference | PINN prediction | absolute error  (left: k=1, right: k=5)")

    rect(sl, 0.2, 0.85, 6.3, 5.9, WHITE)
    txt(sl, "k=1  |  MAE = 1.93 C  |  FEM calls: 20 (baseline)",
        0.4, 0.9, 6.0, 0.38, fs=12, bold=True, color=NAVY)
    img(sl, KF / "lshape" / "fig_th1_2d_lshape_bayesian_k1.png", 0.25, 1.3, w=6.2)

    rect(sl, 6.8, 0.85, 6.3, 5.9, WHITE)
    txt(sl, "k=5  |  MAE = 5.23 C  |  FEM calls: 4  (80% saved)",
        7.0, 0.9, 6.0, 0.38, fs=12, bold=True, color=GREEN)
    img(sl, KF / "lshape" / "fig_th1_2d_lshape_bayesian_k5.png", 6.85, 1.3, w=6.2)

    footer(sl, "Error stays local to re-entrant corner even at k=5  \u2014  no global spread of error observed")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 3D RESULTS TABLE (best per k)
# ══════════════════════════════════════════════════════════════════════════

def s09_3d_table(prs):
    sl = blank(prs)
    header(sl, "3D Domain Results \u2014 Best Result per k",
           "Lowest MAE across all 3D domains and architectures for each k value")

    rect(sl, 0.3, 0.85, 12.73, 2.75, WHITE)
    txt(sl, "Best result for each k  (all 3D domains and architectures)",
        0.5, 0.9, 12.3, 0.38, fs=13, bold=True, color=NAVY)
    draw_table(sl,
        ["k", "Best domain", "Best arch", "MAE (C)", "FEM saved", "FEM calls"],
        [("1 (baseline)", "L-shape 3D", "NSGA-II",      "3.43",         "0% (baseline)", "20"),
         ("2",            "L-shape 3D", "NSGA-II/III",  "\u2605 5.97",  "50%",           "10"),
         ("3",            "L-shape 3D", "Bayesian/TPE", "\u2605 7.08",  "67%",           "7"),
         ("4",            "L-shape 3D", "Bayesian/TPE", "\u2605 8.73",  "75%",           "5"),
         ("5",            "L-shape 3D", "Bayesian/TPE", "\u2605 9.93",  "80%",           "4")],
        x=0.4, y=1.35, col_widths=[1.3, 1.7, 1.7, 1.2, 1.5, 1.2], row_h=0.38, fs=11)

    rect(sl, 0.3, 3.75, 12.73, 2.9, WHITE)
    txt(sl, "Full breakdown (mean MAE in C)",
        0.5, 3.8, 12.3, 0.38, fs=12, bold=True, color=NAVY)
    draw_table(sl,
        ["Domain / Arch", "k=1 (0%)", "k=2 (50%)", "k=3 (67%)", "k=4 (75%)", "k=5 (80%)"],
        [("Rect. prism / Bayesian","8.01","8.80","9.77","11.13","13.01"),
         ("Cylinder / Bayesian",   "6.43","6.64","7.98","9.48", "11.90"),
         ("Stacked / Bayesian",    "7.58","9.25","11.19","11.37","12.26"),
         ("L-shape 3D / Bayesian", "8.41","7.18","7.08","8.73", "9.93"),
         ("L-shape 3D / NSGA-II",  "3.43","5.97","13.12","9.28","27.05")],
        x=0.4, y=4.25, col_widths=[2.4, 1.4, 1.4, 1.4, 1.4, 1.4], row_h=0.38, fs=10)

    footer(sl, "L-shape 3D / NSGA-II best at k=1 (3.43 C)  |  Bayesian most stable at k=3..5  |  "
           "k=2-3 practical sweet spot for 3D step-skipping")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 3D HEAT MAPS (Cylinder k=1 and k=5)
# ══════════════════════════════════════════════════════════════════════════

def s10_3d_heatmaps(prs):
    sl = blank(prs)
    header(sl, "3D Cylinder Heat Maps \u2014 Bayesian/TPE",
           "FEM reference | PINN prediction | absolute error  (left: k=1, right: k=5)")

    rect(sl, 0.2, 0.85, 6.3, 5.9, WHITE)
    txt(sl, "k=1  |  MAE = 6.43 C  |  FEM calls: 20 (baseline)",
        0.4, 0.9, 6.0, 0.38, fs=12, bold=True, color=NAVY)
    img(sl, KF / "fig_th1_3d_cylinder_bayesian_k1.png", 0.25, 1.3, w=6.2)

    rect(sl, 6.8, 0.85, 6.3, 5.9, WHITE)
    txt(sl, "k=5  |  MAE = 11.90 C  |  FEM calls: 4  (80% saved)",
        7.0, 0.9, 6.0, 0.38, fs=12, bold=True, color=GREEN)
    img(sl, KF / "fig_th1_3d_cylinder_bayesian_k5.png", 6.85, 1.3, w=6.2)

    footer(sl, "Main error concentrates at top/bottom edges where boundary curvature is highest  \u2014  "
           "global cooling pattern well preserved at k=5")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SUBFRAME DOMAIN OVERVIEW
# ══════════════════════════════════════════════════════════════════════════

def s11_subframe_domain(prs):
    sl = blank(prs)
    header(sl, "Subframe-like Surrogate Domain",
           "Academic 3D surrogate: external cooling body + support pads + cavity regions")

    # full-width image (3D geometry + 2D top-view layout side by side)
    rect(sl, 0.2, 0.85, 12.93, 4.05, WHITE)
    img(sl, SFDIR / "subframe_reference_domain_3d_2d_.png", 0.25, 0.88, w=12.83)

    # feature bullets below the image
    rect(sl, 0.2, 5.05, 6.4, 1.65, WHITE)
    txt(sl, "Domain features", 0.35, 5.08, 6.0, 0.35, fs=12, bold=True, color=NAVY)
    bullets(sl, [
        "External cooling body (teal) — main quenched surface",
        "Support / contact pads (dark red) — localised transitions",
        "Simplified cavity regions (orange) — enclosed features",
        "Paper-inspired simplified 3D surrogate for PINN/FEM comparison",
    ], 0.35, 5.45, 6.0, 1.2, fs=11)

    rect(sl, 6.9, 5.05, 6.2, 1.65, WHITE)
    txt(sl, "Two evaluation protocols", 7.05, 5.08, 5.8, 0.35, fs=12, bold=True, color=NAVY)
    bullets(sl, [
        "FEM-anchored hybrid: FEM at every window start, k=1..5 sweep",
        "PINN-only rollout: FEM only at t=0, autonomous advance",
        "All 3 NAS architectures  |  Regional MAE: external / support / cavity",
    ], 7.05, 5.45, 5.8, 1.2, fs=11)

    footer(sl, "Left panel: 3D perspective view  |  Right panel: 2D top-view layout  \u2014 "
           "heterogeneous multi-region geometry tests framework scalability")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SUBFRAME HYBRID RESULTS
# ══════════════════════════════════════════════════════════════════════════

def s12_hybrid(prs):
    sl = blank(prs)
    header(sl, "Subframe \u2014 FEM-Anchored Hybrid Sweep",
           "Fixed k=1..5 sweep  |  FEM anchor at every window start")

    rect(sl, 0.3, 0.85, 5.8, 5.85, WHITE)
    txt(sl, "Results by architecture (recommended k)", 0.5, 0.9, 5.4, 0.35, fs=12, bold=True, color=NAVY)
    draw_table(sl,
        ["Architecture", "Rec. k", "Mean MAE", "FEM saved", "Run (min)"],
        [("Bayesian/TPE", "4", "5.56 C", "75%",        "2.6"),
         ("NSGA-II",      "1", "5.03 C", "0% (k=1)",   "10.3"),
         ("\u2605 NSGA-III","5","7.51 C", "80%",        "2.3")],
        x=0.4, y=1.35, col_widths=[1.7, 0.7, 1.2, 1.2, 0.75], row_h=0.42, fs=11)

    txt(sl, "NSGA-III k=5 field at t=30 s", 0.5, 2.75, 5.4, 0.35, fs=12, bold=True, color=NAVY)
    img(sl, SFDIR / "subframe_hybrid_nsga3_k5_field_application.png", 0.4, 3.15, w=5.6)

    rect(sl, 6.4, 0.85, 6.6, 5.85, WHITE)
    txt(sl, "k-sweep curve + cavity slice", 6.6, 0.9, 6.2, 0.35, fs=12, bold=True, color=NAVY)
    img(sl, FIG / "subframe_hybrid_k_sweep_paper.png", 6.5, 1.3, w=6.3)
    img(sl, SFDIR / "subframe_hybrid_nsga3_k5_cavity_slice.png", 6.5, 4.0, w=6.3)

    footer(sl, "Best hybrid result: NSGA-III k=5  \u2014  80% FEM saving at 7.51 C mean MAE  |  "
           "Bayesian k=4: 75% saving at 5.56 C (safer option)")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SUBFRAME PINN-ONLY RESULTS
# ══════════════════════════════════════════════════════════════════════════

def s13_pinn_only(prs):
    sl = blank(prs)
    header(sl, "Subframe \u2014 PINN-Only Adaptive Rollout",
           "FEM used only at t=0  |  all subsequent ICs from previous PINN prediction")

    rect(sl, 0.3, 0.85, 5.8, 5.85, WHITE)
    txt(sl, "Rollout results", 0.5, 0.9, 5.4, 0.35, fs=12, bold=True, color=NAVY)
    draw_table(sl,
        ["Architecture", "Final MAE", "Rel. L2", "FEM saving"],
        [("Bayesian/TPE",     "30.57 C", "0.246", "35%"),
         ("\u2605 NSGA-II",   "25.08 C", "0.192", "35%"),
         ("NSGA-III",         "27.58 C", "0.209", "35%")],
        x=0.4, y=1.35, col_widths=[1.7, 1.2, 0.95, 1.0], row_h=0.42, fs=11)

    txt(sl, "Regional MAE \u2014 NSGA-II", 0.5, 2.75, 5.4, 0.35, fs=12, bold=True, color=NAVY)
    draw_table(sl,
        ["Region", "MAE", "Bias"],
        [("External body", "22.88 C", "+2.95 C"),
         ("Support pads",  "58.17 C", "-53.34 C"),
         ("Cavity",        "36.76 C", "-14.08 C")],
        x=0.4, y=3.22, col_widths=[1.8, 1.2, 1.0], row_h=0.40, fs=11)

    txt(sl, "k schedule (all architectures):", 0.5, 4.3, 5.4, 0.35, fs=11, bold=True, color=NAVY)
    txt(sl, "k = 1-1-1-1-1-1-1-2-2-2-3-3-1",
        0.5, 4.7, 5.4, 0.38, fs=14, bold=True, color=ROSE, align=PP_ALIGN.CENTER)
    txt(sl, "Small k during fast early cooling, wider windows later",
        0.5, 5.15, 5.4, 0.35, fs=11, color=DGRAY, italic=True)
    img(sl, FIG / "subframe_pinn_only_adaptive_k_schedule.png", 0.4, 5.55, w=5.6)

    rect(sl, 6.4, 0.85, 6.6, 5.85, WHITE)
    txt(sl, "PINN-only NSGA-II field + cavity slice", 6.6, 0.9, 6.2, 0.35, fs=12, bold=True, color=NAVY)
    img(sl, SFDIR / "subframe_pinn_only_nsga2_field_application.png", 6.5, 1.3, w=6.3)
    img(sl, SFDIR / "subframe_pinn_only_nsga2_cavity_slice.png", 6.5, 4.0, w=6.3)

    footer(sl, "PINN-only: 1 FEM solve total  \u2014  NSGA-II achieves 25.08 C final MAE  |  "
           "support/cavity regions dominate the error budget")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — PINN-ONLY vs HYBRID COMPARISON
# ══════════════════════════════════════════════════════════════════════════

def s14_comparison(prs):
    sl = blank(prs)
    header(sl, "PINN-Only vs. FEM-Anchored Hybrid \u2014 Side by Side",
           "Both protocols on the same subframe-like surrogate domain")

    rect(sl, 0.2, 0.85, 12.93, 5.75, WHITE)
    img(sl, FIG / "subframe_v8_pinn_only_vs_hybrid.png", 0.3, 0.9, w=12.73)

    footer(sl, "Top rows: PINN-only (FEM at t=0 only)  |  Bottom rows: FEM-anchored hybrid  |  "
           "Left: temperature field at t=30 s  |  Right: absolute error map")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — PROTOCOL SUMMARY TABLE
# ══════════════════════════════════════════════════════════════════════════

def s15_summary_table(prs):
    sl = blank(prs)
    header(sl, "Protocol Comparison \u2014 All Key Results",
           "Canonical domains + subframe domain, both protocols")

    rect(sl, 0.2, 0.85, 12.93, 5.9, WHITE)
    txt(sl, "Summary of best results", 0.4, 0.9, 12.5, 0.38, fs=13, bold=True, color=NAVY)
    draw_table(sl,
        ["Protocol", "Domain", "Architecture", "k", "MAE (C)", "FEM saved"],
        [("2D canonical",     "L-shape",     "Bayesian/TPE", "1",  "1.93",         "0% (baseline)"),
         ("2D canonical",     "L-shape",     "Bayesian/TPE", "3",  "\u2605 3.25",  "67%"),
         ("2D canonical",     "L-shape",     "Bayesian/TPE", "5",  "5.23",         "80%"),
         ("3D canonical",     "L-shape 3D",  "NSGA-II",      "1",  "3.43",         "0% (baseline)"),
         ("3D canonical",     "Cylinder",    "Bayesian/TPE", "3",  "\u2605 7.98",  "67%"),
         ("3D canonical",     "L-shape 3D",  "Bayesian/TPE", "5",  "9.93",         "80%"),
         ("Subframe hybrid",  "Subframe",    "Bayesian/TPE", "4",  "5.56 (mean)",  "75%"),
         ("Subframe hybrid",  "Subframe",    "NSGA-III",     "5",  "\u2605 7.51 (mean)", "80%"),
         ("PINN-only rollout","Subframe",    "NSGA-II",      "adp","25.08 (final)","35%  (1 FEM solve)")],
        x=0.3, y=1.35, col_widths=[2.0, 1.5, 1.7, 0.6, 2.0, 2.3], row_h=0.40, fs=10)

    footer(sl, "Subframe hybrid: up to 80% FEM saving at <10 C mean error  |  "
           "PINN-only: fully autonomous, 25 C MAE with single FEM solve at t=0")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONCLUSIONS
# ══════════════════════════════════════════════════════════════════════════

def s16_conclusions(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.33, 7.5, NAVY)
    rect(sl, 0, 5.6, 13.33, 0.07, CYAN)

    txt(sl, "Conclusions", 0.6, 0.25, 12.1, 0.65, fs=32, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

    conc = [
        "NAS-PINN MSWP reduces FEM anchor calls by 50-80% on all 8 domains tested; "
        "Bayesian/TPE most consistent for canonical 2D/3D, NSGA-II/III better for subframe",
        "Best 2D: L-shape/Bayesian k=5 reaches 5.23 C at 80% FEM saving; "
        "k=3 gives 3.25 C at 67% saving (practical sweet spot)",
        "Best 3D: L-shape 3D/NSGA-II at k=1 gives 3.43 C; Bayesian/TPE cylinder "
        "at k=3 gives 7.98 C with 67% FEM saving",
        "Subframe hybrid: NSGA-III k=5 achieves 7.51 C mean MAE with 80% FEM saving; "
        "Bayesian k=4 gives 5.56 C at 75% saving",
        "PINN-only rollout: NSGA-II reaches 25.08 C final MAE using FEM only at t=0; "
        "k schedule adapts automatically to physical cooling stages",
        "IC-consistent output layer and Fourier embeddings are essential components; "
        "support/cavity regions remain dominant error sources in autonomous rollout",
    ]

    tb = sl.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.1), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for c in conc:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = "\u2714  " + c
        p.font.size = Pt(12.5)
        p.font.color.rgb = WHITE
        p.space_before = Pt(5)
        first = False

    txt(sl, "References", 0.6, 5.35, 12.1, 0.35, fs=13, bold=True, color=CYAN)
    txt(sl,
        "[1] Mortensen et al. (2026) Int.J.Adv.Manuf.Technol. doi:10.1007/s00170-026-17515-w  "
        "[2] Raissi et al. (2019) J.Comput.Phys.378  "
        "[3] Tancik et al. (2020) NeurIPS  "
        "[4] McClenny & Braga-Neto (2023) J.Comput.Phys.474  "
        "[5] Deb et al. (2002) IEEE TEC  "
        "[6] Bergstra et al. (2011) NeurIPS",
        0.6, 5.75, 12.1, 0.75, fs=10, color=MGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════

def main():
    print("[gen_pptx_full] Generating NAS_PINN_Thermal_Quenching.pptx ...")
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    steps = [
        ("01 Title",                  s01_title),
        ("02 Motivation",             s02_motivation),
        ("03 Physical Problem",       s03_problem),
        ("04 Architecture",           s04_architecture),
        ("05 MSWP Timestep",          s05_timestep),
        ("06 k-Skip Framework",       s06_kskip),
        ("07 2D Results Table",       s07_2d_table),
        ("08 2D Heat Maps",           s08_2d_heatmaps),
        ("09 3D Results Table",       s09_3d_table),
        ("10 3D Heat Maps",           s10_3d_heatmaps),
        ("11 Subframe Domain",        s11_subframe_domain),
        ("12 Subframe Hybrid",        s12_hybrid),
        ("13 Subframe PINN-Only",     s13_pinn_only),
        ("14 Comparison",             s14_comparison),
        ("15 Summary Table",          s15_summary_table),
        ("16 Conclusions",            s16_conclusions),
    ]

    for name, fn in steps:
        print(f"  slide {name} ...")
        fn(prs)

    prs.save(str(OUT))
    print(f"\n  Saved: {OUT}")
    print("Done.")


if __name__ == "__main__":
    main()
