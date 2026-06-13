#!/usr/bin/env python3
"""
NAS-PINN Defense PPTX v4 — 30 slides
Title · Outline · Intro · Theory(6) · ProbStmt · Method(8) · Results(7) · Disc(2) · Concl(3) · Thanks
Run:  python3 make_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIGURES = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Figures")
OUTPUT  = os.path.join(os.path.dirname(os.path.abspath(__file__)), "defense_final.pptx")

C_PRI  = RGBColor(26,  74,  92)
C_ACC  = RGBColor(31, 115, 158)
C_LTBG = RGBColor(245, 248, 250)
C_TG   = RGBColor(80,  80,  80)
C_RUL  = RGBColor(200, 215, 222)
C_WHI  = RGBColor(255, 255, 255)
C_DRK  = RGBColor(20,  20,  20)
C_LBLU = RGBColor(220, 235, 242)
C_GRN  = RGBColor(34, 139,  34)
C_RED  = RGBColor(160,  30,  30)
C_AMB  = RGBColor(150,  90,   0)
C_YLW  = RGBColor(255, 248, 220)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 29
_fn = [0]

def ns(): return prs.slides.add_slide(BLANK)

def rect(sl, l, t, w, h, fill=None, line=None, lw=Pt(0.5)):
    sh = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else:    sh.line.fill.background()

def _p(tf, i, text, sz, bold=False, col=C_DRK, align=PP_ALIGN.LEFT, italic=False, spb=0):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = align
    if spb: p.space_before = Pt(spb)
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = col

def txt(sl, l, t, w, h, lines, wrap=True):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    for i, ln in enumerate(lines):
        _p(tf, i, ln[0], ln[1], ln[2], ln[3], ln[4],
           ln[5] if len(ln)>5 else False, ln[6] if len(ln)>6 else 0)

def t1(sl, text, l, t, w, h, sz=14, bold=False, col=C_DRK, align=PP_ALIGN.LEFT, italic=False):
    txt(sl, l, t, w, h, [(text, sz, bold, col, align, italic)])

def img(sl, fname, l, t, w, h, sub=""):
    p = os.path.join(FIGURES, sub, fname) if sub else os.path.join(FIGURES, fname)
    if os.path.exists(p):
        sl.shapes.add_picture(p, Inches(l), Inches(t), Inches(w), Inches(h))

def notes(sl, en, tr=""):
    sl.notes_slide.notes_text_frame.text = (
        f"{en}\n\n--- TURKCE ---\n{tr}" if tr else en)

def header(sl, title, sub=None):
    rect(sl, 0, 0, 13.33, 0.93, fill=C_PRI)
    t1(sl, title, 0.35, 0.07, 12.5, 0.55, sz=22, bold=True, col=C_WHI)
    if sub: t1(sl, sub, 0.35, 0.61, 12.5, 0.31, sz=10.5, col=C_RUL)
    rect(sl, 0, 0.93, 13.33, 0.04, fill=C_ACC)

def footer(sl):
    _fn[0] += 1
    rect(sl, 0, 7.17, 13.33, 0.33, fill=RGBColor(240,243,245))
    t1(sl, "Omer Cetinkaya  |  NAS-PINN Thesis Defense  |  UiA June 2026",
       0.3, 7.19, 11.3, 0.28, sz=8.5, col=C_TG)
    t1(sl, f"{_fn[0]} / {TOTAL}", 12.1, 7.19, 1.0, 0.28, sz=8.5, col=C_TG, align=PP_ALIGN.RIGHT)

def secdiv(sl, num, title, sub):
    rect(sl, 0, 0, 13.33, 7.5, fill=C_PRI)
    rect(sl, 0, 3.4, 13.33, 0.06, fill=C_ACC)
    t1(sl, f"{num:02d}", 0.7, 1.5, 2.0, 1.8, sz=80, bold=True, col=RGBColor(50,100,120))
    t1(sl, title, 2.9, 2.45, 9.5, 1.2, sz=40, bold=True, col=C_WHI)
    t1(sl, sub,   2.9, 3.75, 9.5, 0.7, sz=15, col=C_RUL)

def card(sl, l, t, w, h, title, items, hbg=C_ACC, hfg=C_WHI, bbg=C_LTBG, tsz=13, fsz=11.5):
    rect(sl, l, t, w, 0.38, fill=hbg)
    rect(sl, l, t+0.38, w, h-0.38, fill=bbg)
    t1(sl, title, l+0.1, t+0.04, w-0.2, 0.33, sz=tsz, bold=True, col=hfg)
    lines = [(b, fsz, False, C_DRK, PP_ALIGN.LEFT, False, 3 if i else 0) for i,b in enumerate(items)]
    txt(sl, l+0.12, t+0.44, w-0.24, h-0.55, lines)

def keybx(sl, l, t, w, h, text, sz=12, bg=C_YLW, border=C_AMB):
    rect(sl, l, t, w, h, fill=bg, line=border)
    t1(sl, text, l+0.12, t+0.08, w-0.24, h-0.16, sz=sz, italic=True, col=C_DRK)

def stbl(sl, l, t, w, h, hdrs, rows, hbg=C_PRI, hfg=C_WHI, sz=10):
    nr,nc = len(rows)+1,len(hdrs)
    tbl = sl.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    def _c(cell, text, bold=False, fg=C_DRK, bg=C_WHI, align=PP_ALIGN.CENTER):
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=fg
    for j,h2 in enumerate(hdrs): _c(tbl.cell(0,j), h2, bold=True, fg=hfg, bg=hbg)
    for i,row in enumerate(rows):
        bg = C_LTBG if i%2==0 else C_WHI
        for j,v in enumerate(row):
            fg  = C_GRN if "+" in v and "C" in v and v[0].isdigit() else C_GRN if "check" in v.lower() else C_RED if "x" in v.lower() and len(v)<5 else C_DRK
            fg  = C_GRN if "✓" in v else C_RED if "✗" in v else C_DRK
            bld = (j==0 or "✓" in v or "✗" in v)
            _c(tbl.cell(i+1,j), v, bold=bld, fg=fg, bg=bg)

def htbl(sl, l, t, w, h, hdrs, rows, sz=9.5):
    nr,nc = len(rows)+1,len(hdrs)
    tbl = sl.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(h)).table
    def _c(cell, text, bold=False, fg=C_DRK, bg=C_WHI, align=PP_ALIGN.CENTER):
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=fg
    for j,h2 in enumerate(hdrs): _c(tbl.cell(0,j), h2, bold=True, fg=C_WHI, bg=C_PRI)
    for i,row in enumerate(rows):
        for j,v in enumerate(row):
            TICK,CROSS = "✓","✗"
            if CROSS in v:   bg2,fg,b = RGBColor(255,220,220), C_RED, True
            elif TICK in v:  bg2,fg,b = RGBColor(210,230,245), RGBColor(31,100,176), True
            elif j==0:       bg2,fg,b = C_LTBG, C_PRI, True
            else:            bg2,fg,b = C_WHI, C_DRK, False
            _c(tbl.cell(i+1,j), v, bold=b, fg=fg, bg=bg2)

# ---- SLIDES ----------------------------------------------------------------

# S01 TITLE
sl=ns()
rect(sl,0,0,13.33,7.5,fill=C_PRI)
rect(sl,0,1.85,13.33,3.85,fill=C_WHI)
rect(sl,0,1.85,13.33,0.05,fill=C_ACC); rect(sl,0,5.65,13.33,0.05,fill=C_ACC)
t1(sl,"MASTERS THESIS DEFENSE",0,0.62,13.33,0.5,sz=13,col=C_RUL,align=PP_ALIGN.CENTER)
t1(sl,"NAS-Guided FEM-Anchored PINNs",0.5,2.05,12.33,0.8,sz=32,bold=True,col=C_PRI,align=PP_ALIGN.CENTER)
t1(sl,"for Transient Heat Transfer Simulation",0.5,2.83,12.33,0.52,sz=22,col=C_ACC,align=PP_ALIGN.CENTER)
txt(sl,0.5,3.42,12.33,0.85,[("Multi-Step Window Prediction for Industrial Water Quenching of A356 Aluminium Automotive Subframes",13,False,C_TG,PP_ALIGN.CENTER,True)])
txt(sl,0.5,4.35,12.33,0.95,[("Omer Cetinkaya",18,True,C_PRI,PP_ALIGN.CENTER),("Supervisor: Prof. Turgay Celik   |   University of Agder (UiA), Norway   |   June 2026",11,False,C_TG,PP_ALIGN.CENTER,False,5)])
rect(sl,0,6.8,13.33,0.7,fill=C_ACC)
t1(sl,"Department of Engineering Sciences   |   Faculty of Engineering and Science",0,6.89,13.33,0.5,sz=11,col=C_WHI,align=PP_ALIGN.CENTER)
notes(sl,"Welcome. Hybrid FEM-PINN reduces solver calls 65-85% while keeping engineering accuracy.","Hos geldiniz. Hibrit FEM-PINN yaklasimi cagrilari yüzde 65-85 azaltiyor.")

# S02 OUTLINE
sl=ns(); header(sl,"Outline"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,1.0,1.12,11.33,0.72,fill=C_ACC)
t1(sl,"Key Question: How many FEM solver calls can be replaced while keeping temperature error within an engineering accuracy band?",1.2,1.18,11.0,0.6,sz=12.5,bold=True,col=C_WHI,align=PP_ALIGN.CENTER)
secs=[("01","Introduction","Motivation  Problem"),("02","Theory","Heat Eq  FEM  PINN  NAS\nBayesian  NSGA-II/III"),("03","Methodology","FEM to NAS-PINN to Framework\nTraining  MSWP"),("04","Results","2D  3D  Thermal Fin\nMAE  L2  Heatmaps"),("05","Conclusion","Summary  Future work")]
bw,gap=2.35,0.18; sx=(13.33-5*bw-4*gap)/2
for i,(n,ttl,sub) in enumerate(secs):
    x=sx+i*(bw+gap)
    rect(sl,x,2.05,bw,4.75,fill=C_LTBG,line=C_RUL); rect(sl,x,2.05,bw,0.52,fill=C_PRI)
    t1(sl,n,x+0.1,2.09,0.55,0.44,sz=20,bold=True,col=C_RUL)
    t1(sl,ttl,x+0.1,2.62,bw-0.2,0.45,sz=13,bold=True,col=C_PRI)
    t1(sl,sub,x+0.1,3.12,bw-0.2,3.5,sz=10.5,col=C_TG)
notes(sl,"Five sections including Theory.","Bes bolum, Theory dahil.")

# S03 INTRODUCTION
sl=ns(); header(sl,"Introduction: Why This Problem Matters","A356 aluminium subframes  Water quenching  FEM cost"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
for i,(ttl,buls) in enumerate([("The Process",["A356 aluminium subframes heated to 540 C","then rapidly cooled in water (quenching).","","Non-uniform cooling causes thermal gradients","and mechanical distortion (up to 3 mm).","","Mortensen et al. 2026: rack design","reduces distortion from 3 mm to 1 mm."]),("The Challenge",["Full thermo-mechanical FEM (StaMiSim):","  3M finite elements  27 domains","  2 min per FEM time step","","Design optimisation:","100 runs x 40 min = 67 hours","Too slow for routine design use."]),("This Thesis",["FEM-anchored NAS-PINN hybrid:","  FEM anchors at selected time steps","  NAS-PINN predicts between anchors","  Skip factor k controls FEM budget","","Goal:","65-85% fewer FEM calls","within engineering accuracy band"])]):
    card(sl,0.3+i*4.22,1.12,4.0,5.75,ttl,buls,fsz=12.5,tsz=14)
notes(sl,"A356 quenching costs 67 hours for 100-design study. We replace most FEM steps with PINN.","A356 quenching 100 tasarim icin 67 saat. FEM adimlarini PINN ile degistiriyoruz.")

# S04 SECTION THEORY
sl=ns()
secdiv(sl,2,"Theoretical Background","Heat Equation  FEM  PINN  NAS  Bayesian/TPE  NSGA-II/III")
notes(sl,"Theory covers all building blocks.","Theory tum yapi taslarini kapsiyor.")

# S05 HEAT EQ + FEM
sl=ns(); header(sl,"Theory: Heat Equation and FEM Reference Solver","Governing physics  Robin BC  Crank-Nicolson  Our solver vs Mortensen FEM"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,0.3,1.12,6.0,5.75,fill=C_LTBG,line=C_RUL)
txt(sl,0.45,1.22,5.7,5.5,[("Governing equation:",14,True,C_PRI,PP_ALIGN.LEFT),("rho cp dT/dt = nabla.(kT nabla T)",16,False,C_DRK,PP_ALIGN.CENTER,False,7),("rho=density  cp=specific heat  kT=conductivity",10.5,False,C_TG,PP_ALIGN.CENTER),("",4,False,C_DRK,PP_ALIGN.LEFT),("Initial condition:",12,True,C_ACC,PP_ALIGN.LEFT,False,4),("T(x,0) = T0 = 540 C  (uniform)",13,False,C_DRK,PP_ALIGN.CENTER,False,4),("",4,False,C_DRK,PP_ALIGN.LEFT),("Robin boundary condition:",12,True,C_ACC,PP_ALIGN.LEFT,False,4),("-kT nabla T.n = h(T - Tw)",14,False,C_DRK,PP_ALIGN.CENTER,False,4),("h=HTC  Tw=20C  n=outward normal",10.5,False,C_TG,PP_ALIGN.CENTER,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Dimensionless: theta=(T-Tw)/(T0-Tw)",12,True,C_PRI,PP_ALIGN.LEFT,False,4),("d theta/dt = alpha nabla^2 theta",14,False,C_DRK,PP_ALIGN.CENTER,False,4),("theta=1 at start  theta=0 at equilibrium",10.5,False,C_TG,PP_ALIGN.CENTER,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Biot Bi=hL/kT >> 1:",12,True,C_PRI,PP_ALIGN.LEFT,False,3),("Internal conduction is bottleneck.",11,False,C_DRK,PP_ALIGN.LEFT,False,2),("Steep gradients = hard prediction.",11,True,C_RED,PP_ALIGN.LEFT,False,2)])
rect(sl,6.6,1.12,6.4,5.75,fill=C_LTBG,line=C_RUL)
txt(sl,6.75,1.22,6.1,5.5,[("FEM: how it works:",14,True,C_PRI,PP_ALIGN.LEFT),("Domain to elements to linear system per step:",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("(C + dt K) T_{n+1} = C T_n + dt f",13,False,C_DRK,PP_ALIGN.CENTER,False,4),("C=capacity  K=conductance  f=boundary",10.5,False,C_TG,PP_ALIGN.CENTER,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Our Crank-Nicolson FD solver:",12,True,C_PRI,PP_ALIGN.LEFT,False,4),("  dt_int=0.5s  unconditionally stable",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("  Snapshot every 1.5s  21 total",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("  Each snapshot = one FEM call",11.5,True,C_ACC,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("NOT Mortensen FEM!",12,True,C_RED,PP_ALIGN.LEFT,False,4),("Our solver: thermal only, structured grid.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("Fast. Used as training ground truth.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Mortensen FEM (StaMiSim):",12,True,C_RED,PP_ALIGN.LEFT,False,4),("Full thermo-mechanical. 3M elements.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("27 coupled domains. 2 min per call.",11.5,True,C_RED,PP_ALIGN.LEFT,False,2)])
notes(sl,"Left: governing equation, Robin BC, dimensionless form. Bi>>1: internal conduction bottleneck. Right: FEM, Crank-Nicolson stable. Our FD solver is NOT Mortensen FEM — ours is fast thermal-only, Mortensen is full commercial.","Sol: fizik denklemleri. Sag: FEM. Bizim FD cozucumuz Mortensen FEM degil — bizimki hizli termal, Mortensen tam ticari.")

# S06 PINN
sl=ns(); header(sl,"Theory: Physics-Informed Neural Networks (PINNs)","Loss function  IC-consistent output  Window decomposition  Fourier  Adaptive weights"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
card(sl,0.3,1.12,4.0,2.85,"What is a PINN?",["Neural network trained to satisfy PDE,","not just fit data.","","L = wpde Lpde + wic Lic + wbc Lbc","","Lpde = MSE(dN/dt - alpha nablaSquared N)","via automatic differentiation.","L=0 means network IS exact PDE solution."],fsz=11.5,tsz=13)
card(sl,4.55,1.12,4.0,2.85,"IC-Consistent Output Form",["T_hat(x,t) = T_ic(x) + tau N_phi(x,t)","","At tau=0: T_hat = T_ic EXACTLY","for any network output.","IC is HARD-CODED. wic L_ic = 0 always.","Network learns only the CHANGE from IC.","Easier task."],fsz=11.5,tsz=13)
card(sl,8.8,1.12,4.23,2.85,"Window Decomposition",["Long time  gradient pathology:","IC gradient (t=0) vs PDE gradient (all t)","have different magnitudes.","Optimizer ignores small-gradient terms.","","FIX: one PINN per short window.","FEM anchor resets error each window.","Drift bounded."],fsz=11.5,tsz=13)
rect(sl,0.3,4.12,12.73,0.9,fill=C_YLW,line=C_AMB)
txt(sl,0.45,4.18,12.5,0.78,[("Self-Adaptive Weights: wbc and wend are LEARNABLE (updated by optimizer). wpde and wic are FIXED.",11.5,True,C_DRK,PP_ALIGN.LEFT),("Endpoint loss L_end: penalises drift at t=t_end of window (temporal end, NOT spatial boundary).",11,False,C_DRK,PP_ALIGN.LEFT,False,3)])
rect(sl,0.3,5.18,12.73,1.65,fill=C_LTBG,line=C_RUL)
txt(sl,0.45,5.24,12.3,1.52,[("Fourier Feature Embedding  gamma(x) = [cos(2pi B x), sin(2pi B x)]  B~N(0,sigma2):",12,True,C_ACC,PP_ALIGN.LEFT),("Neural networks learn LOW-frequency patterns first (spectral bias). Steep gradients at fin-tip/corners are HIGH-frequency  plain PINN underfits these regions.",11,False,C_DRK,PP_ALIGN.LEFT,False,3),("Fourier maps coordinates through random sinusoids  network sees high-freq content from start. BUT: applied ONLY for Bayesian/TPE on Thermal Fin. Compact NSGA nets: too small  Fourier hurts them.",11,False,C_DRK,PP_ALIGN.LEFT,False,2)])
notes(sl,"PINNs: physics in loss, no labelled data needed. IC-consistent form hard-codes IC. Window decomposition prevents gradient pathology. FEM anchor resets drift. Fourier embedding fixes spectral bias for Bayesian only. Self-adaptive weights auto-balance loss terms.","PINN'ler: fizik kayipla kodlanir, etiketli veri gerekmez. IC-tutarli form IC'yi sabit kodlar. Pencere ayristirmasi gradyan patolojisini onler. FEM capalama sapmay sifirlar. Fourier yalnizca Bayesian icin.")

# S07 NAS + CHAIN
sl=ns(); header(sl,"Theory: NAS and FEM-NAS-PINN Chain","Why NAS for PINNs  Search space  How the three components connect"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,0.3,1.12,5.5,5.75,fill=C_LTBG,line=C_RUL)
txt(sl,0.45,1.22,5.3,5.5,[("Why NAS?",14,True,C_PRI,PP_ALIGN.LEFT),("PINNs extremely sensitive to architecture.",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("Manual tuning = slow, not reproducible.",12,True,C_RED,PP_ALIGN.LEFT,False,2),("NAS automates: search ONCE on rectangle",12,True,C_GRN,PP_ALIGN.LEFT,False,3),("then transfer to all 10 geometries.",12,True,C_GRN,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Search space:",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("Layers: 2-6    Neurons: 32-256",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("Activations: ReLU / tanh / SiLU / GELU",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Three architectures found:",13,True,C_PRI,PP_ALIGN.LEFT,False,4),("Bayesian/TPE  5x151 ReLU   93K params",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("NSGA-II       3x153 tanh   48K params",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("NSGA-III      3x75  tanh   12K params",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2)])
rect(sl,6.1,1.12,6.93,5.75,fill=C_LTBG,line=C_RUL)
t1(sl,"FEM to NAS to PINN Chain:",6.25,1.22,6.63,0.35,sz=13,bold=True,col=C_PRI)
for y,bg,ttl,body in [(1.68,C_ACC,"1. FEM Reference Solver","Provides 21 snapshots as ground truth.\nUsed as IC for each PINN window."),(2.75,RGBColor(45,130,100),"2. NAS Search (once on rectangle)","150-300 arch candidates evaluated.\nSelects 3 best via validation MAE."),(3.82,RGBColor(100,60,140),"3. NAS-PINN (all 10 domains)","Best arch trained per window.\nT_hat = T_ic + tau N_phi (IC-consistent)."),(4.89,C_PRI,"4. MSWP Controller","Decides FEM or PINN per window.\nSkip factor k: fixed or adaptive.")]:
    rect(sl,6.25,y,6.63,0.9,fill=bg)
    txt(sl,6.4,y+0.04,6.35,0.82,[(ttl,11,True,C_WHI,PP_ALIGN.LEFT),(body,10,False,C_WHI,PP_ALIGN.LEFT,False,3)])
notes(sl,"NAS: search once on rectangle, deploy on 10 domains. Chain: FEM anchors, NAS finds arch, PINN trained per window, MSWP controls budget.","NAS: dikdortgene bir kez ara, 10 alana deploy et. Zincir: FEM, NAS, PINN, MSWP.")

# S08 BAYESIAN + NSGA
sl=ns(); header(sl,"Theory: Optimisers — Bayesian/TPE vs NSGA-II vs NSGA-III","What  How  Why chosen  Result and contribution of each"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
card(sl,0.3,1.12,4.25,5.75,"Bayesian/TPE",["Builds PROBABILISTIC MODEL of which","archs give low MAE.","","TPE: maintains TWO densities:","  l(x) = good configs","  g(x) = bad configs","Next trial = argmax l(x)/g(x)","","150 trials on rectangle. Optuna.","","WHY: sample-efficient, handles","mixed search space (int+categorical).","","RESULT: 5x151 ReLU 93K","k=5 on 7/8 domains. Best overall."],hbg=C_PRI,tsz=12,fsz=10.5)
card(sl,4.8,1.12,4.15,5.75,"NSGA-II",["GENETIC ALGORITHM:","Each generation:","  Evaluate: train PINN  MAE+params","  Select: non-dominated sorting","  Crossover: mix parent archs","  Mutate: random changes","","TWO OBJECTIVES:","  Minimise MAE  (accuracy)","  Minimise params  (model size)","NSGA-II: crowding distance.","15 gen, population 20.","","RESULT: 3x153 tanh 48K","k=3 safe. Fails Cylinder k=5."],hbg=C_ACC,tsz=12,fsz=10.5)
card(sl,9.2,1.12,3.83,5.75,"NSGA-III",["Same genetic algorithm as NSGA-II","but uses REFERENCE POINTS","instead of crowding distance.","Better spread for 3+ objectives.","","15 gen, population 24.","pymoo library.","","WHY vs NSGA-II:","More systematic diversity,","better for multi-objective.","","RESULT: 3x75 tanh 12K","Most compact architecture.","Beats Bayesian on L-Shape k=1.","Safe on Thermal Fin k=4 (75%)."],hbg=RGBColor(45,130,100),tsz=12,fsz=10.5)
notes(sl,"Bayesian: probabilistic model, l/g ratio, 150 trials, sample-efficient. Result: 93K ReLU, best overall. NSGA-II: genetic algorithm, two objectives (MAE+params), crowding distance. Result: 48K tanh, reliable k=3. NSGA-III: reference points, better spread. Result: 12K tanh, most compact, good on irregular shapes.","Bayesian: olasiliksal model, l/g oran, 150 deneme. NSGA-II: genetik, iki amac, yogunluk mesafesi. NSGA-III: referans noktalari, daha iyi yayilim.")

# S09 PROBLEM STATEMENT
sl=ns(); header(sl,"Problem Statement","Skip factor k  Three prediction modes  Engineering accuracy band"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,0.5,1.12,12.33,0.72,fill=C_PRI)
t1(sl,"How many FEM solver calls can be replaced by PINN predictions without exceeding an engineering temperature accuracy band?",0.7,1.18,12.0,0.6,sz=14,bold=True,col=C_WHI,align=PP_ALIGN.CENTER)
img(sl,"fig_mswp_three_modes.png",0.3,2.0,5.6,4.9)
rect(sl,6.2,2.0,3.0,3.1,fill=C_LTBG,line=C_RUL)
txt(sl,6.35,2.05,2.8,3.0,[("Skip Factor k",13,True,C_PRI,PP_ALIGN.CENTER),("k  FEM calls  Saving",11,True,C_TG,PP_ALIGN.LEFT,False,8),("1      20       0%",11,False,C_DRK,PP_ALIGN.LEFT,False,3),("2      10      50%",11,False,C_DRK,PP_ALIGN.LEFT,False,3),("3       7      65%",11,False,C_DRK,PP_ALIGN.LEFT,False,3),("4       5      75%",11,False,C_DRK,PP_ALIGN.LEFT,False,3),("5       4      80%",11,True,C_GRN,PP_ALIGN.LEFT,False,3)])
rect(sl,9.5,2.0,3.53,3.1,fill=C_LTBG,line=C_RUL)
txt(sl,9.65,2.05,3.28,3.0,[("3 Prediction Modes",13,True,C_PRI,PP_ALIGN.CENTER),("1. FEM-anchored fixed k",12,True,C_ACC,PP_ALIGN.LEFT,False,8),("   Anchor every k-th step",10.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("2. FEM-anchored adaptive k",12,True,C_ACC,PP_ALIGN.LEFT,False,5),("   k adjusts to error signal",10.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("3. PINN-only stress test",12,True,C_ACC,PP_ALIGN.LEFT,False,5),("   No FEM correction at all",10.5,False,C_DRK,PP_ALIGN.LEFT,False,2)])
rect(sl,6.2,5.35,6.83,0.95,fill=C_YLW,line=C_AMB)
t1(sl,"10-15C: engineering criterion defined in THIS THESIS. DeltaT~520C so 15C~2-3% relative. NOT from Mortensen. See notes for jury answer.",6.35,5.38,6.65,0.85,sz=10,italic=True,col=C_DRK)
notes(sl,"JURY QUESTION 15C threshold:\nThis is an engineering criterion defined in THIS THESIS. Mortensen criterion is distortion (mm) not temperature. DeltaT=520C so 15C=2-3% relative.\nBEST ANSWER: This is an engineering acceptance criterion defined in this thesis not a universal threshold from the literature.\n--- TURKCE ---\n15C jury sorusu: Bu tezde tanimlanmis bir kriter. Mortensen kriteri distorsiyon. 15C=yuzde 2-3 goreceli hata.\nSOYLENECEK: Bu tezde tanimlanmis bir muhendislik kabul kriteridir literaturden dogrudan alinmis evrensel bir esik degildir.","")

# S10 SECTION METHODOLOGY
sl=ns()
secdiv(sl,3,"Methodology","FEM Solver  to  NAS-PINN  to  Framework  to  Training  to  MSWP")
notes(sl,"Method order: FEM first, MSWP last.","Yontem sirasi: once FEM, en son MSWP.")

# S11 FEM PROCESS
sl=ns(); header(sl,"Methodology: FEM Reference Solver","Crank-Nicolson  21 snapshots  Ground truth"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl,"fig7_fem_process.png",0.3,1.12,7.0,5.95)
txt(sl,7.6,1.22,5.35,5.75,[("What FEM gives us:",14,True,C_PRI,PP_ALIGN.LEFT),("Accurate T(x,t) at every point.",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("21 snapshots = ground truth.",12,True,C_ACC,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Setup:",13,True,C_PRI,PP_ALIGN.LEFT,False,4),("  dt_int=0.5s  Crank-Nicolson",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("  Export every dt_a=1.5s",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("  21 snapshots t=0..30s",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("  Each snapshot = one FEM call",11.5,True,C_ACC,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Bottleneck on real mesh:",13,True,C_RED,PP_ALIGN.LEFT,False,4),("  2 min per snapshot = 40 min per run",11.5,True,C_RED,PP_ALIGN.LEFT,False,3),("  100 runs = 67 hours",11.5,True,C_RED,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Our strategy:",13,True,C_GRN,PP_ALIGN.LEFT,False,4),("Call FEM only every k-th step.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("PINN predicts steps in between.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2)])
notes(sl,"FEM is ground truth. Real mesh: 2 min per call. Reduce calls by factor k.","FEM referans cozumumuz. 2 dk per cagri. k faktoru ile azaltiyoruz.")

# S12 NAS-PINN ARCH
sl=ns(); header(sl,"Methodology: NAS-PINN Architecture","IC-consistent MLP  Three discovered archs  Deploy everywhere"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl,"nas_pinn_framework.png",0.3,1.12,6.55,5.95)
txt(sl,7.1,1.22,5.85,2.0,[("Search space:",14,True,C_PRI,PP_ALIGN.LEFT),("Layers 2-6  Neurons 32-256 (uniform)",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("Activations: ReLU tanh SiLU GELU",12,False,C_DRK,PP_ALIGN.LEFT,False,3)])
rect(sl,7.1,3.3,5.85,3.45,fill=C_LTBG,line=C_RUL); rect(sl,7.1,3.3,5.85,0.42,fill=C_PRI)
t1(sl,"Three Architectures Found",7.25,3.33,5.55,0.38,sz=13,bold=True,col=C_WHI,align=PP_ALIGN.CENTER)
txt(sl,7.18,3.8,5.65,2.85,[("Strategy           Network     Params",11,True,C_TG,PP_ALIGN.LEFT),("Bayesian/TPE  5x151 ReLU  93K  150 trials",11.5,False,C_DRK,PP_ALIGN.LEFT,False,5),("NSGA-II         3x153 tanh  48K  15gen/20",11.5,False,C_DRK,PP_ALIGN.LEFT,False,4),("NSGA-III         3x75  tanh  12K  15gen/24",11.5,False,C_DRK,PP_ALIGN.LEFT,False,4),("",5,False,C_DRK,PP_ALIGN.LEFT),("NAS run ONCE on rectangle",12,True,C_ACC,PP_ALIGN.LEFT,False,4),("deployed on all 10 geometries.",12,False,C_DRK,PP_ALIGN.LEFT)])
notes(sl,"NAS finds archs once, deploys on 10 domains.","NAS bir kez arar, 10 alana deploy eder.")

# S13 FRAMEWORK
sl=ns(); header(sl,"Methodology: Full Framework - FEM + NAS-PINN Together","How the three components connect as a pipeline"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl,"fig6_naspinn_process.png",0.3,1.12,7.35,5.95)
txt(sl,7.9,1.22,5.05,5.75,[("Three components:",14,True,C_PRI,PP_ALIGN.LEFT),("",4,False,C_DRK,PP_ALIGN.LEFT),("1. FEM Reference Solver",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("Provides T_FEM(ti) at anchor steps.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("Called every k-th snapshot only.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("2. NAS-PINN Predictor",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("IC = T_FEM from FEM anchor.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("Trained per window independently.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("T_hat = T_ic + tau N_phi(x,t)",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("3. MSWP Controller",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("Decides: FEM or PINN?",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("Skip factor k: fixed or adaptive.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2)])
notes(sl,"Full pipeline: FEM anchors to PINN predicts between to MSWP controls budget.","Tam boru hatti: FEM capa, PINN tahmin, MSWP butce kontrolu.")

# S14 TRAINING PROTOCOL
sl=ns(); header(sl,"Methodology: Training Protocol - Adam + L-BFGS","Two-phase optimizer  Self-adaptive weights  Budget per domain"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,0.3,1.12,5.85,5.75,fill=C_LTBG,line=C_RUL)
txt(sl,0.45,1.22,5.6,5.55,[("Phase 1 - Adam:",13,True,C_ACC,PP_ALIGN.LEFT),("Cosine annealing lr0=8e-4 to lrmin=1e-5",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("Robust global convergence.",12,False,C_DRK,PP_ALIGN.LEFT,False,2),("800-2000 epochs (domain-dependent).",12,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Phase 2 - L-BFGS:",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("Quasi-Newton fine-tuning near optimum.",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("50-100 steps. Sharpens steep regions.",12,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Self-adaptive weights:",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("wbc, wend: LEARNABLE",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("wpde, wic: fixed",12,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Key insight:",13,True,C_PRI,PP_ALIGN.LEFT,False,4),("Training cost proportional to 20/k windows.",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("k=5 trains 5x faster than k=1.",12,True,C_GRN,PP_ALIGN.LEFT,False,2)])
rect(sl,6.4,1.12,6.6,0.38,fill=C_PRI)
t1(sl,"Training Budget per Window (V100 GPU)",6.55,1.15,6.35,0.34,sz=11.5,bold=True,col=C_WHI)
htbl(sl,6.4,1.52,6.6,4.3,["Domain","Arch.","Adam","L-BFGS","n_dom","n_bc","Time"],[["2D","Bayesian","800","50","1500","300","~10s"],["2D","NSGA-II/III","1500","50","1500","300","~16s"],["3D","Bayesian","800","100","12000","2400","~25s"],["3D","NSGA-II/III","1500","100","12000","2400","~35s"],["T.Fin","Bayesian*","2000","100","12000","2400","~43s"],["T.Fin","NSGA-II/III","2000","100","12000","2400","~31s"]],sz=10)
keybx(sl,6.4,5.95,6.6,0.85,"*Bayesian/TPE uses Fourier embedding on Thermal Fin only.",sz=11)
notes(sl,"Two-phase: Adam global + L-BFGS fine-tuning. Training scales with window count 20/k. k=5 trains faster.","Iki asama: Adam + L-BFGS. Egitim 20/k pencere sayisiyla olcekleniyor.")

# S15 MSWP LAST
sl=ns(); header(sl,"Methodology: MSWP - The Controller (LAST Component)","Fixed-k algorithm  Adaptive-k  Error-driven k adjustment"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl,"FEM_Proses.png",0.3,1.12,6.55,5.95)
txt(sl,7.1,1.22,5.85,5.75,[("Algorithm 1 - Fixed k:",14,True,C_PRI,PP_ALIGN.LEFT),("1. IC = T_FEM(ti)   [FEM anchor]",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("2. Train PINN over [ti, ti+k dt]",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("3. Predict T_hat(ti+1)",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("4. epsilon = MAE(T_hat, T_FEM)",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("5. Repeat for all 20/k windows",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("",5,False,C_DRK,PP_ALIGN.LEFT),("Algorithm 3 - Adaptive k:",14,True,C_PRI,PP_ALIGN.LEFT,False,4),("Promote: eps<tau_up 2 windows to k+1",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("Demote:  eps>tau_down         to k-1",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("Caps: 2 early / 4 mid / 5 late",12,False,C_TG,PP_ALIGN.LEFT,False,3),("",5,False,C_DRK,PP_ALIGN.LEFT),("Hardware:",11.5,True,C_ACC,PP_ALIGN.LEFT,False,4),("V100-SXM3 32GB  Xeon 8168  512GB",10.5,False,C_DRK,PP_ALIGN.LEFT,False,2)])
notes(sl,"MSWP is LAST component: controls FEM vs PINN per window.","MSWP son bilesen: her pencere icin FEM mi PINN mi karari.")

# S16 2D DOMAINS
sl=ns(); header(sl,"Methodology: 2D Benchmark Domains","Rectangle  Circle  L-Shape   Adaptive-k adds: Square  Flower"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl,"fig2_thesis_2d.png",0.3,1.12,7.0,5.95)
txt(sl,7.6,1.22,5.35,5.75,[("Fixed-k canonical (3 domains):",14,True,C_PRI,PP_ALIGN.LEFT),("  Rectangle - simplest, convex",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("  Circle - curved boundary normals",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  L-Shape - re-entrant corner",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  (gradient concentration point)",11,False,C_TG,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Adaptive-k adds (4 domains):",13,True,C_PRI,PP_ALIGN.LEFT,False,4),("  Square - smooth, easy baseline",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  Flower - curved multi-lobe",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("",4,False,C_DRK,PP_ALIGN.LEFT),("All 2D domains:",12,True,C_ACC,PP_ALIGN.LEFT,False,4),("T0=540C  Tw=20C  30-second quench",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("10 seeds per domain arch k",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("Threshold for 2D: 10 C",11.5,True,C_PRI,PP_ALIGN.LEFT,False,2)])
notes(sl,"3 canonical 2D for fixed-k (threshold 10C). 4 for adaptive-k.","Sabit-k icin 3 kanonik 2B (esik 10C). Adaptif-k icin 4.")

# S17 3D DOMAINS
sl=ns(); header(sl,"Methodology: 3D and Thermal Fin Benchmark Domains","4 canonical 3D shapes  Thermal Fin (hardest closest to real subframe)"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl,"fig3_thesis_3d.png",0.3,1.12,5.5,4.2)
img(sl,"fig4_thermalfin.png",6.0,1.12,4.5,4.2)
for x,ttl,sub in [(0.3,"Canonical 3D","Rectangular Prism  Cylinder\nStacked Cubes  L-Shape Prism"),(6.0,"Thermal Fin (3D, hardest)","RBniCS benchmark to 3D transient\n5 sub-domains  22 Robin surfaces")]:
    rect(sl,x,5.37,5.2,1.38,fill=C_LTBG,line=C_RUL)
    t1(sl,ttl,x+0.1,5.41,5.0,0.38,sz=13,bold=True,col=C_PRI)
    t1(sl,sub,x+0.1,5.8,5.0,0.9,sz=10.5,col=C_TG)
txt(sl,10.7,1.22,2.33,4.9,[("Why Thermal Fin is hardest:",12,True,C_RED,PP_ALIGN.LEFT),("",3,False,C_DRK,PP_ALIGN.LEFT),("Fin TIP:",11,True,C_DRK,PP_ALIGN.LEFT,False,3),("540 to 100C in first 1.5s",11,False,C_DRK,PP_ALIGN.LEFT,False,2),("Extreme gradient.",11,True,C_RED,PP_ALIGN.LEFT,False,2),("",3,False,C_DRK,PP_ALIGN.LEFT),("22 Robin surfaces.",11,False,C_DRK,PP_ALIGN.LEFT,False,3),("Base-fin junctions:",11,False,C_DRK,PP_ALIGN.LEFT,False,2),("singularities.",11,False,C_TG,PP_ALIGN.LEFT,False,1),("",3,False,C_DRK,PP_ALIGN.LEFT),("Closest proxy",11,True,C_ACC,PP_ALIGN.LEFT,False,3),("to real A356",11,True,C_ACC,PP_ALIGN.LEFT,False,1),("subframe.",11,True,C_ACC,PP_ALIGN.LEFT,False,1),("Threshold: 15C",11,True,C_PRI,PP_ALIGN.LEFT,False,3)])
notes(sl,"4 canonical 3D (threshold 15C). Thermal Fin hardest: extreme fin-tip gradient, 22 surfaces.","4 kanonik 3B (esik 15C). Termal Fin en zor: fin ucu asin gradyan, 22 yuzey.")

# S18 SECTION RESULTS
sl=ns()
secdiv(sl,4,"Results","2D Fixed-k  Adaptive-k  3D  Thermal Fin  PINN-only")
notes(sl,"Results follow thesis chapter order.","Sonuclar tez bolum sirasini takip ediyor.")

# S19 2D MAE TABLE 4.1
sl=ns(); header(sl,"Results: 2D Fixed-k - Table 4.1 MAE + L2 Figure","10-seed mean+-std  Threshold 10C  Why MAE and L2 rank architectures differently"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,0.3,1.12,7.5,0.38,fill=C_PRI)
t1(sl,"Table 4.1 - Mean-Window MAE (C)  10 seeds  Threshold 10C",0.45,1.15,7.3,0.34,sz=11,bold=True,col=C_WHI)
htbl(sl,0.3,1.52,7.5,3.42,["Domain","Arch.","k=1","k=2","k=3","k=4","k=5"],[["Rectangle","Bayes","2.12+-0.23","2.81+-0.30","3.19+-0.57","3.69+-0.84","4.03+-0.66"],["","NSG-II","3.44+-0.49","5.40+-0.82","6.76+-1.14","8.33+-1.65","9.66+-1.30"],["","NSG-III","3.52+-0.48","5.11+-0.70","6.10+-0.45","7.95+-0.89","9.28+-1.64"],["Circle","Bayes","2.75+-0.33","3.58+-1.17","4.93+-1.95","6.13+-3.36","10.01+-5.16"],["","NSG-II","4.89+-1.50","7.58+-2.86","7.34+-3.63","9.31+-4.61","14.97+-7.81"],["","NSG-III","4.01+-1.16","5.83+-2.41","6.82+-3.55","9.45+-5.96","11.33+-7.58"],["L-Shape","Bayes","2.29+-0.24","3.41+-0.37","3.69+-0.66","4.74+-1.01","5.41+-1.06"],["","NSG-II","3.50+-0.24","6.61+-0.89","8.53+-0.99","11.33+-1.98","14.37+-2.10"],["","NSG-III","2.87+-0.46","5.13+-0.65","7.30+-1.47","9.88+-1.90","11.29+-3.04"]],sz=8.8)
img(sl, "l2_all_k_2d.png", 8.05,1.12,5.0,3.05, sub="seed_results")
rect(sl,8.05,4.27,5.0,0.38,fill=RGBColor(80,40,120))
t1(sl,"Why MAE and L2 rank differently:",8.2,4.3,4.8,0.34,sz=12,bold=True,col=C_WHI)
txt(sl,8.1,4.72,4.85,2.1,[("MAE = absolute error in C.",10.5,False,C_DRK,PP_ALIGN.LEFT),("L2 = ||T_hat-T_FEM||2 / ||T_FEM||2",10.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("(normalised by field magnitude)",10,False,C_TG,PP_ALIGN.LEFT,False,1),("",4,False,C_DRK,PP_ALIGN.LEFT),("If errors are in HOT regions:",10.5,True,C_ACC,PP_ALIGN.LEFT,False,3),("large ||T_FEM|| denominator",10,False,C_DRK,PP_ALIGN.LEFT,False,2),("divides them away in L2.",10,False,C_DRK,PP_ALIGN.LEFT,False,2),("L2 looks small even if MAE is large.",10,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("NSGA: lower L2 but higher MAE",10.5,True,C_RED,PP_ALIGN.LEFT,False,3),("than Bayesian on some domains.",10,False,C_DRK,PP_ALIGN.LEFT,False,2),("MAE is the engineering criterion.",10.5,True,C_GRN,PP_ALIGN.LEFT,False,2)])
notes(sl,"Table 4.1. Threshold 10C for 2D. Bayesian k=5 safe on Rectangle and L-Shape. Circle: curved normals cause high variance. MAE vs L2 discrepancy: L2 normalised by ||T_FEM||2. NSGA errors in HOT regions where T_FEM large - L2 denominator divides them away but MAE stays absolute. Example: Square NSGA-II MAE=1.768 > Bayesian 1.641 but L2=0.015 < 0.017. MAE is the engineering criterion.","Tablo 4.1. 2B esik 10C. Bayesian Dikdortgen ve L-Sekil k=5 guvenli. MAE vs L2: L2 normallendirilmis, NSGA sicak bolgede hata yaparsa buyuk payda L2'yi kucultse de MAE mutlak kalir.")

# S20 2D ADAPTIVE TABLE 4.2
sl=ns(); header(sl,"Results: 2D Adaptive-k - Table 4.2  (NAS Architecture Transfer)","k-schedule physics-driven  Why 83-85%  MAE vs L2 again"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,0.3,1.12,7.5,0.38,fill=C_PRI)
t1(sl,"Table 4.2 - Adaptive-k Results  10 seeds  FEM saved = 83-85%",0.45,1.15,7.3,0.34,sz=11,bold=True,col=C_WHI)
htbl(sl,0.3,1.52,7.5,2.82,["Domain","Architecture","MAE (C)","L2","Saving","k-schedule"],[["Square","Bayesian/TPE","1.64+-0.10","0.017","85%","1-1-1-1-2-3-4-5-5"],["Square","NSGA-II","1.77+-0.06","0.015*","85%","1-1-1-1-2-3-4-5-5"],["Square","NSGA-III","1.72+-0.04","0.015*","85%","1-1-1-1-2-3-4-5-5"],["Circle","Bayesian/TPE","1.53+-0.10","0.017","85%","1-1-1-1-2-3-4-5-5"],["L-Shape","Bayesian/TPE","1.36+-0.03","0.017","85%","1-1-1-1-2-2-3-4-5"],["L-Shape","NSGA-II","2.44+-0.05","0.023","83%","1-1-1-1-1-2-3-4-5-5"],["Flower","Bayesian/TPE","1.74+-0.10","0.015","85%","1-1-1-1-2-3-4-5-5"]],sz=9)
txt(sl,8.05,1.12,5.0,5.75,[("Why adaptive-k achieves 83-85%?",14,True,C_PRI,PP_ALIGN.LEFT),("",4,False,C_DRK,PP_ALIGN.LEFT),("k-schedule 1-1-1-1-2-3-4-5-5:",13,True,C_ACC,PP_ALIGN.LEFT,False,3),("Windows 1-4  (t=0..6s): k=1",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  RAPID cooling. FEM every step.",11,False,C_TG,PP_ALIGN.LEFT,False,1),("Windows 5-8 (t=6..18s): k=2,3,4",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  Cooling slows. Gradual extension.",11,False,C_TG,PP_ALIGN.LEFT,False,1),("Window 9  (t=18..30s): k=5",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  Slow approach to equilibrium.",11,False,C_TG,PP_ALIGN.LEFT,False,1),("",4,False,C_DRK,PP_ALIGN.LEFT),("Physics drives the schedule!",12,True,C_GRN,PP_ALIGN.LEFT,False,3),("Same schedule for ALL architectures.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("Controller responds to COOLING,",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("not architecture artefacts.",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("* NSGA lower L2 but higher MAE:",12,True,C_AMB,PP_ALIGN.LEFT,False,3),("Errors in HOT zones, L2 normalises.",11,False,C_DRK,PP_ALIGN.LEFT,False,2)])
keybx(sl,0.3,4.48,7.5,0.72,"* NSGA lower L2 but higher MAE than Bayesian on Square/Circle. Reason: NSGA errors in HIGH-TEMPERATURE regions. Large T_FEM denominator reduces L2. MAE absolute - stays higher. Engineering criterion = MAE.",sz=10)
notes(sl,"Table 4.2. All archs 83-85% saving, MAE < 2.5C. k-schedule physics-driven: k=1 when rapid cooling t=0..6s, gradually k=5 as T approaches equilibrium. Same for all archs: physics not architecture. * MAE vs L2 on Square/Circle: NSGA errors in hot regions, large norm reduces L2 but MAE absolute stays higher.","Tablo 4.2. Tum mimariler yuzde 83-85, MAE<2.5C. k-programi fizik gudumlu. * MAE vs L2: NSGA sicak bolgede hata, buyuk payda L2 azaltir, MAE mutlak kalir.")

# S21 2D PINN-ONLY
sl=ns(); header(sl,"Results: 2D - PINN-Only Stress Test","What happens without FEM correction  Error amplification"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl,"thermal_fin_pinn_only_best_clean_heatmaps.png",0.3,1.12,5.5,3.8)
rect(sl,0.3,5.0,5.5,1.85,fill=C_LTBG,line=C_RUL); rect(sl,0.3,5.0,5.5,0.4,fill=C_PRI)
t1(sl,"Anchored vs PINN-only comparison",0.45,5.03,5.3,0.36,sz=12,bold=True,col=C_WHI)
stbl(sl,0.3,5.42,5.5,1.4,["Domain","Anchored","PINN-only","Factor"],[["2D Square","1.6C","29.6C","x18"],["TF NSGA-III","14.1C","59.2C","x4"],["TF Bayesian","12.6C",">200C","diverges"]],sz=10.5)
txt(sl,6.1,1.22,6.9,5.75,[("PINN-only = no FEM corrections.",14,True,C_PRI,PP_ALIGN.LEFT),("",4,False,C_DRK,PP_ALIGN.LEFT),("Why errors explode x4 to x18?",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("  PINN output used as IC for next window",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("  Errors ACCUMULATE - no FEM reset",11.5,True,C_RED,PP_ALIGN.LEFT,False,2),("  After 20 windows: catastrophic drift",11.5,True,C_RED,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Why Bayesian fails fastest?",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("  More expressive capacity (93K ReLU)",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("  Learns and amplifies errors faster",11.5,True,C_RED,PP_ALIGN.LEFT,False,2),("  without periodic correction.",11.5,True,C_RED,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Conclusion:",13,True,C_PRI,PP_ALIGN.LEFT,False,4),("FEM anchoring is STRUCTURALLY",13,True,C_ACC,PP_ALIGN.LEFT,False,3),("REQUIRED - not optional.",13,True,C_ACC,PP_ALIGN.LEFT,False,2),("This is FEM-REDUCTION,",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("not FEM-replacement.",12,False,C_DRK,PP_ALIGN.LEFT,False,2)])
notes(sl,"PINN-only: errors explode 4-18x without FEM correction. Bayesian fails fastest (more expressive = faster drift). FEM anchoring structurally required.","PINN-only: FEM duzeltmesi olmadan hatalar 4-18x. Bayesian en hizli (daha ifadeli = daha hizli sapma). FEM capalama zorunlu.")

# S22 3D TABLE 4.3 + ADAPTIVE
sl=ns(); header(sl,"Results: 3D - Table 4.3 Fixed-k + Adaptive-k","4 domains  Bayesian dominates  Why NSGA-II fails on Cylinder  3D adaptive lower than 2D"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,0.3,1.12,7.5,0.38,fill=C_PRI)
t1(sl,"Table 4.3 - MAE (C)  3D Canonical  Threshold 15C",0.45,1.15,7.3,0.34,sz=11,bold=True,col=C_WHI)
htbl(sl,0.3,1.52,7.5,3.52,["Domain","Arch.","k=1","k=2","k=3","k=4","k=5"],[["Rect.","Bayes","8.42+-1.40","8.47+-1.33","8.74+-1.96","9.43+-1.59","9.88+-1.60"],["Rect.","NSG-II","5.47+-1.31","7.62+-2.02","9.91+-2.11","11.01+-2.37","13.42+-4.57"],["Cyl.","Bayes","5.26+-0.92","6.51+-1.30","7.25+-0.99","8.59+-1.00","9.58+-1.07"],["Cyl.","NSG-II","5.79+-1.85","7.96+-3.30","11.36+-5.63","12.67+-8.21","19.21+-12.28"],["Cyl.","NSG-III","5.79+-2.53","7.48+-2.60","8.51+-4.45","12.06+-7.21","9.60+-2.99"],["Stack.","Bayes","10.51+-1.33","11.35+-2.32","11.84+-1.63","13.03+-2.34","14.07+-2.97"],["Stack.","NSG-II","8.69+-3.60","11.33+-3.94","13.31+-4.97","15.24+-4.59","18.19+-5.33"],["L-Sh.","Bayes","10.20+-1.40","12.10+-2.67","12.06+-2.44","12.19+-2.10","13.96+-2.15"],["L-Sh.","NSG-III","4.41+-1.01","7.50+-1.97","8.94+-2.92","10.33+-1.94","12.26+-2.22"]],sz=8.2)
txt(sl,8.05,1.12,5.0,2.88,[("Bayesian k=5 ALL 4 domains",12,True,C_GRN,PP_ALIGN.LEFT),("80% FEM saving",12,True,C_GRN,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("NSGA-II Cylinder k=5 fails:",12,True,C_RED,PP_ALIGN.LEFT,False,3),("sigma=12.28C - Seed C diverges.",11,False,C_DRK,PP_ALIGN.LEFT,False,2),("Curved normals + tanh saturation",11,False,C_TG,PP_ALIGN.LEFT,False,2),("+ compact loss landscape.",11,False,C_TG,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("NSGA-III L-Shape k=1:",12,True,C_ACC,PP_ALIGN.LEFT,False,3),("4.41 < Bayesian 10.20C",11,False,C_DRK,PP_ALIGN.LEFT,False,2),("tanh matches edge gradient.",11,False,C_TG,PP_ALIGN.LEFT,False,2)])
rect(sl,8.05,4.1,5.0,0.38,fill=C_PRI)
t1(sl,"3D Adaptive-k  tau_up=7C tau_down=10C",8.2,4.13,4.8,0.34,sz=11,bold=True,col=C_WHI)
stbl(sl,8.05,4.5,5.0,2.28,["Domain","Arch.","MAE","Saving","Max k"],[["Rect.","Bayesian","9.30C","20%","k=3"],["Cylinder","Bayesian","6.66C","35%","k=3"],["Cylinder","NSGA-II","13.21C","0%","k=1"],["Stacked","Bayesian","8.92C","15%","k=2"],["L-Shape","NSGA-II","13.11C","30%","k=3"]],sz=9)
keybx(sl,8.05,6.83,5.0,0.85,"3D adaptive: only 20-35% saving vs 83-85% 2D. Base MAE at k=1 is 5-10C near tau_up=7C so controller rarely promotes.",sz=10)
notes(sl,"Table 4.3. Bayesian k=5 all 4 3D domains (80%). NSGA-II Cylinder k=5 fails: high seed variance sigma=12.3C, Seed C diverges. Curved normals + tanh saturation. NSGA-III beats Bayesian on L-Shape prism at k=1 (4.41 < 10.20C): tanh bias matches edge gradient. 3D adaptive: only 20-35% saving because base MAE k=1 is 5-10C near tau_up=7C so controller rarely promotes.","Tablo 4.3. Bayesian tum 4 3B alanda k=5 (yuzde 80). NSGA-II Silindir k=5 basarisiz. 3B adaptif: yalnizca yuzde 20-35 tasarruf cunku baz MAE k=1'de tau_up'a yakin.")

# S23 THERMAL FIN
sl=ns(); header(sl,"Results: Thermal Fin - MAE and L2 (Appendix Table A.7)","k=1..5  Bayesian+Fourier reaches k=5 safely  Heatmaps and seed bars"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl,"thermal_fin_bayesian_fixed_k_clean_heatmaps.png",0.3,1.12,4.8,3.25)
img(sl, "seed_bar_thermal_fin.png", 0.3,4.5,4.8,2.38, sub="seed_results")
rect(sl,5.35,1.12,4.5,0.38,fill=C_PRI)
t1(sl,"MAE (C)  (*Bayesian uses Fourier)",5.5,1.15,4.3,0.34,sz=10.5,bold=True,col=C_WHI)
htbl(sl,5.35,1.52,4.5,1.88,["Architecture","k=1","k=2","k=3","k=4","k=5"],[["Bayes/TPE*","5.10+-0.15","6.84+-0.08","8.75+-0.18","10.82+-0.22","12.55+-0.29"],["NSGA-II","8.98+-0.49","11.03+-0.64","12.49+-0.29","14.14+-0.50","15.56+-0.48"],["NSGA-III","8.97+-0.86","11.53+-0.53","12.51+-0.59","14.40+-0.36","15.74+-0.46"]],sz=9.5)
rect(sl,5.35,3.52,4.5,0.38,fill=C_PRI)
t1(sl,"L2 Error  (practical threshold <0.10)",5.5,3.55,4.3,0.34,sz=10.5,bold=True,col=C_WHI)
htbl(sl,5.35,3.92,4.5,1.48,["Architecture","k=1","k=3","k=5"],[["Bayes/TPE*","0.049","0.062","0.081"],["NSGA-II","0.068","0.079","0.083"],["NSGA-III","0.068","0.081","0.084"]],sz=10)
img(sl, "l2_all_k_thermal_fin.png", 5.35,5.52,4.5,2.35, sub="seed_results")
txt(sl,10.05,1.22,3.0,5.75,[("Key findings:",13,True,C_PRI,PP_ALIGN.LEFT),("",3,False,C_DRK,PP_ALIGN.LEFT),("Bayesian+Fourier:",12,True,C_GRN,PP_ALIGN.LEFT,False,3),("k=5 to 12.55C",11.5,True,C_GRN,PP_ALIGN.LEFT,False,2),("80% FEM saving.",11,False,C_DRK,PP_ALIGN.LEFT,False,1),("sigma=0.29C - very",11,False,C_TG,PP_ALIGN.LEFT,False,1),("reproducible.",11,False,C_TG,PP_ALIGN.LEFT,False,1),("",3,False,C_DRK,PP_ALIGN.LEFT),("NSGA-II and III:",12,True,C_AMB,PP_ALIGN.LEFT,False,3),("k=4 to 14.1/14.4C",11.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("75% saving.",11,False,C_DRK,PP_ALIGN.LEFT,False,1),("k=5 to 15.6/15.7C",11.5,True,C_RED,PP_ALIGN.LEFT,False,2),("Just 0.6-0.7C over!",10.5,False,C_TG,PP_ALIGN.LEFT,False,1),("",3,False,C_DRK,PP_ALIGN.LEFT),("All L2 < 0.10",12,True,C_GRN,PP_ALIGN.LEFT,False,3),("at all k values.",11,False,C_DRK,PP_ALIGN.LEFT,False,2),("",3,False,C_DRK,PP_ALIGN.LEFT),("Fourier:",11,True,C_ACC,PP_ALIGN.LEFT,False,3),("Reduces Bayes by",10.5,False,C_DRK,PP_ALIGN.LEFT,False,2),("1.8-4.4C",10.5,True,C_GRN,PP_ALIGN.LEFT,False,1),("Hurts NSGA",10.5,True,C_RED,PP_ALIGN.LEFT,False,2)])
notes(sl,"Appendix Table A.7. Bayesian+Fourier: k=5 safely (12.55C, sigma=0.29). NSGA safe at k=4, marginally exceed 15C at k=5. All L2 < 0.10. Heatmaps show visual quality at k=5. Seed bars: Bayesian reproducible, NSGA wider at k=4-5.","Ek Tablo A.7. Bayesian k=5 guvenli. NSGA k=4'te guvenli, k=5'te hafifce asiyor. Tum L2 < 0.10.")

# S24 SECTION DISCUSSION
sl=ns()
secdiv(sl,5,"Discussion","Cross-domain  Practical guidance  Industrial impact")
notes(sl,"Discussion.","Tartisma.")

# S25 CROSS-DOMAIN
sl=ns(); header(sl,"Discussion: Cross-Domain Analysis and Practical Guide","Best arch per use case  Industrial impact  NAS transfer"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
img(sl, "l2_all_k_2d.png", 0.3,1.12,5.8,3.88, sub="seed_results")
rect(sl,0.3,5.1,5.8,0.38,fill=C_PRI)
t1(sl,"Recommended config per geometry",0.45,5.13,5.6,0.34,sz=11.5,bold=True,col=C_WHI)
stbl(sl,0.3,5.5,5.8,1.35,["Geometry","Best arch.","Best k","Saving"],[["Simple 2D","any (adaptive)","adaptive","83-85%"],["Complex 3D","Bayesian/TPE","k=5","80%"],["Thermal Fin","Bayesian+Fourier","k=5","80%"]],sz=10)
txt(sl,6.35,1.22,6.65,5.75,[("Bayesian/TPE dominates:",14,True,C_PRI,PP_ALIGN.LEFT),("  k=5 on 7/8 canonical domains <=15C",12,False,C_DRK,PP_ALIGN.LEFT,False,4),("  NAS on rectangle: transfers to all 10",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  geometries without re-training.",11,False,C_TG,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("When NSGA-III adds value:",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("  L-Shape prism k=1: 4.41 < 10.20C",11.5,False,C_DRK,PP_ALIGN.LEFT,False,3),("  tanh matches edge gradient.",11,False,C_TG,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Industrial impact (k=4):",13,True,C_ACC,PP_ALIGN.LEFT,False,4),("  40 min/run to ~10 min/run",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  100 design runs:",12,False,C_DRK,PP_ALIGN.LEFT,False,3),("  67 h to 17 h  (50 h saved)",13,True,C_ACC,PP_ALIGN.LEFT,False,2),("  Training amortised in <1 extra run",11,False,C_TG,PP_ALIGN.LEFT,False,3)])
notes(sl,"Bayesian best overall. NAS transfers. 2D: adaptive-k. 3D/TF: Bayesian k=4-5. Industrial: 50 hours saved.","Bayesian genel en iyi. NAS aktariliyor. 2B: adaptif-k. 3B/TF: Bayesian k=4-5. Endustri: 50 saat.")

# S26 SECTION CONCLUSION
sl=ns()
secdiv(sl,6,"Conclusion","Summary  Limitations  Future Steps")
notes(sl,"Conclusion.","Sonuc.")

# S27 SUMMARY
sl=ns(); header(sl,"Conclusion: Summary of Findings"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
for i,(ttl,buls) in enumerate([("Main Results",["65-85% FEM saving within accuracy band","Bayesian/TPE: k=5 on 7/8 canonical domains","Adaptive-k: 83-85% on all 2D domains","Thermal Fin: Bayesian+Fourier k=5 to 12.55C"]),("Architecture",["Bayesian/TPE most robust (93K ReLU)","NSGA-III: value on irregular shapes low k","Fourier embedding: helps large nets only","NAS: run once deploy everywhere"]),("FEM Anchoring",["PINN-only: 4-18x higher errors or diverge","FEM anchoring is STRUCTURALLY required","This is FEM-REDUCTION not FEM-replacement","Periodic reset prevents error drift"]),("Practical Impact",["k=4 saves 30 min per run on real mesh","100 design runs: 67 h to 17 h (50 h saved)","Training amortised after less than 1 run","5-step roadmap to real subframe"])]):
    row,col=divmod(i,2)
    card(sl,0.3+col*6.36,1.12+row*2.97,6.17,2.78,ttl,buls,fsz=12,tsz=13)
notes(sl,"65-85% saving on 10 domains. Bayesian best. FEM anchoring required. 67h to 17h.","10 alanda yuzde 65-85 tasarruf. Bayesian en iyi. FEM capalama zorunlu. 67h-17h.")

# S28 LIMITATIONS + FUTURE
sl=ns(); header(sl,"Conclusion: Limitations and 5-Step Deployment Roadmap"); footer(sl)
rect(sl,0,0.97,13.33,6.2,fill=C_WHI)
rect(sl,0.3,1.12,5.85,5.75,fill=C_LTBG,line=C_RUL); rect(sl,0.3,1.12,5.85,0.4,fill=C_RED)
t1(sl,"Current Limitations",0.45,1.15,5.55,0.36,sz=13,bold=True,col=C_WHI)
txt(sl,0.45,1.6,5.55,5.15,[("Constant h:",12,True,C_DRK,PP_ALIGN.LEFT),("Real quenching: 5-regime boiling curve.",11.5,False,C_TG,PP_ALIGN.LEFT,False,2),("h varies 20x - not modelled here.",11.5,False,C_TG,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Surrogate geometry:",12,True,C_DRK,PP_ALIGN.LEFT,False,4),("Not the real 3M-element Mortensen mesh.",11.5,False,C_TG,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Single FEM solver:",12,True,C_DRK,PP_ALIGN.LEFT,False,4),("No Abaqus StaMiSim validation.",11.5,False,C_TG,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("Constant material properties:",12,True,C_DRK,PP_ALIGN.LEFT,False,4),("kT rho cp not temperature-dependent.",11.5,False,C_TG,PP_ALIGN.LEFT,False,2),("",4,False,C_DRK,PP_ALIGN.LEFT),("10-15C acceptance band:",12,True,C_DRK,PP_ALIGN.LEFT,False,4),("Engineering criterion (this thesis).",11.5,False,C_TG,PP_ALIGN.LEFT,False,2),("NOT a universal literature threshold.",11.5,True,C_RED,PP_ALIGN.LEFT,False,2)])
rect(sl,6.6,1.12,6.35,5.75,fill=C_LTBG,line=C_RUL); rect(sl,6.6,1.12,6.35,0.4,fill=C_PRI)
t1(sl,"5-Step Deployment Roadmap",6.75,1.15,6.05,0.36,sz=13,bold=True,col=C_WHI)
for i,(step,nt) in enumerate([("Apply to Mortensen mesh","infrastructure is ready"),("Add temperature-dependent h(T)","5-regime boiling curve"),("Architecture-aware adaptive-k thresholds","per-arch tau_up tau_down"),("Residual-adaptive collocation","focus near steep gradients"),("Validate against Abaqus StaMiSim","end-to-end verification")]):
    y=1.65+i*0.93
    rect(sl,6.7,y,0.52,0.52,fill=C_ACC)
    t1(sl,str(i+1),6.7,y+0.05,0.52,0.42,sz=16,bold=True,col=C_WHI,align=PP_ALIGN.CENTER)
    txt(sl,7.33,y,5.45,0.88,[(step,12,True,C_DRK,PP_ALIGN.LEFT),(nt,10.5,False,C_TG,PP_ALIGN.LEFT,False,2)])
rect(sl,6.6,6.43,6.35,0.39,fill=C_LBLU)
t1(sl,"Each step independent - no core pipeline restructuring.",6.72,6.45,6.1,0.35,sz=9.5,italic=True,col=C_PRI,align=PP_ALIGN.CENTER)
notes(sl,"Limits: constant h, surrogate geometry, no StaMiSim. 10-15C is thesis design decision. 5 independent steps.","Sinirlar: sabit h, basit geometri, StaMiSim yok. 10-15C tez tasarim karari. 5 bagimsiz adim.")

# S29 THANK YOU  (was S30)
sl=ns()
rect(sl,0,0,13.33,7.5,fill=C_PRI)
rect(sl,0,2.45,13.33,2.55,fill=C_WHI)
rect(sl,0,2.45,13.33,0.06,fill=C_ACC); rect(sl,0,4.94,13.33,0.06,fill=C_ACC)
t1(sl,"Thank You",0,1.0,13.33,1.2,sz=54,bold=True,col=C_WHI,align=PP_ALIGN.CENTER)
t1(sl,"Questions?",0,2.1,13.33,0.44,sz=22,col=C_RUL,align=PP_ALIGN.CENTER)
t1(sl,'"FEM-anchored NAS-PINNs replace 65-85% of FEM calls while maintaining engineering accuracy - validated on 10 domains."',1.0,2.65,11.33,2.05,sz=15,italic=True,col=C_PRI,align=PP_ALIGN.CENTER)
rect(sl,0,6.5,13.33,1.0,fill=C_ACC)
t1(sl,"Omer Cetinkaya   |   University of Agder, Norway   |   June 2026",0,6.7,13.33,0.5,sz=13,col=C_WHI,align=PP_ALIGN.CENTER)
notes(sl,"KEY for 15C jury question: This is an engineering acceptance criterion defined in this thesis not a universal threshold from the literature. DeltaT=520C so 15C = 2-3 percent relative error.","15C jury sorusu icin: Bu tezde tanimlanmis bir muhendislik kabul kriteridir literaturden dogrudan alinmis evrensel bir esik degildir.")

prs.save(OUTPUT)
print(f"Saved  : {OUTPUT}")
print(f"Slides : {len(prs.slides)}  (TOTAL set to {TOTAL})")
if len(prs.slides) != TOTAL:
    print(f"  NOTE: slide count {len(prs.slides)} != TOTAL {TOTAL}")
