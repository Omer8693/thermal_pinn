from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

try:
    from PIL import Image as PILImage
except ImportError:  # pragma: no cover
    PILImage = None


ROOT = Path(__file__).resolve().parent
FIG = ROOT / "Figures"
REPORT_FIG = ROOT / "Raports_figures"
OUT = ROOT / "NAS_PINN_Full_Report_Order_EN_TR_Notes.pptx"

PRIMARY = RGBColor(151, 66, 62)      # light academic red
ACCENT = RGBColor(220, 118, 93)      # warm coral accent
LIGHT = RGBColor(255, 247, 245)      # pale rose paper
DARK = RGBColor(48, 36, 34)
MUTED = RGBColor(116, 88, 82)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(76, 126, 108)       # muted technical green
RED = RGBColor(190, 78, 72)
LINE = RGBColor(235, 202, 195)
CARD = RGBColor(255, 252, 251)
WARM_NOTE = RGBColor(255, 235, 229)
COOL_PANEL = RGBColor(246, 241, 234)


def add_bg(slide, color=WHITE):
    s = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def add_notes(slide, note):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = note
    p.font.name = "Aptos"
    p.font.size = Pt(12)


def add_text(slide, x, y, w, h, value, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = value
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb


def add_title(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(0.66))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    accent = slide.shapes.add_shape(1, 0, Inches(0.66), Inches(13.333), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    add_text(slide, 0.45, 0.11, 12.3, 0.38, title, 19, True, WHITE)
    if subtitle:
        add_text(slide, 0.48, 0.78, 12.0, 0.28, subtitle, 9, False, MUTED)


def add_footer(slide, idx):
    tb = add_text(slide, 0.45, 7.15, 12.4, 0.18, f"NAS-PINN thesis defense | {idx}", 7, False, MUTED)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_bullets(slide, x, y, w, h, items, size=13):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for n, item in enumerate(items):
        p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(3)
    return tb


def add_card(slide, x, y, w, h, head, body, fill=CARD):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = LINE
    add_text(slide, x + 0.16, y + 0.12, w - 0.32, 0.28, head, 11, True, PRIMARY)
    if isinstance(body, list):
        add_bullets(slide, x + 0.16, y + 0.48, w - 0.32, h - 0.56, body, 10)
    else:
        add_text(slide, x + 0.16, y + 0.48, w - 0.32, h - 0.56, body, 10)


def add_step_card(slide, x, y, w, h, num, head, body):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = CARD
    s.line.color.rgb = LINE
    badge = slide.shapes.add_shape(1, Inches(x + 0.18), Inches(y + 0.18), Inches(0.42), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    add_text(slide, x + 0.18, y + 0.25, 0.42, 0.18, str(num), 9, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, x + 0.72, y + 0.16, w - 0.9, 0.34, head, 12, True, PRIMARY)
    add_text(slide, x + 0.22, y + 0.72, w - 0.44, h - 0.86, body, 9, False, DARK)


def add_metric(slide, x, y, w, h, label, value, note=""):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = WARM_NOTE
    s.line.color.rgb = LINE
    add_text(slide, x + 0.12, y + 0.12, w - 0.24, 0.2, label, 8, True, MUTED, PP_ALIGN.CENTER)
    add_text(slide, x + 0.12, y + 0.36, w - 0.24, 0.34, value, 15, True, PRIMARY, PP_ALIGN.CENTER)
    if note:
        add_text(slide, x + 0.12, y + 0.78, w - 0.24, 0.2, note, 7, False, MUTED, PP_ALIGN.CENTER)


def add_image(slide, name, x, y, w, h, frame=True):
    path = REPORT_FIG / name
    if not path.exists():
        path = FIG / name
    if not path.exists():
        path = ROOT / name
    if not path.exists():
        add_card(slide, x, y, w, h, "Missing report figure", name, RGBColor(252, 242, 242))
        return None
    if frame:
        panel = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = LINE
    pad = 0.08
    ix, iy, iw, ih = x + pad, y + pad, w - 2 * pad, h - 2 * pad
    if PILImage is not None:
        with PILImage.open(path) as im:
            px_w, px_h = im.size
        ratio = min(iw / px_w, ih / px_h)
        draw_w, draw_h = px_w * ratio, px_h * ratio
        ix += (iw - draw_w) / 2
        iy += (ih - draw_h) / 2
        iw, ih = draw_w, draw_h
    return slide.shapes.add_picture(str(path), Inches(ix), Inches(iy), width=Inches(iw), height=Inches(ih))


def add_table(slide, x, y, w, h, headers, rows, size=9, header_fill=PRIMARY):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Aptos"
            p.font.size = Pt(size)
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER
    return table


def new_slide(prs, idx, title, note, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT)
    add_title(s, title, subtitle)
    add_footer(s, idx)
    add_notes(s, note)
    return s


def section_slide(prs, idx, title, subtitle, note):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_text(s, 0.8, 2.55, 11.9, 0.8, title, 34, True, WHITE, PP_ALIGN.CENTER)
    add_text(s, 1.2, 3.45, 10.9, 0.44, subtitle, 16, False, RGBColor(255, 226, 218), PP_ALIGN.CENTER)
    add_footer(s, idx)
    add_notes(s, note)
    return s


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    i = 1

    def ph(slide, x, y, w, h, label, detail=""):
        body = detail if detail else "Insert directly from the thesis/report."
        add_card(slide, x, y, w, h, label, body, WARM_NOTE)

    # 1 Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    add_text(s, 0.8, 1.0, 11.9, 1.0, "NAS-Guided FEM-Anchored PINNs", 36, True, WHITE, PP_ALIGN.CENTER)
    add_text(s, 1.2, 2.05, 10.9, 0.55, "Multi-Step Window Prediction for Transient A356 Water Quenching", 18, False, RGBColor(255, 226, 218), PP_ALIGN.CENTER)
    add_text(s, 1.4, 3.15, 10.5, 0.65, "A FEM accelerator inspired by NAS-PINN and motivated by industrial subframe simulations", 17, False, WHITE, PP_ALIGN.CENTER)
    add_text(s, 1.4, 4.65, 10.5, 0.55, "Omer Cetinkaya | University of Agder | 2026", 15, False, RGBColor(248, 225, 204), PP_ALIGN.CENTER)
    add_image(s, "visual_abstract/visual_abstract.png", 4.25, 5.35, 4.8, 1.35, frame=False)
    add_notes(s, "EN: In this project, we study transient water quenching of A356 aluminium and ask whether a FEM simulation can be accelerated without removing FEM completely. The main idea is to keep FEM as a reliable correction source at anchor points, and use NAS-guided PINNs to predict the intermediate time windows. So the work is not a full FEM replacement; it is a FEM-call reduction strategy for an industrial heat-treatment problem.\nTR: Bu projede A356 alüminyumun su verme sırasındaki geçici ısı transferi problemini ele aldık. Amacımız FEM'i tamamen ortadan kaldırmak değil; FEM'i güvenilir anchor/correction noktası olarak tutup aradaki zaman adımlarını NAS-PINN ile tahmin ederek FEM çağrı sayısını azaltmaktır. Yani bu çalışma bir FEM replacement değil, endüstriyel quenching problemi için FEM hızlandırma yaklaşımıdır.")
    i += 1

    # 2 Structure
    s = new_slide(prs, i, "Presentation Structure", "EN: This slide shows the logic of the presentation. We first introduce the industrial A356 quenching problem and the two papers that guide the thesis: Mortensen et al. for the physical problem and Wang and Zhong for NAS-PINN. Then we explain the theory, the FEM-anchored NAS-PINN method, the benchmark domains, the results, and finally the discussion and conclusion. The results follow the report structure: tables and metrics first, then figures and interpretation.\nTR: Bu slaytta sunumun akışını anlat. Önce endüstriyel A356 quenching problemi ve çalışmanın dayandığı iki temel kaynak tanıtılıyor: fiziksel problem için Mortensen et al., yöntem temeli için Wang and Zhong NAS-PINN. Daha sonra teori, FEM-anchored NAS-PINN yöntemi, benchmark domainler, sonuçlar, discussion ve conclusion geliyor. Results kısmında rapordaki sıraya uyuyoruz: önce tablo ve metrikler, sonra figürler ve yorum.")
    add_step_card(s, 0.35, 0.95, 3.05, 1.8, 1, "Intro", "Reference paper context; A356 quenching problem; why we used the same physical/process data; delayed access to real industrial data; simplified-to-complex surrogate geometry path.")
    add_step_card(s, 3.55, 0.95, 3.05, 1.8, 2, "Theory", "Transient heat equation; FEM; PINN; NAS-PINN; Bayesian/TPE, NSGA-II and NSGA-III optimisation.")
    add_step_card(s, 6.75, 0.95, 3.05, 1.8, 3, "Method", "FEM process; NAS-PINN process; MSWP; 2D NAS transfer domains; canonical 2D domains; 3D and Thermal Fin domains.")
    add_step_card(s, 9.95, 0.95, 3.0, 1.8, 4, "Results", "2D NAS transfer/adaptive-k; canonical 2D; canonical 3D; Thermal Fin. Each section: MAE/L2 tables first, then figures and heat maps.")
    add_step_card(s, 0.35, 3.35, 3.95, 2.45, 5, "Discussion", "2D adaptive-k interpretation; canonical 2D interpretation; canonical 3D interpretation; Thermal Fin interpretation; cross-domain analysis and MSWP significance.")
    add_step_card(s, 4.65, 3.35, 3.95, 2.45, 6, "Conclusion", "Summary of results; limitations; future work; closing remarks. The final message is that the framework is a FEM accelerator, not a full FEM replacement.")
    add_step_card(s, 8.95, 3.35, 3.95, 2.45, 7, "References", "Mortensen et al. for the industrial quenching problem and Wang and Zhong for the NAS-PINN methodological foundation.")
    i += 1

    # 3 Intro opening
    s = new_slide(prs, i, "Intro: What We Do in This Project", "EN: Here I introduce the main problem and the motivation. The reference paper studies water quenching of A356 aluminium automotive subframes, where rapid cooling creates strong temperature gradients. These gradients are important because they later drive residual stress and distortion. In this thesis, we used the same physical process setting and asked a different computational question: can we skip selected FEM time-step solves by using a FEM-anchored NAS-PINN? Since the real industrial data arrived late, we first tested the idea on controlled surrogate geometries, starting simple and then increasing geometric complexity.\nTR: Burada ana problemi ve motivasyonu anlatıyorum. Referans makale A356 alüminyum automotive subframe parçalarının su verme sırasında soğumasını inceliyor. Hızlı soğuma güçlü sıcaklık gradyanları oluşturuyor ve bu gradyanlar daha sonra residual stress ve distortion hesabı için kritik hale geliyor. Bu tezde aynı fiziksel/process değerlerini kullanarak farklı bir hesaplama sorusu sorduk: bazı FEM zaman adımlarını FEM-anchored NAS-PINN ile atlayabilir miyiz? Gerçek endüstriyel veri geç geldiği için yöntemi önce kontrollü surrogate geometrilerde, basitten daha karmaşığa doğru test ettik.")
    add_text(s, 0.75, 0.95, 11.9, 0.45, "The project begins with an industrial A356 aluminium quenching problem and asks whether FEM time-step calls can be reduced without losing useful temperature accuracy.", 16, True, PRIMARY, PP_ALIGN.CENTER)
    add_card(s, 0.65, 1.65, 3.55, 3.1, "Reference paper", "The reference study models water quenching of cast A356 automotive subframes. Its FEM simulation gives the industrial motivation: temperature gradients during quenching drive distortion.")
    add_card(s, 4.45, 1.65, 3.55, 3.1, "Our approach", "We keep the same physical process setting, but test a different route: FEM provides anchor snapshots and NAS-PINN predicts selected intermediate windows.", WARM_NOTE)
    add_card(s, 8.25, 1.65, 3.9, 3.1, "Validation path", "Because the real industrial data arrived late, the workflow starts with controlled surrogate geometries and then moves toward more complex 3D and Thermal Fin benchmarks.")
    add_metric(s, 1.0, 5.35, 1.85, 1.05, "Material", "A356", "aluminium")
    add_metric(s, 3.15, 5.35, 1.85, 1.05, "Initial T", "540 C", "hot part")
    add_metric(s, 5.3, 5.35, 1.85, 1.05, "Bath T", "20 C", "water")
    add_metric(s, 7.45, 5.35, 1.85, 1.05, "Quench", "30 s", "transient")
    add_metric(s, 9.6, 5.35, 1.85, 1.05, "Goal", "skip FEM", "with anchors")
    i += 1

    # 3 Context A356
    s = new_slide(prs, i, "Intro: From the Reference Paper to This Thesis", "EN: This slide connects the problem statement to the method. As shown in Figure 1.1, the MSWP framework has three operating modes: fixed-k, adaptive-k, and PINN-only. The goal is to reduce repeated FEM calls while still keeping FEM anchor points for correction. The numerical values at the bottom show the physical setup inherited from the quenching context: the part starts at 540 C, the water bath is 20 C, and the transient cooling lasts 30 seconds. For 2D benchmarks we use h = 5000 W/(m2K), and for 3D/Thermal Fin we use h = 4000 W/(m2K).\nTR: Bu slaytta problem statement ile yöntemi birbirine bağlıyor. Figure 1.1'de görüldüğü gibi MSWP framework üç çalışma moduna sahip: fixed-k, adaptive-k ve PINN-only. Amaç, FEM anchor noktalarını korurken tekrarlanan FEM çağrılarını azaltmak. Alttaki sayısal değerler quenching probleminden alınan fiziksel setup'ı gösteriyor: parça 540 C'den başlıyor, su banyosu 20 C ve transient cooling 30 saniye sürüyor. 2D benchmarklarda h = 5000 W/(m2K), 3D ve Thermal Fin için h = 4000 W/(m2K) kullanıyoruz.")
    add_text(s, 0.75, 0.95, 11.9, 0.48, "We follow the industrial quenching problem from the reference paper, then test whether FEM time-step calls can be reduced with a FEM-anchored NAS-PINN strategy.", 16, True, PRIMARY, PP_ALIGN.CENTER)
    add_image(s, "fig_mswp_three_modes.png", 0.65, 1.65, 5.4, 3.25)
    add_card(s, 6.35, 1.62, 5.95, 1.25, "Problem statement", "Figure 1.1 shows the three operating modes of the MSWP framework. Starting from the reference paper, the thesis keeps the same quenching setting but tests whether selected FEM steps can be skipped with NAS-PINN predictions.", WARM_NOTE)
    add_card(s, 6.35, 3.15, 5.95, 1.05, "Why surrogate geometries first?", "The real industrial data arrived late, so the method was validated first on controlled geometries: NAS-PINN source domains, canonical 2D domains, then 3D and Thermal Fin benchmarks.")
    add_metric(s, 0.75, 5.35, 2.05, 1.05, "Initial temperature", "540 C", "A356 part")
    add_metric(s, 3.0, 5.35, 2.05, 1.05, "Water bath", "20 C", "quench medium")
    add_metric(s, 5.25, 5.35, 2.05, 1.05, "Simulation time", "30 s", "transient cooling")
    add_metric(s, 7.5, 5.35, 2.05, 1.05, "2D h", "5000", "W/(m2K)")
    add_metric(s, 9.75, 5.35, 2.05, 1.05, "3D/Fin h", "4000", "W/(m2K)")
    i += 1

    # 4 Main source article
    s = new_slide(prs, i, "Main Source Article and Why It Was Chosen", "EN: This slide explains why the Mortensen et al. article is the main industrial reference. The paper is not chosen randomly. It belongs to the ongoing project context, and the supervisor asked us to investigate an alternative computational approach for the same type of FEM-based quenching model. Therefore, we use the article as the physical reference: A356 aluminium, water quenching, transient temperature field, and distortion relevance. What changes in this thesis is the computational strategy: instead of running FEM at every step, we test whether FEM anchors plus NAS-PINN predictions can reduce the number of solver calls.\nTR: Bu slaytta Mortensen et al. makalesinin neden ana endüstriyel referans olduğunu açıklıyorum. Bu makale rastgele seçilmedi; devam eden proje bağlamına ait ve danışman hocam bu FEM tabanlı quenching modeline alternatif bir hesaplama yaklaşımı araştırmamı istedi. Bu yüzden makaleyi fiziksel referans olarak kullanıyoruz: A356 alüminyum, water quenching, transient temperature field ve distortion ile ilişkisi. Bu tezde değişen şey hesaplama stratejisi: FEM'i her adımda çalıştırmak yerine, FEM anchor noktaları ve NAS-PINN tahminleriyle solver call sayısını azaltmayı test ediyoruz.")
    add_card(s, 0.65, 1.05, 5.25, 5.7, "Industrial source", "Mortensen et al. (2026) study water quenching of cast A356 automotive subframes with a validated FEM distortion simulation. This article defines the industrial problem that motivates the thesis.")
    add_card(s, 6.25, 1.05, 6.1, 2.55, "What this thesis inherits", "The thesis takes the material context and process scale from that work: A356 aluminium, T0 = 540 C, Tw = 20 C, a 30-second transient quench, and an engineering temperature-error target.")
    add_card(s, 6.25, 3.95, 6.1, 2.8, "Why this article?", "It is not a random source. The work is connected to an ongoing project, and the thesis studies an alternative FEM-acceleration route for the model discussed with the supervisor.", WARM_NOTE)
    i += 1

    # 5 Method source article
    s = new_slide(prs, i, "Method Source: NAS-PINN", "EN: This slide introduces the methodological source. Wang and Zhong's NAS-PINN paper provides the idea that the PINN architecture should not be selected manually; instead, Neural Architecture Search can choose depth, width, and activation for PDE solving. In this thesis, we adapt that idea to transient quenching. We use the NAS-PINN transfer architecture from the paper and compare it with our own NAS-found architectures: Bayesian/TPE, NSGA-II, and NSGA-III. The results section will show how these architectures behave across simple and complex geometries.\nTR: Bu slaytta yöntem kaynağını tanıtıyorum. Wang ve Zhong'un NAS-PINN makalesi, PINN mimarisinin manuel seçilmesi yerine Neural Architecture Search ile depth, width ve activation seçilmesi fikrini veriyor. Bu tezde bu fikri transient quenching problemine uyarlıyoruz. NAS-PINN makalesindeki transfer mimarisini kullanıyoruz ve bunu kendi NAS ile bulduğumuz Bayesian/TPE, NSGA-II ve NSGA-III mimarileriyle karşılaştırıyoruz. Results kısmında bu mimarilerin basit ve karmaşık geometrilerde nasıl davrandığını göstereceğiz.")
    add_card(s, 0.65, 1.05, 5.4, 5.7, "NAS-PINN source", "Wang and Zhong (2024) provide the methodological starting point: a PINN is not chosen manually, but selected through neural architecture search for solving PDEs.")
    add_card(s, 6.35, 1.05, 5.95, 2.65, "How this thesis adapts it", "The thesis keeps the NAS-PINN idea but changes the deployment setting: FEM anchors are retained, the network predicts intermediate time windows, and MSWP controls solver-call reduction.")
    add_image(s, "fig1_naspinn_2d.png", 6.35, 4.0, 5.95, 2.5)
    i += 1

    # 6 Research question
    s = new_slide(prs, i, "Research Question and Contributions", "EN: This slide states the central research question. We ask how many FEM solver calls can be replaced by PINN predictions while keeping the temperature error within useful engineering limits. I should explain the threshold carefully: Mortensen et al. do not give a universal PINN error threshold. Their work is an industrial FEM distortion study, validated with production-part measurements, where the temperature field is important because thermal gradients drive residual stress and final dimensional distortion. In our thesis, we only solve the thermal surrogate problem, so we translate that industrial tolerance idea into a temperature-field acceptance band. The full quench range is 540 C - 20 C = 520 C. A 15 C temperature error is about 2.9% of that range, which is close to a 3% engineering band. Therefore, 15 C is used for the harder 3D and Thermal Fin cases, where the geometry is closer to the target industrial complexity. For simpler canonical 2D domains, we use a stricter 10 C band because they are controlled, lower-dimensional validation cases. So the threshold is not arbitrary: it is a practical bridge from Mortensen's distortion-tolerance context to our thermal prediction problem, adjusted by geometry complexity.\nTR: Bu slaytta temel araştırma sorusunu veriyorum. Hata threshold kısmını özellikle iyi açıklamalıyım: Mortensen et al. doğrudan evrensel bir PINN hata limiti vermiyor. Onların çalışması üretim parçası ölçümleriyle doğrulanan endüstriyel bir FEM distortion çalışmasıdır. Burada sıcaklık alanı önemlidir çünkü thermal gradients residual stress ve final dimensional distortion'ı belirler. Bizim tezimiz ise mechanical distortion'ı değil, önce thermal surrogate problemi çözüyor. Bu yüzden Mortensen'deki endüstriyel tolerans fikrini sıcaklık alanı için bir engineering acceptance band olarak yorumladık. Toplam quench aralığı 540 C - 20 C = 520 C. 15 C hata bu aralığın yaklaşık 2.9%'udur, yani yaklaşık 3% engineering band'a karşılık gelir. Bu nedenle daha zor ve endüstriyel probleme daha yakın olan 3D ve Thermal Fin için 15 C kullandık. Daha basit canonical 2D domainler için ise 10 C daha sıkı bir band olarak kullanıldı, çünkü bunlar kontrollü ve düşük boyutlu validation case'lerdir. Yani threshold keyfi değil; Mortensen'deki distortion-tolerance bağlamından bizim thermal prediction problemimize kurulan pratik bir köprüdür ve geometri karmaşıklığına göre ayarlanmıştır.")
    add_text(s, 0.9, 1.1, 11.6, 0.9, "How many FEM solver calls can be replaced by PINN predictions while keeping temperature error within engineering limits?", 24, True, PRIMARY, PP_ALIGN.CENTER)
    add_card(s, 0.75, 2.55, 3.05, 1.85, "We built", "A FEM-anchored NAS-PINN framework that predicts intermediate thermal fields while FEM snapshots remain the correction source.")
    add_card(s, 3.95, 2.55, 3.05, 1.85, "We tested", "Fixed-k, adaptive-k and PINN-only modes were evaluated from simple 2D domains to complex 3D and Thermal Fin cases.")
    add_card(s, 7.15, 2.55, 3.05, 1.85, "We measured", "MAE, relative L2 error, skip-factor schedule and FEM-call reduction are reported together.")
    add_card(s, 10.35, 2.55, 2.55, 1.85, "Error limits", "Geometry-dependent thresholds: 10 C for canonical 2D; 15 C for 3D and Thermal Fin.", WARM_NOTE)
    add_card(s, 1.05, 5.0, 11.3, 1.25, "Engineering tolerance calculation", "Delta T = T0 - Tw = 540 C - 20 C = 520 C     |     15 C / 520 C = 0.0288 = 2.9% approx. 3%     |     10 C / 520 C = 1.9%", COOL_PANEL)
    i += 1

    # 8 Theory
    section_slide(prs, i, "Chapter 2: Theoretical Background", "Heat transfer, FEM, PINN, NAS-PINN and optimisation", "EN: This section gives the theoretical background needed to understand the method. I will first explain the heat transfer model, then FEM and PINNs, then why NAS and optimisation are used for architecture selection. The goal is not to present theory separately from the thesis, but to show which theoretical tools are used in the framework.\nTR: Bu bölüm yöntemi anlamak için gereken teorik altyapıyı veriyor. Önce heat transfer modelini, sonra FEM ve PINN kavramlarını, ardından NAS ve optimizasyonların mimari seçiminde neden kullanıldığını anlatacağım. Amaç teoriyi ayrı bir bölüm gibi sunmak değil; bu framework içinde hangi teorik araçların kullanıldığını göstermektir.")
    i += 1

    s = new_slide(prs, i, "Heat Transfer Model", "EN: This slide explains the physical equation behind the thesis. We solve transient heat transfer: the temperature field changes over time while the part cools from the hot initial state toward the water bath temperature. The boundary condition is convective cooling, represented by h(T - Tw). This equation is used twice: FEM solves it numerically to create reference snapshots, and PINN uses it inside the loss function as a physics constraint.\nTR: Bu slaytta tezin fiziksel denklemini açıklıyor. Geçici heat transfer çözüyoruz: parça sıcak başlangıç durumundan su banyosu sıcaklığına doğru soğurken sıcaklık alanı zamanla değişiyor. Boundary condition convective cooling olarak yazılıyor ve h(T - Tw) ile temsil ediliyor. Bu denklem iki yerde kullanılıyor: FEM bu denklemi sayısal olarak çözüp reference snapshot üretiyor; PINN ise aynı denklemi loss function içinde physics constraint olarak kullanıyor.")
    add_text(s, 0.75, 1.05, 12.0, 0.45, "The whole framework is built on transient heat transfer during water quenching.", 17, True, PRIMARY, PP_ALIGN.CENTER)
    add_card(s, 0.75, 1.8, 4.0, 2.45, "Governing equation", "rho cp dT/dt = div(kT grad T)\n\nThis describes how the temperature field evolves inside the solid part.", WARM_NOTE)
    add_card(s, 4.95, 1.8, 4.0, 2.45, "Boundary condition", "-kT grad(T) . n = h (T - Tw)\n\nThis represents heat leaving the hot metal surface into the quench water.")
    add_card(s, 9.15, 1.8, 3.65, 2.45, "Initial condition", "T(x,0) = T0\n\nIn this thesis: T0 = 540 C and Tw = 20 C.")
    add_card(s, 1.0, 5.0, 5.85, 1.25, "How FEM uses it", "FEM discretises the domain and solves this equation to produce reference temperature snapshots.")
    add_card(s, 7.05, 5.0, 5.35, 1.25, "How PINN uses it", "PINN predicts T(x,t) and penalises violations of the PDE, boundary condition and initial condition.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "FEM, PINN and NAS-PINN", "EN: On this slide I explain three important words. FEM is the numerical reference solver. It solves the heat equation on a mesh and gives trusted temperature snapshots. PINN is a neural network, but it is not only fitted to data. It also uses the heat equation in the loss. NAS-PINN means that we do not choose the neural network by hand. We search for the best architecture. In this thesis, FEM stays as the anchor, PINN predicts the missing windows, and NAS helps choose a better PINN.\nTR: Bu slaytta üç önemli kavramı basit anlatıyorum. FEM sayısal reference solver'dır. Heat equation'ı mesh üzerinde çözer ve güvenilir temperature snapshots veriyorum. PINN bir neural network'tür ama sadece data fit etmez; heat equation'ı da loss içinde kullanır. NAS-PINN ise neural network mimarisini elle seçmemek demektir; iyi mimariyi search ile buluruz. Bu tezde FEM anchor olarak kalır, PINN aradaki windowları tahmin eder, NAS ise daha iyi PINN seçmeye yardım ediyor.")
    add_card(s, 0.65, 1.0, 3.9, 4.55, "FEM", "What is it?\nA numerical solver for the heat equation.\n\nHow we used it?\nIt gives reliable anchor snapshots.\n\nKey point: FEM is not removed.")
    add_card(s, 4.75, 1.0, 3.9, 4.55, "PINN", "What is it?\nA neural network for T(x,t).\n\nHow it learns?\nIt uses data error + heat-equation error.\n\nKey point: physics is inside the loss.", WARM_NOTE)
    add_card(s, 8.85, 1.0, 3.75, 4.55, "NAS-PINN", "What is it?\nPINN + architecture search.\n\nWhat is searched?\nDepth, width and activation.\n\nKey point: we do not choose the network manually.")
    add_card(s, 0.9, 5.85, 11.6, 0.75, "Defense message", "We used FEM for trust, PINN for prediction, and NAS to choose a better neural architecture for different geometries.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Optimisation Methods Used", "EN: This slide explains the NAS optimisation methods. Bayesian/TPE is a sample-efficient single-objective search, mainly seeking low MAE. NSGA-II and NSGA-III are multi-objective evolutionary methods: they look at accuracy together with model compactness and diversity. We used all three because the best architecture is not only the most accurate one; it must also remain stable when the geometry becomes more complex and the prediction window becomes longer.\nTR: Bu slaytta NAS optimizasyon yöntemlerini açıklıyor. Bayesian/TPE sample-efficient single-objective search yöntemidir ve temel olarak düşük MAE arar. NSGA-II ve NSGA-III multi-objective evolutionary yöntemlerdir; accuracy ile birlikte model compactness ve solution diversity'yi de dikkate alır. Üçünü de kullanmamızın nedeni şu: en iyi mimari sadece en düşük hatayı veren mimari değildir; geometri karmaşıklaştığında ve prediction window uzadığında da stabil kalmalıdır.")
    add_table(s, 0.55, 1.0, 12.2, 5.85, ["Method", "Search idea", "Why used here", "Architecture"], [
        ["Bayesian/TPE", "Single-objective probabilistic search", "Find an accurate architecture with limited trials", "5 x 151, ReLU"],
        ["NSGA-II", "Multi-objective evolutionary search", "Balance accuracy and compactness", "3 x 153, tanh"],
        ["NSGA-III", "Reference-point many-objective search", "Promote diverse compact Pareto solutions", "3 x 75, tanh"],
    ], size=10)
    add_card(s, 0.75, 6.08, 11.85, 0.7, "Key message", "The optimisation methods were not random choices; they test different architecture-search philosophies under the same FEM-anchored quenching framework.", WARM_NOTE)
    i += 1

    # Method
    section_slide(prs, i, "Chapter 3: Methodology", "Reference problem -> MSWP -> NAS-PINN -> domains -> metrics", "EN: Here, I follow the report order. First, I define the reference industrial problem from Mortensen et al. and what this thesis takes or simplifies. Then I explain the MSWP framework: FEM reference solver, skip factor, and operating modes. After that, I explain the NAS-PINN formulation, architecture search, benchmark domains, and evaluation metrics.\nTR: Bu bölüm rapordaki sırayı takip ediyorum. Önce Mortensen et al. referans endüstriyel problemi ve bu tezde neleri aldığımız/neleri sadeleştirdiğimiz anlatılıyor. Sonra MSWP framework geliyor: FEM reference solver, skip factor ve operating modes. Daha sonra NAS-PINN formulation, architecture search, benchmark domainler ve evaluation metrics anlatılıyor.")
    i += 1

    s = new_slide(prs, i, "Reference Problem and Thesis Simplification", "EN: This slide corresponds to the first Methodology section in the report. Mortensen et al. solve the full industrial A356 quenching and distortion problem with a thermo-elastoviscoplastic FEM model. In this thesis, we focus on the thermal sub-problem: predicting the transient temperature field. We keep the material and process setting, but simplify the industrial model so the FEM-anchored NAS-PINN method can be validated systematically before moving to the real mesh.\nTR: Bu slaytta rapordaki Methodology bölümünün ilk kısmına karşılık geliyor. Mortensen et al. tam endüstriyel A356 quenching ve distortion problemini thermo-elastoviscoplastic FEM ile çözüyor. Bu tezde biz thermal sub-problem'e odaklanıyoruz: transient temperature field tahmini. Material ve process setting korunuyor, fakat gerçek endüstriyel model sadeleştiriliyor; böylece FEM-anchored NAS-PINN yöntemi gerçek mesh'e geçmeden önce sistematik olarak doğrulanabiliyor.")
    add_card(s, 0.65, 1.05, 3.85, 5.55, "Mortensen et al.", "Full industrial simulation: A356 subframe, coupled thermal-mechanical FEM, temperature-dependent boiling curve, industrial validation with CMM measurements.")
    add_card(s, 4.75, 1.05, 3.85, 5.55, "This thesis", "Thermal-only surrogate: transient heat equation, constant material properties, controlled 2D/3D/Thermal Fin geometries, comparison to own FEM reference.", WARM_NOTE)
    add_card(s, 8.85, 1.05, 3.65, 5.55, "Why this step matters", "The NAS-PINN must first prove that it can reproduce the thermal field before it can be connected to a full industrial distortion pipeline.")
    i += 1

    s = new_slide(prs, i, "MSWP Framework: Working Logic", "EN: In this slide, I explain the MSWP working logic step by step. Step 1: FEM creates anchor temperature fields. Step 2: the framework chooses a skip factor k. Step 3: NAS-PINN predicts the skipped windows instead of calling FEM every time. Step 4: a later FEM anchor corrects the trajectory. My main point is that MSWP reduces FEM calls, but it does not remove FEM.\nTR: Bu slaytta MSWP çalışma mantığını adım adım anlatıyorum. Step 1: FEM anchor temperature field oluşturuyor. Step 2: Framework skip factor k seçiyor. Step 3: NAS-PINN, FEM'i her adımda çağırmak yerine skipped window'ları tahmin ediyor. Step 4: Daha sonraki FEM anchor trajectory'yi düzeltiyor. Benim ana mesajım: MSWP FEM call sayısını azaltır, ama FEM'i tamamen kaldırmaz.")
    add_image(s, "fig6_naspinn_process.png", 0.55, 1.0, 6.6, 5.65)
    add_card(s, 7.45, 1.0, 4.8, 1.15, "Step 1: FEM anchors", "FEM gives trusted reference states at selected times.")
    add_card(s, 7.45, 2.35, 4.8, 1.15, "Step 2: choose k", "The skip factor decides how many FEM calls are skipped.", WARM_NOTE)
    add_card(s, 7.45, 3.7, 4.8, 1.15, "Step 3: PINN predicts", "NAS-PINN fills the skipped thermal windows.")
    add_card(s, 7.45, 5.05, 4.8, 1.15, "Step 4: correct", "The next FEM anchor limits drift and resets the field.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "FEM Process: How It Works", "EN: In this slide, I explain the general FEM process step by step. Step 1: I define the physical domain, material properties, initial condition and boundary conditions. Step 2: I discretise the domain into grid or mesh nodes. Step 3: the heat equation is solved in time with the numerical solver. Step 4: the solver returns the temperature field at each saved time. Here, I use this FEM result as the trusted reference for the rest of the thesis.\nTR: Bu slaytta genel FEM prosesinin nasıl çalıştığını adım adım anlatıyorum. Step 1: Physical domain, material properties, initial condition ve boundary conditions tanımlanır. Step 2: Domain grid veya mesh node'larına ayrılır. Step 3: Heat equation zaman içinde numerical solver ile çözülür. Step 4: Solver her kayıt zamanında temperature field verir. Bu tezde bu FEM sonucu trusted reference olarak kullanıyorum.")
    add_card(s, 0.55, 1.0, 3.65, 1.25, "Step 1: physics", "Define geometry, material properties, initial condition and boundary conditions.")
    add_card(s, 0.55, 2.55, 3.65, 1.25, "Step 2: discretise", "Split the domain into grid or mesh nodes.", WARM_NOTE)
    add_card(s, 0.55, 4.1, 3.65, 1.25, "Step 3: solve", "Advance the heat equation in time with the numerical solver.")
    add_card(s, 0.55, 5.65, 3.65, 0.75, "Step 4: output", "Save temperature fields as reference snapshots.", COOL_PANEL)
    add_image(s, "FEM_Proses.png", 4.55, 1.0, 8.0, 5.75)
    i += 1

    s = new_slide(prs, i, "FEM Solver and Mesh Configuration", "EN: In this slide, I show Appendix Table B.1 from the report. This table explains the FEM reference grids and active nodes used for each benchmark domain. I use it here because after explaining the FEM process, I need to show what the solver actually runs on. My main point is that FEM is the trusted reference, and these grids define the numerical resolution behind that reference.\nTR: Bu slaytta rapordaki Appendix Table B.1'i gösteriyorum. Bu tablo her benchmark domain için FEM reference grid ve active node sayılarını açıklar. FEM prosesini anlattıktan sonra solver'ın gerçekten hangi grid üzerinde çalıştığını göstermek için buraya koyuyorum. Benim ana mesajım: FEM trusted reference'tır ve bu gridler o reference'ın numerical resolution kısmını tanımlar.")
    add_image(s, "Table B.1.png", 0.75, 1.0, 7.2, 5.7)
    add_card(s, 8.25, 1.2, 4.3, 1.35, "Step after FEM process", "After defining the FEM workflow, I show the actual grids and active nodes.", WARM_NOTE)
    add_card(s, 8.25, 2.85, 4.3, 1.35, "Why it matters", "The reference error depends on the numerical grid quality and domain resolution.")
    add_card(s, 8.25, 4.5, 4.3, 1.35, "How I use it", "These FEM fields become the baseline for MAE, L2 and anchor snapshots.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "FEM vs PINN: Runtime Logic", "EN: In this slide, I explain the cost logic clearly. A single FEM run can be fast for a small benchmark, but FEM must solve the numerical system again when I change the geometry, boundary condition, material setting, time window or design case. So repeated FEM use can become expensive. A PINN can be slower at the first training stage, because it has to learn the physics and the field. But after training, prediction is only a neural-network inference step, so repeated evaluation is much faster. This is why my framework is useful when I need many repeated thermal predictions, not only one isolated solve.\nTR: Bu slaytta cost mantığını net açıklıyorum. Küçük bir benchmark için tek FEM run hızlı görünebilir. Ama geometri, boundary condition, material setting, time window veya design case değişirse FEM numerical system'i tekrar çözmek zorundadır. Bu yüzden tekrar tekrar FEM kullanımı maliyetli olabilir. PINN ilk training aşamasında FEM'den uzun sürebilir, çünkü physics ve field öğrenilir. Ama training bittikten sonra prediction sadece neural-network inference olur ve tekrar kullanım çok daha hızlıdır. Bu yüzden framework özellikle çok sayıda tekrar thermal prediction gerektiğinde faydalıdır.")
    add_table(s, 0.65, 1.05, 12.0, 3.1, ["Method", "First use", "Repeated use", "Why it matters"], [
        ["FEM", "Can be fast for one small case", "Must solve again for each new case", "Cost grows with many runs"],
        ["PINN", "Training can be slower", "Inference is very fast", "Good for repeated prediction"],
        ["MSWP", "Uses FEM anchors", "Skips selected FEM calls", "Keeps trust while reducing cost"],
    ], size=9)
    add_card(s, 0.8, 4.55, 3.7, 1.55, "My simple message", "FEM is reliable, but repeated solves can be costly.", WARM_NOTE)
    add_card(s, 4.75, 4.55, 3.7, 1.55, "PINN trade-off", "PINN pays a training cost first, then reuses the trained model quickly.")
    add_card(s, 8.7, 4.55, 3.55, 1.55, "Why MSWP", "I keep FEM anchors, but replace selected intermediate solves with PINN predictions.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Skip Factor and Operating Modes", "EN: This slide explains how FEM calls are reduced. The skip factor k defines how many 1.5 s anchor intervals one PINN window covers. At k=1, there is no FEM saving because FEM is called at every anchor interval. At k=5, the PINN covers five intervals per FEM call, replacing four out of five FEM calls, which corresponds to 80% saving. The report evaluates fixed-k, adaptive-k and PINN-only to test both safe acceleration and the failure mode without anchoring.\nTR: Bu slaytta FEM çağrılarının nasıl azaltıldığını anlatıyor. Skip factor k, bir PINN window'un kaç tane 1.5 s anchor interval kapsadığını belirler. k=1 iken FEM her anchor intervalda çağrılır, yani saving yoktur. k=5 iken PINN bir FEM çağrısı başına beş interval kapsar; beş FEM çağrısından dördü atlanır ve bu 80% saving'e karşılık gelir. Raporda fixed-k, adaptive-k ve PINN-only test edilerek hem güvenli acceleration hem de anchoring kaldırıldığında oluşan failure mode inceleniyor.")
    add_image(s, "fig_mswp_three_modes.png", 0.65, 1.0, 5.8, 5.65)
    add_card(s, 6.8, 1.0, 5.6, 1.25, "Formula", "Delta t_k = k * Delta t_a, with Delta t_a = 1.5 s")
    add_table(s, 6.8, 2.55, 5.6, 2.45, ["k", "Windows", "FEM saving"], [["1", "20", "0%"], ["2", "10", "50%"], ["3", "7", "65%"], ["4", "5", "75%"], ["5", "4", "80%"]], size=10)
    add_card(s, 6.8, 5.35, 5.6, 0.95, "Operating modes", "Fixed-k tests constant windows, adaptive-k changes k during the transient, and PINN-only removes FEM correction after the first state.", WARM_NOTE)
    i += 1

    s = new_slide(prs, i, "NAS-PINN Process", "EN: In this slide, I explain the NAS-PINN process step by step. Step 1: I start from one FEM anchor field at the beginning of a window. Step 2: the PINN predicts the temperature field inside that window. Step 3: the physics-informed loss checks the heat equation, boundary behaviour and endpoint supervision. Step 4: NAS searches different neural architectures and selects the architecture that gives the best accuracy-efficiency trade-off.\nTR: Bu slaytta NAS-PINN prosesini adım adım anlatıyorum. Step 1: Bir window başında FEM anchor field ile başlıyorum. Step 2: PINN o window içinde temperature field tahmin ediyor. Step 3: Physics-informed loss heat equation, boundary behaviour ve endpoint supervision'ı kontrol ediyor. Step 4: NAS farklı neural architecture'ları arıyor ve en iyi accuracy-efficiency trade-off veren architecture'ı seçiyor.")
    add_image(s, "Single-Window FEM-Anchored NAS-PINN Framework.png", 0.55, 1.0, 6.65, 5.65)
    add_card(s, 7.45, 1.0, 4.8, 1.15, "Step 1: anchor", "Use the FEM field at the start of one window.")
    add_card(s, 7.45, 2.35, 4.8, 1.15, "Step 2: prediction", "PINN predicts T(x,t) inside this window.", WARM_NOTE)
    add_card(s, 7.45, 3.7, 4.8, 1.15, "Step 3: loss", "Physics and endpoint terms control the prediction.")
    add_card(s, 7.45, 5.05, 4.8, 1.15, "Step 4: NAS", "Architecture search chooses the best network design.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Architecture Search and Training Protocol", "EN: This slide belongs to the NAS-PINN Architecture Search part of the report. We search architectures once on the 2D rectangle validation setting, then transfer the selected architectures to all benchmark domains and skip factors. This reflects the practical aim: we do not want to rerun NAS for every geometry. Bayesian/TPE, NSGA-II and NSGA-III produce different architecture profiles, and training uses Adam followed by L-BFGS for each prediction window.\nTR: Bu slaytta rapordaki NAS-PINN Architecture Search kısmına karşılık geliyor. Architecture search 2D rectangle validation setting üzerinde bir kez yapılıyor; sonra seçilen mimariler tüm benchmark domainlere ve skip factor değerlerine transfer ediliyor. Bu pratik hedefi yansıtıyor: her geometri için NAS'i yeniden çalıştırmak istemiyoruz. Bayesian/TPE, NSGA-II ve NSGA-III farklı architecture profile üretiyor; her prediction window için training Adam ardından L-BFGS ile yapılıyor.")
    add_table(s, 0.65, 1.0, 12.0, 3.25, ["Strategy", "Selected architecture", "Role in method"], [["Bayesian/TPE", "5 layers x 151, ReLU", "High-capacity accuracy-focused model"], ["NSGA-II", "3 layers x 153, tanh", "Accuracy-size trade-off"], ["NSGA-III", "3 layers x 75, tanh", "Compact Pareto-diverse model"]], size=11)
    add_card(s, 0.9, 4.65, 3.75, 1.5, "Search setting", "Search once on the 2D rectangle validation domain, then transfer architectures.")
    add_card(s, 4.9, 4.65, 3.75, 1.5, "Training", "Each window is trained with Adam followed by L-BFGS.", WARM_NOTE)
    add_card(s, 8.9, 4.65, 3.55, 1.5, "Why transfer?", "Industrial use requires architecture robustness across geometries.")
    i += 1

    s = new_slide(prs, i, "NAS Hyperparameter Search Space", "EN: In this slide, I show Appendix Table A.1 from the report. After explaining the architecture search protocol, I show the actual search space: number of layers, neurons, activation functions and Fourier embedding. This is important because Bayesian/TPE, NSGA-II and NSGA-III all search the same space; only the search strategy changes.\nTR: Bu slaytta rapordaki Appendix Table A.1'i gösteriyorum. Architecture search protocol'ü anlattıktan sonra gerçek search space'i gösteriyorum: layer sayısı, neuron sayısı, activation functions ve Fourier embedding. Bu önemlidir çünkü Bayesian/TPE, NSGA-II ve NSGA-III aynı search space'i arar; sadece search strategy değişir.")
    add_image(s, "Table A.1.png", 1.0, 1.15, 6.6, 3.6)
    add_card(s, 8.0, 1.2, 4.35, 1.25, "Same search space", "All NAS strategies use the same candidate choices.", WARM_NOTE)
    add_card(s, 8.0, 2.8, 4.35, 1.25, "What changes?", "Bayesian/TPE, NSGA-II and NSGA-III search this space differently.")
    add_card(s, 8.0, 4.4, 4.35, 1.25, "Why include it?", "It makes the architecture comparison fair and reproducible.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Benchmark Domains: 2D Domain Groups", "EN: This slide follows the Benchmark Domains section of the report. There are two different 2D groups, and they should not be mixed. Level 1 is the NAS transfer and adaptive-k study: Square, Circle, L-Shape and Flower. Level 2 is the canonical 2D fixed-k sweep: Rectangle, Circle and L-Shape. Circle and L-Shape are shared; Square and Flower appear only in the adaptive-k study.\nTR: Bu slaytta rapordaki Benchmark Domains bölümünü takip ediyorum. Burada iki farklı 2D grup var ve karıştırılmamalı. Level 1, NAS transfer ve adaptive-k çalışmasıdır: Square, Circle, L-Shape ve Flower. Level 2, canonical 2D fixed-k sweep çalışmasıdır: Rectangle, Circle ve L-Shape. Circle ve L-Shape iki grupta ortak; Square ve Flower sadece adaptive-k çalışmasında yer alır.")
    add_image(s, "fig1_naspinn_2d.png", 0.55, 1.0, 5.75, 2.65)
    add_image(s, "fig2_thesis_2d.png", 6.75, 1.0, 5.95, 2.65)
    add_card(s, 0.55, 4.0, 5.75, 2.15, "Level 1: NAS transfer + adaptive-k", "Domains: Square, Circle, L-Shape, Flower. Purpose: test NAS-PINN transfer architecture and our NAS-found architectures under adaptive-k. Values: T0=540 C, Tw=20 C, h=5000, kT=160.", WARM_NOTE)
    add_card(s, 6.75, 4.0, 5.95, 2.15, "Level 2: canonical 2D fixed-k sweep", "Domains: Rectangle, Circle, L-Shape. Purpose: compare Bayesian/TPE, NSGA-II and NSGA-III over k=1...5. Values: T0=540 C, Tw=20 C, h=5000, kT=150.")
    i += 1

    s = new_slide(prs, i, "Benchmark Domains: 3D and Thermal Fin", "EN: This slide continues the Benchmark Domains section. Level 3 contains the canonical 3D domains: Rectangular Prism, Cylinder, Stacked Cubes and L-Shape Prism. These test volume, curvature, interfaces and reentrant edges. Level 4 is Thermal Fin, kept separate because it is the hardest benchmark with sharp fin gradients and a more demanding 3D transient adaptation.\nTR: Bu slaytta Benchmark Domains bölümünün devamıdır. Level 3 canonical 3D domainleri içerir: Rectangular Prism, Cylinder, Stacked Cubes ve L-Shape Prism. Bunlar volume, curvature, interface ve reentrant edge etkilerini test eder. Level 4 Thermal Fin'dir; ayrı tutulur çünkü sharp fin gradients içeren ve 3D transient'e uyarlanmış en zor benchmarktır.")
    add_image(s, "fig3_thesis_3d.png", 0.75, 1.0, 5.95, 3.0)
    add_image(s, "fig4_thermalfin.png", 7.0, 1.0, 5.6, 3.0)
    add_card(s, 0.9, 4.45, 5.55, 1.75, "Level 3: canonical 3D", "Rectangular, Cylinder, Stacked Cubes and L-Shape Prism test whether the method remains stable when geometric complexity increases.")
    add_card(s, 6.85, 4.45, 5.55, 1.75, "Level 4: Thermal Fin", "Thermal Fin is the most demanding benchmark; it tests anchoring, architecture capacity and Fourier features under sharp local gradients.", WARM_NOTE)
    i += 1

    s = new_slide(prs, i, "Evaluation Metrics", "EN: This slide closes the Methodology chapter. I evaluate both accuracy and computational saving. Mean-window MAE gives the average Celsius error, and relative L2 checks the spatial field. I also report FEM saving because runtime is part of the research question. FEM is reliable, but repeated FEM calls across many windows or design cases can become expensive. So I always connect accuracy with how many FEM calls are avoided.\nTR: Bu slaytta Methodology chapter'ı kapatıyorum. Hem accuracy hem computational saving değerlendiriyorum. Mean-window MAE average Celsius error verir, relative L2 spatial field'ı kontrol eder. FEM saving de raporlanır çünkü runtime research question'ın parçasıdır. FEM güvenilirdir, ama birçok window veya design case için tekrar tekrar çalıştırıldığında maliyetli olabilir. Bu yüzden accuracy ile kaç FEM call atlandığını birlikte anlatıyorum.")
    add_card(s, 0.75, 1.15, 3.8, 4.95, "MAE", "Mean absolute temperature error in degrees Celsius. Primary metric for engineering interpretation.", WARM_NOTE)
    add_card(s, 4.85, 1.15, 3.8, 4.95, "Relative L2", "Normalised field-level error. Useful for comparing spatial prediction quality across windows and domains.")
    add_card(s, 8.95, 1.15, 3.55, 4.95, "FEM saving", "Percentage of FEM calls replaced by PINN windows. Reported together with accuracy to show the trade-off.")
    i += 1

    # Results
    section_slide(prs, i, "Chapter 4: Results", "Report order: 2D adaptive-k -> canonical 2D -> 3D -> Thermal Fin", "EN: Here, I follow the exact order of the report. For each benchmark, the presentation now shows the main MAE table first, then the L2 error table, and only after that the figures and heat maps. I do not compress many report figures into one slide; each important table or figure is given enough space to be readable.\nTR: Burada rapordaki sırayı takip ediyorum. Her benchmark için önce ana MAE tablosu, sonra L2 error tablosu, ardından figürler ve heat map'ler gösteriyorum. Çok sayıda rapor figürünü tek slayta sıkıştırmıyorum; her önemli tablo veya figür okunabilir olacak kadar büyük veriyorum.")
    i += 1

    # 2D adaptive-k
    s = new_slide(prs, i, "2D Adaptive-k: Table 4.1 Setup", "EN: In this slide, I show Table 4.1, the physical and training setup for the 2D adaptive-k experiments. This benchmark group uses Square, Circle, L-Shape and Flower domains. It is the NAS architecture transfer and adaptive-k test group. The setup table is shown separately so the parameter choices are readable before discussing the numerical results.\nTR: Bu slaytta 2D adaptive-k deneyleri için Table 4.1'i gösteriyorum. Bu benchmark grubu Square, Circle, L-Shape ve Flower domainlerini içerir. Bu grup NAS architecture transfer ve adaptive-k testidir. Setup tablosunu ayrı veriyoruz ki parameter choices okunabilir olsun, sonra numerical results'a geçiyoruz.")
    add_image(s, "4.1.png", 2.15, 1.0, 8.9, 4.8)
    add_card(s, 1.0, 6.0, 11.35, 0.75, "Why this matters", "These parameters define the controlled 2D adaptive-k benchmark before comparing architectures and L2 field errors.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "2D Adaptive-k: MAE and L2 Summary", "EN: In this slide, I read Table 4.2 together with Appendix Table D.5. The left table gives the main result: mean-window MAE, FEM saving and k-schedule. The right table checks L2 field error. MAE is averaged over the prediction windows. L2 is the field-level error at the reported final snapshot, t=30 s. Together, I show that the method is accurate on average and also keeps the spatial field close to FEM.\nTR: Bu slaytta Table 4.2 ile Appendix Table D.5'i birlikte okuyorum. Soldaki tabloda main result'ı anlatıyorum: mean-window MAE, FEM saving ve k-schedule. Sağdaki tabloda L2 field error'ı kontrol ediyorum. MAE prediction window'lar üzerinden ortalamadır. L2 ise t=30 s final snapshot'taki field-level error'dır. Birlikte, method'un hem average olarak hem spatial field olarak FEM'e yakın olduğunu gösteriyorum.")
    add_text(s, 0.65, 0.86, 5.8, 0.24, "Table 4.2: mean-window MAE + FEM saving", 10, True, PRIMARY)
    add_text(s, 6.75, 0.86, 5.5, 0.24, "Appendix D.5: relative L2 at t=30 s", 10, True, PRIMARY)
    add_image(s, "4.2.png", 0.45, 1.15, 6.05, 4.25)
    add_image(s, "Table D.5.png", 6.65, 1.15, 5.55, 4.25)
    add_card(s, 0.65, 5.65, 3.65, 0.95, "MAE", "Average Celsius error over windows. It answers: how large is the temperature error on average?", WARM_NOTE)
    add_card(s, 4.55, 5.65, 3.65, 0.95, "L2", "Relative field error at t=30 s. It answers: is the full spatial field close to FEM?")
    add_card(s, 8.45, 5.65, 3.45, 0.95, "Key result", "Low MAE and low L2: adaptive-k saves FEM calls without losing field quality.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "2D Adaptive-k: Figure 4.1 FEM-PINN Curves", "EN: This slide connects the four adaptive-k domains. I found that PINN follows the FEM mean temperature curve closely in all domains. The first windows have higher error because cooling is fastest at the start. Square and Circle are easier. L-Shape and Flower are harder because their boundaries create local gradients. This supports the table result: adaptive-k gives high FEM saving while keeping the field close to FEM.\nTR: Bu slaytta dört adaptive-k domaini birlikte anlatır. PINN'in tüm domainlerde FEM mean temperature curve'ü yakından takip ettiğini bulduk. İlk windowlarda hata daha yüksek çünkü soğuma başlangıçta en hızlıdır. Square ve Circle daha kolaydır. L-Shape ve Flower daha zordur çünkü boundary yapıları lokal gradient oluşturur. Bu, tablo sonucunu destekler: adaptive-k yüksek FEM saving verirken field FEM'e yakın kalıyor.")
    add_image(s, "fem_pinn_adaptive_k_square.png", 0.55, 0.95, 3.0, 2.15)
    add_image(s, "fem_pinn_adaptive_k_circle.png", 3.75, 0.95, 3.0, 2.15)
    add_image(s, "fem_pinn_adaptive_k_lshape.png", 6.95, 0.95, 3.0, 2.15)
    add_image(s, "fem_pinn_adaptive_k_flower.png", 10.15, 0.95, 2.75, 2.15)
    add_card(s, 0.8, 3.65, 3.9, 1.5, "What I found", "FEM and PINN curves are close. Error is highest in early windows.", WARM_NOTE)
    add_card(s, 4.95, 3.65, 3.9, 1.5, "Why it happens", "Early cooling has steep gradients. Complex boundaries make local errors harder.")
    add_card(s, 9.1, 3.65, 3.05, 1.5, "Result link", "I use this to explain why adaptive-k works and why L2 stays low.", COOL_PANEL)
    i += 1

    # canonical 2D
    s = new_slide(prs, i, "Canonical 2D: Setup and Adaptive-k Schedule", "EN: In this slide, I place the canonical 2D setup next to the new adaptive-k schedule table. The setup tells what was tested. The schedule table tells how k changes when the controller is used. This makes the result section easier to read before the MAE and L2 comparison.\nTR: Bu slaytta canonical 2D setup tablosunu yeni adaptive-k schedule tablosu ile yan yana koyar. Setup neyin test edildiğini söyler. Schedule tablosu controller kullanıldığında k'nin nasıl değiştiğini söyler. Böylece MAE ve L2 karşılaştırmasından önce result section daha kolay okunur.")
    add_text(s, 0.85, 0.86, 4.7, 0.24, "Table 4.3: canonical 2D setup", 10, True, PRIMARY)
    add_text(s, 5.95, 0.86, 6.1, 0.24, "New schedule table: adaptive-k behaviour", 10, True, PRIMARY)
    add_image(s, "4.3.png", 0.75, 1.15, 4.8, 3.75)
    add_image(s, "adaptive_k_2d_canonical_schedules_table.png", 5.85, 1.15, 6.2, 3.75)
    add_card(s, 1.0, 5.35, 11.35, 0.9, "Why this helps", "The setup gives the experiment rules. The schedule table shows when FEM calls are skipped. Read them together before judging MAE and L2.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Canonical 2D: MAE and L2 Summary", "EN: In this slide, I read Table 4.4 together with Appendix Table D.13. The left table reports mean-window MAE for k=1 to k=5. The right table reports relative L2 error at t=30 s. MAE shows average temperature accuracy, while L2 checks the spatial field shape. Together, I show that Bayesian/TPE is the most stable option as k increases.\nTR: Bu slaytta Table 4.4 ile Appendix Table D.13'ü birlikte okuyorum. Soldaki tablo k=1'den k=5'e mean-window MAE veriyorum. Sağdaki tablo t=30 s'deki relative L2 error veriyorum. MAE ile average temperature accuracy'yi, L2 ile spatial field shape'i kontrol ediyorum. Birlikte Bayesian/TPE'nin k arttıkça en stable seçenek olduğunu gösteriyorum.")
    add_text(s, 0.65, 0.86, 5.8, 0.24, "Table 4.4: mean-window MAE", 10, True, PRIMARY)
    add_text(s, 6.75, 0.86, 5.5, 0.24, "Appendix D.13: relative L2 at t=30 s", 10, True, PRIMARY)
    add_image(s, "4.4.png", 0.45, 1.15, 6.05, 4.25)
    add_image(s, "Table D.13.png", 6.65, 1.15, 5.55, 4.25)
    add_card(s, 0.65, 5.65, 3.65, 0.95, "MAE", "Average error across windows. It is the engineering temperature error.", WARM_NOTE)
    add_card(s, 4.55, 5.65, 3.65, 0.95, "L2", "Field error at t=30 s. It checks if the full temperature map is right.")
    add_card(s, 8.45, 5.65, 3.45, 0.95, "Key result", "Bayesian/TPE stays strongest; high k is harder for complex geometry.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Canonical 2D: MAE Variability and L2 Curves", "EN: In this slide, I connect the MAE graph and the L2 graph. On the left, seed variability shows how stable the training is. On the right, L2 curves show how field error changes with time. I found that Bayesian/TPE is more stable, and the early windows are harder because the cooling gradient is steep. So the two graphs tell the same story: accuracy depends on both geometry and window length.\nTR: Bu slaytta MAE grafiği ile L2 grafiğini birbirine bağlar. Solda seed variability training'in ne kadar stabil olduğunu gösteriyorum. Sağda L2 curves field error'ın zamanla nasıl değiştiğini gösteriyorum. Bayesian/TPE'nin daha stabil olduğunu ve early windows'un daha zor olduğunu bulduk çünkü cooling gradient diktir. Yani iki grafik aynı hikayeyi anlatır: accuracy hem geometriye hem window length'e bağlıdır.")
    add_image(s, "seed_results/seed_bar_2d.png", 0.55, 0.95, 5.9, 3.75)
    add_image(s, "seed_results/l2_all_k_2d.png", 6.85, 0.95, 5.9, 3.75)
    add_card(s, 0.75, 5.25, 3.75, 1.0, "MAE graph", "Shows mean error and seed stability.", WARM_NOTE)
    add_card(s, 4.75, 5.25, 3.75, 1.0, "L2 graph", "Shows full-field error over time.")
    add_card(s, 8.75, 5.25, 3.1, 1.0, "Together", "Bayesian/TPE is robust; early windows are hardest.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Canonical 2D Heat Maps: Best Geometry per Architecture", "EN: In this slide, I do not show every 2D heat map. Instead, I choose the strongest example for each architecture. For Bayesian/TPE, I show Rectangle because it has the lowest mean-window MAE in the canonical 2D table. For NSGA-II, I also show Rectangle because it is the most stable low-error 2D case for that architecture. For NSGA-III, I show L-Shape because its compact architecture gives its strongest 2D field result there. My point is to use the heat maps as representative evidence, not as a full appendix repeat.\nTR: Bu slaytta bütün 2D heat map'leri göstermiyorum. Her mimari için en güçlü örnek geometriyi seçiyorum. Bayesian/TPE için Rectangle gösteriyorum çünkü canonical 2D tabloda en düşük mean-window MAE burada. NSGA-II için de Rectangle gösteriyorum çünkü bu mimari için en stabil düşük hata örneği budur. NSGA-III için L-Shape gösteriyorum çünkü compact architecture burada en güçlü 2D field sonucunu veriyor. Benim amacım heat map'leri full appendix tekrarından ziyade representative evidence olarak kullanmak.")
    add_text(s, 0.75, 0.86, 3.6, 0.24, "Bayesian/TPE: Rectangle", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 4.75, 0.86, 3.6, 0.24, "NSGA-II: Rectangle", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 8.75, 0.86, 3.6, 0.24, "NSGA-III: L-Shape", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_image(s, "2d_bestk_bayesian_rectangle.png", 0.55, 1.15, 3.95, 4.75)
    add_image(s, "2d_bestk_nsga2_rectangle.png", 4.55, 1.15, 3.95, 4.75)
    add_image(s, "2d_bestk_nsga3_lshape.png", 8.55, 1.15, 3.95, 4.75)
    add_card(s, 0.8, 6.1, 3.7, 0.65, "Why these?", "One best representative geometry per architecture.", WARM_NOTE)
    add_card(s, 4.75, 6.1, 3.7, 0.65, "How to read", "MAE tells average error; heat map shows where the field error appears.")
    add_card(s, 8.7, 6.1, 3.55, 0.65, "Message", "The best visual cases still need L2/MAE support.", COOL_PANEL)
    i += 1

    # 3D
    # 3D
    s = new_slide(prs, i, "3D Results: Table 4.5 Setup", "EN: In this slide, I show Table 4.5, the setup for the canonical 3D benchmarks. The domains are Rectangular Prism, Cylinder, Stacked Cubes and L-Shape Prism. These domains test volume, curvature, material/interface-like complexity and reentrant edges.\nTR: Bu slaytta canonical 3D benchmarkları için Table 4.5'i veriyorum. Domainler Rectangular Prism, Cylinder, Stacked Cubes ve L-Shape Prism'dir. Bu domainler volume, curvature, interface-like complexity ve reentrant edge etkilerini test eder.")
    add_image(s, "4.5.png", 2.0, 1.0, 9.1, 4.9)
    add_card(s, 1.0, 6.1, 11.35, 0.65, "Key point", "3D and Thermal Fin use the 15 C threshold because the geometry is more complex.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "3D Results: MAE and L2 Summary", "EN: In this slide, I read Table 4.6 together with Appendix Table D.21. The left table gives mean-window MAE for the 3D benchmarks. The right table gives relative L2 error at t=30 s. This matters because a model can have acceptable average MAE but still show spatial field differences. Together, I show that 3D difficulty depends strongly on geometry.\nTR: Bu slaytta Table 4.6 ile Appendix Table D.21'i birlikte okuyorum. Soldaki tablo 3D benchmarklar için mean-window MAE veriyorum. Sağdaki tablo t=30 s'deki relative L2 error veriyorum. Bu önemlidir çünkü model average MAE olarak acceptable olabilir ama spatial field difference gösterebilir. Bu iki tabloyla 3D zorluğun geometriye güçlü şekilde bağlı olduğunu gösteriyorum.")
    add_text(s, 0.65, 0.86, 5.8, 0.24, "Table 4.6: mean-window MAE", 10, True, PRIMARY)
    add_text(s, 6.75, 0.86, 5.5, 0.24, "Appendix D.21: relative L2 at t=30 s", 10, True, PRIMARY)
    add_image(s, "4.6.png", 0.45, 1.15, 6.05, 4.25)
    add_image(s, "Table D.21.png", 6.65, 1.15, 5.55, 4.25)
    add_card(s, 0.65, 5.65, 3.65, 0.95, "MAE", "Average Celsius error over 3D windows. Lower is better.", WARM_NOTE)
    add_card(s, 4.55, 5.65, 3.65, 0.95, "L2", "Field error at t=30 s. It checks the 3D temperature volume.")
    add_card(s, 8.45, 5.65, 3.45, 0.95, "Key result", "Bayesian/TPE is robust, but curved and stacked geometries are harder.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "3D Results: Adaptive-k Table and Schedule", "EN: In this slide, I put the 3D adaptive-k result table and the adaptive-k schedule image side by side. The table gives the numerical result. The schedule image shows how k changes during the run. Read together, they show that the 3D controller is safe but conservative because 3D field errors are close to the promotion threshold.\nTR: Bu slaytta 3D adaptive-k result tablosunu ve adaptive-k schedule görselini yan yana koyar. Tablo numerical result veriyorum. Schedule görseli run sırasında k'nin nasıl değiştiğini gösteriyorum. Birlikte okunduğunda 3D controller'ın güvenli ama conservative olduğunu gösteriyorum, çünkü 3D field errors promotion threshold'a yakındır.")
    add_text(s, 0.75, 0.86, 5.35, 0.24, "Table 4.7: adaptive-k numerical result", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 6.7, 0.86, 5.55, 0.24, "Adaptive-k schedule image", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_image(s, "4.7.png", 0.55, 1.15, 5.65, 3.95)
    add_image(s, "adaptive_k_3d_schedules_table.png", 6.55, 1.15, 5.65, 3.95)
    add_card(s, 0.75, 5.45, 3.7, 0.9, "Table", "Shows MAE and FEM saving for adaptive-k.", WARM_NOTE)
    add_card(s, 4.65, 5.45, 3.7, 0.9, "Schedule", "Shows how the controller changes k over time.")
    add_card(s, 8.55, 5.45, 3.45, 0.9, "Message", "3D adaptive-k is safe, but conservative.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "3D Results: MAE Variability and L2 Curves", "EN: In this slide, I connect the 3D MAE variability and L2 curves. I found that 3D error is not only a dimension problem. Geometry also matters. Cylinder, Stacked Cubes and L-Shape Prism create different gradient patterns. The MAE graph shows stability across seeds, and the L2 graph shows how the full field changes over time.\nTR: Bu slaytta 3D MAE variability ile L2 curves'i birlikte anlatır. 3D error sadece dimension problemi değildir; geometri de önemlidir. Cylinder, Stacked Cubes ve L-Shape Prism farklı gradient patternları oluşturur. MAE grafiği seed stability'yi, L2 grafiği ise full field error'ın zamanla nasıl değiştiğini gösteriyorum.")
    add_image(s, "seed_results/seed_bar_3d.png", 0.55, 0.95, 5.9, 3.75)
    add_image(s, "seed_results/l2_all_k_3d.png", 6.85, 0.95, 5.9, 3.75)
    add_card(s, 0.75, 5.25, 3.75, 1.0, "MAE graph", "Shows seed sensitivity for each 3D geometry.", WARM_NOTE)
    add_card(s, 4.75, 5.25, 3.75, 1.0, "L2 graph", "Shows field error during the 30 s quench.")
    add_card(s, 8.75, 5.25, 3.1, 1.0, "Together", "3D difficulty is geometry-dependent.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Canonical 3D Heat Maps: Best Geometry per Architecture", "EN: In this slide, I use one representative best 3D heat map for each architecture. For Bayesian/TPE, I show Cylinder because it has one of the strongest 3D field-error results. For NSGA-II, I show L-Shape Prism because its L2 result is very low and it is a strong geometry-specific case. For NSGA-III, I also show L-Shape Prism because the compact Pareto architecture performs especially well there. This slide is easier to defend because each heat map has a clear reason.\nTR: Bu slaytta her mimari için bir representative best 3D heat map kullanıyorum. Bayesian/TPE için Cylinder gösteriyorum çünkü 3D field-error sonuçlarında en güçlü örneklerden biridir. NSGA-II için L-Shape Prism gösteriyorum çünkü L2 sonucu çok düşüktür ve güçlü geometry-specific case'tir. NSGA-III için de L-Shape Prism gösteriyorum çünkü compact Pareto architecture burada özellikle iyi performans verir. Bu slaytı savunmak daha kolaydır çünkü her heat map'in net bir seçilme nedeni var.")
    add_text(s, 0.75, 0.86, 3.6, 0.24, "Bayesian/TPE: Cylinder", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 4.75, 0.86, 3.6, 0.24, "NSGA-II: L-Shape Prism", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 8.75, 0.86, 3.6, 0.24, "NSGA-III: L-Shape Prism", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_image(s, "3d_bestk_bayesian_cylinder.png", 0.55, 1.15, 3.95, 4.75)
    add_image(s, "3d_bestk_nsga2_lshape.png", 4.55, 1.15, 3.95, 4.75)
    add_image(s, "3d_bestk_nsga3_lshape.png", 8.55, 1.15, 3.95, 4.75)
    add_card(s, 0.8, 6.1, 3.7, 0.65, "Why these?", "One strong 3D visual case per architecture.", WARM_NOTE)
    add_card(s, 4.75, 6.1, 3.7, 0.65, "How to read", "Volumetric maps show geometry-specific field behaviour.")
    add_card(s, 8.7, 6.1, 3.55, 0.65, "Message", "3D accuracy depends strongly on geometry.", COOL_PANEL)
    i += 1

    # Thermal Fin
    # Thermal Fin
    s = new_slide(prs, i, "Thermal Fin: MAE and L2 Summary", "EN: In this slide, I read Table 4.8 together with the Thermal Fin L2 summary. The left table gives the main MAE and FEM-saving results. The right table gives relative L2 error at t=30 s. MAE is the average window error, but L2 is the final field error. Together, I show that FEM anchoring works, Fourier helps Bayesian/TPE, and PINN-only is not reliable without correction.\nTR: Bu slaytta Table 4.8 ile Thermal Fin L2 summary'yi birlikte okuyorum. Soldaki tablo main MAE ve FEM-saving sonuçlarını veriyorum. Sağdaki tablo t=30 s'deki relative L2 error veriyorum. MAE average window error'dır, L2 ise final field error'dır. Birlikte FEM anchoring'in çalıştığını, Fourier'in Bayesian/TPE'ye yardım ettiğini ve PINN-only'nin correction olmadan reliable olmadığını gösteriyorum.")
    add_text(s, 0.65, 0.86, 5.8, 0.24, "Table 4.8: MAE + FEM saving overview", 10, True, PRIMARY)
    add_text(s, 6.75, 0.86, 5.5, 0.24, "Thermal Fin: relative L2 at t=30 s", 10, True, PRIMARY)
    add_image(s, "4.8.png", 0.45, 1.15, 6.05, 4.15)
    add_table(s, 6.75, 1.25, 5.45, 2.55, ["Architecture", "k=1", "k=2", "k=3", "k=4", "k=5"], [
        ["Bayesian/TPE*", "0.049", "0.053", "0.062", "0.070", "0.080"],
        ["NSGA-II", "0.054", "0.064", "0.070", "0.077", "0.082"],
        ["NSGA-III", "0.057", "0.066", "0.071", "0.077", "0.085"],
    ], size=7)
    add_card(s, 6.8, 4.05, 5.35, 0.82, "Reading note", "MAE = average Celsius error over windows. L2 = relative field error at t=30 s.", WARM_NOTE)
    add_card(s, 0.65, 5.65, 3.65, 0.95, "MAE", "Shows average thermal accuracy and FEM-saving trade-off.", WARM_NOTE)
    add_card(s, 4.55, 5.65, 3.65, 0.95, "L2", "Shows if the final temperature field shape is close to FEM.")
    add_card(s, 8.45, 5.65, 3.45, 0.95, "Key result", "Bayesian/TPE with Fourier has the best field accuracy.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Thermal Fin: Table 4.9 Adaptive-k", "EN: Here, I use Table 4.9 gives the FEM-anchored adaptive-k results for Thermal Fin. All architectures stay within the 15 C threshold, but FEM saving is 45%, which is lower than fixed k=5. I use this to show that the adaptive controller is safe but conservative for this hard geometry.\nTR: Burada Table 4.9 Thermal Fin için FEM-anchored adaptive-k sonuçlarını veriyorum. Tüm mimariler 15 C threshold içinde kalır, fakat FEM saving 45%'tir; bu fixed k=5'ten daha düşüktür. Bu, adaptive controller'ın bu zor geometri için güvenli ama conservative olduğunu gösteriyorum.")
    add_image(s, "4.9.png", 1.0, 1.1, 10.0, 4.2)
    add_card(s, 11.2, 1.2, 2.0, 3.8, "Key", "Safe, but less aggressive than fixed-k.", WARM_NOTE)
    i += 1

    s = new_slide(prs, i, "Thermal Fin: Table 4.10 PINN-only Fixed-k", "EN: Here, I use Table 4.10 shows the PINN-only fixed-k failure. Without FEM correction, mean-window MAE becomes very large. This proves that the method is not a standalone PINN replacement; it is a FEM-anchored accelerator.\nTR: Burada Table 4.10 PINN-only fixed-k failure sonucunu gösteriyorum. FEM correction olmadan mean-window MAE çok büyür. Bu, yöntemin standalone PINN replacement olmadığını; FEM-anchored accelerator olduğunu kanıtlar.")
    add_image(s, "4.10.png", 1.0, 0.95, 10.0, 5.35)
    add_card(s, 11.2, 1.1, 2.0, 4.6, "Key", "No FEM anchor means error accumulation.", WARM_NOTE)
    i += 1

    s = new_slide(prs, i, "Thermal Fin: Table 4.11 PINN-only Adaptive-k", "EN: Here, I use Table 4.11 shows that adaptive-k also fails in PINN-only mode. The controller cannot fix the absence of FEM correction. Even when it stays at k=1, error remains very high.\nTR: Burada Table 4.11 PINN-only mode içinde adaptive-k'nin de başarısız olduğunu gösteriyorum. Controller, FEM correction yokluğunu düzeltemez. k=1'de kalsa bile error çok yüksek kalır.")
    add_image(s, "4.11.png", 1.0, 1.15, 10.0, 4.2)
    add_card(s, 11.2, 1.2, 2.0, 3.8, "Key", "Adaptive control cannot replace FEM correction.", WARM_NOTE)
    i += 1

    s = new_slide(prs, i, "Thermal Fin: MAE Variability and L2 Curves", "EN: In this slide, I connect Thermal Fin MAE stability and L2 field error. I found that Bayesian/TPE with Fourier is stable and has the lowest L2 values. The L2 curve also shows that larger k makes the task harder because the PINN predicts for a longer time without a new FEM anchor. So the message is clear: Fourier helps, but FEM anchoring is still needed.\nTR: Bu slaytta Thermal Fin MAE stability ile L2 field error'ı birlikte anlatır. Bayesian/TPE with Fourier'in stabil olduğunu ve en düşük L2 değerlerini verdiğini bulduk. L2 curve ayrıca büyük k değerlerinde işin zorlaştığını gösteriyorum çünkü PINN yeni FEM anchor olmadan daha uzun süre tahmin yapar. Mesaj nettir: Fourier yardımcı olur, ama FEM anchoring hâlâ gereklidir.")
    add_image(s, "seed_results/seed_bar_thermal_fin.png", 0.55, 0.95, 5.9, 3.75)
    add_image(s, "seed_results/l2_all_k_thermal_fin.png", 6.85, 0.95, 5.9, 3.75)
    add_card(s, 0.75, 5.25, 3.75, 1.0, "MAE graph", "Fourier makes Bayesian/TPE stable.", WARM_NOTE)
    add_card(s, 4.75, 5.25, 3.75, 1.0, "L2 graph", "Higher k increases field error.")
    add_card(s, 8.75, 5.25, 3.1, 1.0, "Together", "Good result, but FEM anchors are still needed.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Thermal Fin: Fixed-k Heat Maps (2D)", "EN: In this slide, I show the report heat maps for the FEM-anchored fixed-k Thermal Fin test. I compare Bayesian/TPE, NSGA-II and NSGA-III. My main point is visual: the FEM reference, PINN prediction and absolute error stay controlled when FEM anchoring is used. This supports the MAE and L2 tables.\nTR: Bu slaytta FEM-anchored fixed-k Thermal Fin testi için report heat map'lerini gösteriyorum. Bayesian/TPE, NSGA-II ve NSGA-III modellerini karşılaştırıyorum. Benim ana mesajım görseldir: FEM reference, PINN prediction ve absolute error FEM anchoring kullanıldığında kontrollü kalır. Bunu MAE ve L2 tablolarını desteklemek için kullanıyorum.")
    add_text(s, 0.95, 0.86, 3.1, 0.24, "Bayesian/TPE", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 4.8, 0.86, 3.1, 0.24, "NSGA-II", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 8.65, 0.86, 3.1, 0.24, "NSGA-III", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_image(s, "thermal_fin_bayesian_fixed_k_clean_heatmaps.png", 0.45, 1.15, 4.0, 4.55)
    add_image(s, "thermal_fin_nsga2_fixed_k_clean_heatmaps.png", 4.55, 1.15, 4.0, 4.55)
    add_image(s, "thermal_fin_nsga3_fixed_k_clean_heatmaps.png", 8.65, 1.15, 4.0, 4.55)
    add_card(s, 1.0, 6.05, 11.35, 0.65, "Interpretation", "These 2D cross-sections show that FEM anchoring keeps the error pattern local and controlled.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Thermal Fin: Adaptive vs PINN-only Heat Maps (2D)", "EN: In this slide, I show the report comparison between FEM-anchored adaptive-k and PINN-only heat maps. The left side shows controlled error with FEM anchoring. The right side shows error drift without FEM correction. I use this to explain why PINN-only tables have very high MAE.\nTR: Bu slaytta FEM-anchored adaptive-k ile PINN-only heat map karşılaştırmasını gösteriyorum. Sol tarafta FEM anchoring ile kontrollü error gösteriyorum. Sağ tarafta FEM correction olmadan error drift gösteriyorum. Bununla PINN-only tablolarında MAE'nin neden çok yüksek olduğunu açıklıyorum.")
    add_text(s, 1.3, 0.86, 4.7, 0.24, "FEM-anchored adaptive-k", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 7.1, 0.86, 4.7, 0.24, "PINN-only", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_image(s, "thermal_fin_fem_adaptive_clean_heatmaps.png", 0.65, 1.15, 5.8, 4.65)
    add_image(s, "thermal_fin_pinn_only_best_clean_heatmaps.png", 6.85, 1.15, 5.8, 4.65)
    add_card(s, 0.95, 6.05, 5.55, 0.65, "With FEM anchors", "The field is corrected at window boundaries.", WARM_NOTE)
    add_card(s, 6.85, 6.05, 5.55, 0.65, "Without FEM anchors", "Errors accumulate and spread across the field.")
    i += 1

    s = new_slide(prs, i, "Thermal Fin: Fixed-k Volumetric Views (3D)", "EN: In this slide, I show the 3D volumetric Thermal Fin views for fixed-k. The 3D views are important because the fin geometry is not only a 2D cross-section problem. They show where the prediction error appears in the full volume.\nTR: Bu slaytta fixed-k için 3D volumetric Thermal Fin görünümlerini gösteriyorum. 3D view önemlidir çünkü fin geometry sadece 2D cross-section problemi değildir. Full volume içinde prediction error'ın nerede oluştuğunu gösteriyorum.")
    add_text(s, 0.95, 0.86, 3.1, 0.24, "Bayesian/TPE", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 4.8, 0.86, 3.1, 0.24, "NSGA-II", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 8.65, 0.86, 3.1, 0.24, "NSGA-III", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_image(s, "thermal_fin_bayesian_fixed_k_clean_3d.png", 0.45, 1.15, 4.0, 4.55)
    add_image(s, "thermal_fin_nsga2_fixed_k_clean_3d.png", 4.55, 1.15, 4.0, 4.55)
    add_image(s, "thermal_fin_nsga3_fixed_k_clean_3d.png", 8.65, 1.15, 4.0, 4.55)
    add_card(s, 1.0, 6.05, 11.35, 0.65, "Interpretation", "The volumetric view confirms the table result: sharp fin regions are the hardest parts of the domain.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Thermal Fin: Adaptive vs PINN-only Volumetric Views (3D)", "EN: In this slide, I show the 3D adaptive-k and PINN-only comparison. In the FEM-anchored case, the field remains physically plausible. In the PINN-only case, drift becomes visible in the volume. This is the visual evidence for why FEM anchoring is essential.\nTR: Bu slaytta 3D adaptive-k ve PINN-only karşılaştırmasını gösteriyorum. FEM-anchored durumda field physically plausible kalır. PINN-only durumda drift volume içinde görünür. Bu, FEM anchoring'in neden essential olduğunu gösteren visual evidence'tır.")
    add_text(s, 1.3, 0.86, 4.7, 0.24, "FEM-anchored adaptive-k", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 7.1, 0.86, 4.7, 0.24, "PINN-only", 10, True, PRIMARY, PP_ALIGN.CENTER)
    add_image(s, "thermal_fin_fem_adaptive_clean_3d.png", 0.65, 1.15, 5.8, 4.65)
    add_image(s, "thermal_fin_pinn_only_best_clean_3d.png", 6.85, 1.15, 5.8, 4.65)
    add_card(s, 0.95, 6.05, 5.55, 0.65, "With FEM anchors", "Periodic correction keeps the 3D field stable.", WARM_NOTE)
    add_card(s, 6.85, 6.05, 5.55, 0.65, "Without FEM anchors", "Prediction drift becomes visible in the volume.")
    i += 1

    s = new_slide(prs, i, "Thermal Fin: Ablation Tables F.1 and F.2", "EN: This slide explains two method choices. Table F.1 shows that Fourier features reduce MAE for Bayesian/TPE on Thermal Fin. Table F.2 shows that endpoint supervision is essential. If we remove endpoint supervision, the error becomes very large. So I found that Fourier helps the network learn sharp gradients, and endpoint supervision prevents window drift.\nTR: Bu slaytta iki method choice'u açıklıyorum. Table F.1, Thermal Fin üzerinde Fourier features'ın Bayesian/TPE için MAE'yi azalttığını gösteriyorum. Table F.2, endpoint supervision'ın essential olduğunu gösteriyorum. Endpoint supervision kaldırılırsa error çok büyür. Yani Fourier sharp gradients'i öğrenmeye yardım ediyor, endpoint supervision ise window drift'i engeller.")
    add_image(s, "Table F.1.png", 0.75, 1.0, 5.75, 2.2)
    add_image(s, "Table F.2.png", 6.85, 1.0, 5.75, 2.2)
    add_card(s, 0.95, 4.05, 5.45, 1.25, "Fourier result", "It lowers MAE and helps with sharp fin gradients.", WARM_NOTE)
    add_card(s, 6.95, 4.05, 5.45, 1.25, "Endpoint result", "It prevents window drift and keeps training stable.")
    add_card(s, 0.95, 5.75, 11.45, 0.7, "Main message", "These are not cosmetic choices; both support stable FEM-anchored prediction on the hardest benchmark.", COOL_PANEL)
    i += 1

    # Discussion
    section_slide(prs, i, "Chapter 5: Discussion", "Interpret the results in the same order as Chapter 4", "EN: The discussion follows the same order as the Results chapter, but the goal is different. In Results, we report what happened. In Discussion, we explain why it happened: why adaptive-k works in 2D, why Bayesian/TPE is robust, why compact architectures degrade on some geometries, why Thermal Fin is difficult, and why FEM anchoring remains essential.\nTR: Discussion bölümü Results chapter ile aynı sırayı takip eder, fakat amacı farklıdır. Results kısmında ne olduğunu raporladık. Discussion kısmında neden olduğunu açıklıyoruz: adaptive-k neden 2D'de iyi çalıştı, Bayesian/TPE neden robust, compact mimariler bazı geometrilerde neden bozuluyor, Thermal Fin neden zor ve FEM anchoring neden hâlâ gerekli.")
    i += 1

    s = new_slide(prs, i, "Discussion: 2D Adaptive-k and Canonical 2D", "EN: In this slide, I cover the first two discussion sections. In the 2D adaptive-k study, the key point is that the controller responds to the physics of quenching. It keeps k small during the early steep-gradient phase and increases k after the temperature field relaxes. In canonical 2D, Bayesian/TPE is most robust because its larger ReLU architecture can handle longer windows and concentrated gradients better than compact tanh networks. L-Shape is important because the reentrant corner creates a localised gradient, so geometry-specific difficulty becomes visible.\nTR: Bu slaytta ilk iki discussion section'ı kapsıyorum. 2D adaptive-k çalışmasında benim ana noktam şu: controller quenching fiziğine tepki veriyor. Erken steep-gradient fazında k küçük kalıyor, field relax olduktan sonra k artıyor. Canonical 2D'de Bayesian/TPE en robust çünkü büyük ReLU mimarisi uzun windowları ve concentrated gradients'i compact tanh networklere göre daha iyi taşıyor. L-Shape önemli çünkü reentrant corner lokal gradient oluşturuyor ve geometriye bağlı zorluk burada görünür hale geliyor.")
    add_card(s, 0.65, 1.0, 3.9, 5.65, "2D adaptive-k", "What I found: all architectures achieved about 83-85% FEM saving with low MAE.\n\nWhy: the controller keeps short windows during the early rapid cooling phase and promotes k only after the field becomes smoother.", WARM_NOTE)
    add_card(s, 4.75, 1.0, 3.9, 5.65, "Canonical 2D", "What I found: Bayesian/TPE stays strongest as k increases. Circle and L-Shape expose geometry sensitivity at high k.\n\nWhy: compact architectures lose accuracy faster when boundaries are curved or gradients are concentrated.")
    add_card(s, 8.85, 1.0, 3.75, 5.65, "Key explanation", "MAE growth is not only a training issue. It is linked to physical gradient structure, boundary shape and the length of unsupported prediction windows.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Discussion: 3D and Thermal Fin", "EN: In this slide, I cover the 3D and Thermal Fin discussion sections. In 3D, geometric complexity changes the behaviour of each architecture. Bayesian/TPE remains robust because it has enough capacity, while compact NSGA models can become unstable on curved surfaces, interfaces, and stacked geometries. Thermal Fin is the strongest test because it has sharp local gradients at the fins. Bayesian/TPE with Fourier features handles this best. PINN-only fails because the model no longer receives FEM correction, so errors accumulate from one window to the next.\nTR: Bu slaytta 3D ve Thermal Fin discussion sectionlarını kapsıyorum. 3D'de geometri karmaşıklığı her mimarinin davranışını değiştiriyor. Bayesian/TPE yeterli kapasiteye sahip olduğu için robust kalıyor; compact NSGA modeller curved surfaces, interfaces ve stacked geometries üzerinde unstable hale gelebiliyor. Thermal Fin en güçlü testtir çünkü fin bölgelerinde sharp local gradients vardır. Bayesian/TPE with Fourier features bunu en iyi taşır. PINN-only başarısız olur çünkü model FEM correction almaz ve hata bir windowdan diğerine birikir.")
    add_card(s, 0.65, 1.0, 3.9, 5.65, "Canonical 3D", "Bayesian/TPE reaches k=5 on all four 3D domains within the 15 C band. NSGA-II/III are more geometry-sensitive at high k.")
    add_card(s, 4.75, 1.0, 3.9, 5.65, "Thermal Fin", "Bayesian/TPE + Fourier remains within the 15 C band at all k values. NSGA-II/III are acceptable through k=4, but k=5 becomes too demanding.", WARM_NOTE)
    add_card(s, 8.85, 1.0, 3.75, 5.65, "PINN-only", "PINN-only errors are much larger because each imperfect prediction becomes the next initial condition. I use this to confirm that periodic FEM correction is necessary.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Discussion: Cross-Domain Guidance", "EN: This slide summarises the final discussion section: Cross-Domain Analysis and MSWP Significance. Across all domains, Bayesian/TPE is the recommended architecture because it gives the strongest accuracy-efficiency frontier. Adaptive-k is recommended for simple and intermediate geometries because it matches the prediction-window length to the transient difficulty. For complex geometries, fixed k=4 or k=5 with Bayesian/TPE can also be viable. The practical conclusion is that MSWP is most useful when FEM calls are repeated many times and become expensive, while a trained PINN can be reused quickly for intermediate predictions.\nTR: Bu slaytta final discussion section'ı özetliyor: Cross-Domain Analysis and MSWP Significance. Tüm domainler boyunca Bayesian/TPE önerilen mimaridir çünkü en güçlü accuracy-efficiency frontier'ı veriyorum. Adaptive-k basit ve intermediate geometriler için önerilir çünkü prediction-window length'i transient zorlukla eşleştirir. Complex geometrilerde Bayesian/TPE ile fixed k=4 veya k=5 de viable olabilir. Pratik sonuç şudur: FEM çağrıları tekrar tekrar çalıştığında maliyetli olur; trained PINN ise intermediate predictions için hızlı tekrar kullanılabilir. MSWP bu durumda en faydalı hale gelir.")
    add_table(s, 0.65, 1.0, 12.0, 5.9, ["Question", "Discussion answer", "Practical guidance"], [
        ["Which architecture transfers best?", "Bayesian/TPE is strongest in 7 of 8 canonical cases and robust on Thermal Fin.", "Use Bayesian/TPE as default deployment architecture."],
        ["When is adaptive-k useful?", "It protects the early steep-gradient phase and increases k as the field relaxes.", "Use adaptive-k for simple/intermediate geometries."],
        ["Can FEM be removed?", "No. PINN-only accumulates error and fails especially on complex geometries.", "Keep periodic FEM anchoring."],
        ["What is MSWP's role?", "It maps geometry and architecture to maximum useful k.", "Use it as a FEM-reduction strategy, not replacement."],
    ], size=9)
    i += 1

    # Conclusion
    section_slide(prs, i, "Chapter 6: Conclusion", "Summary of results, limitations, future work and closing remarks", "EN: The conclusion follows four report sections: Summary of Results, Limitations, Future Work, and Closing Remarks. My main message is positive but careful: FEM calls can be reduced substantially, but FEM anchoring remains necessary. The next step is moving from controlled benchmark geometries toward the real Mortensen industrial mesh and temperature-dependent boiling curve.\nTR: Conclusion dört rapor section'ını takip eder: Summary of Results, Limitations, Future Work ve Closing Remarks. Ana mesaj olumlu ama dikkatli: FEM call sayısı ciddi şekilde azaltılabilir, fakat FEM anchoring hâlâ gereklidir. Sonraki adım controlled benchmark geometrilerden gerçek Mortensen industrial mesh'e ve temperature-dependent boiling curve'e geçmektir.")
    i += 1

    s = new_slide(prs, i, "Conclusion: Summary of Results", "EN: This slide summarises the main findings. The central question was whether FEM solver calls can be replaced by PINN predictions without losing useful accuracy. The answer is yes, with FEM anchoring. Adaptive-k gives 83-85% saving in the 2D transfer study. Bayesian/TPE reaches k=5 on all canonical 3D domains within the 15 C threshold, replacing 16 of 20 FEM calls. Thermal Fin is hardest, but Bayesian/TPE with Fourier remains within the 15 C band at all k values. PINN-only fails, showing that the framework is a FEM accelerator, not a standalone PINN solver.\nTR: Bu slaytta ana bulguları özetler. Temel soru şuydu: FEM solver calls, useful accuracy kaybedilmeden PINN predictions ile değiştirilebilir mi? Cevap evet, ama FEM anchoring ile. Adaptive-k 2D transfer çalışmasında 83-85% saving sağlıyor. Bayesian/TPE canonical 3D domainlerin tamamında k=5'e kadar 15 C threshold içinde kalıyor ve 20 FEM çağrısından 16'sını değiştiriyor. Thermal Fin en zor case ama Bayesian/TPE with Fourier tüm k değerlerinde 15 C bandında kalıyor. PINN-only başarısız oluyor; bu da frameworkün standalone PINN solver değil, FEM accelerator olduğunu gösteriyor.")
    add_table(s, 0.65, 1.0, 12.0, 5.9, ["Finding", "Evidence", "Meaning"], [
        ["FEM calls can be reduced", "65-85% saving in many anchored cases", "MSWP is effective as an accelerator"],
        ["Bayesian/TPE is most robust", "Strong 2D, 3D and Thermal Fin performance", "Recommended default architecture"],
        ["Adaptive-k helps", "83-85% saving with low MAE in 2D transfer study", "Useful for early steep gradients"],
        ["PINN-only fails", "Thermal Fin best PINN-only is 59.20 C", "FEM correction is necessary"],
    ], size=10)
    i += 1

    s = new_slide(prs, i, "Conclusion: Limitations", "EN: This slide states the limitations clearly. The thesis uses constant material properties and a constant convection coefficient, while real industrial A356 quenching has temperature-dependent properties and a five-regime boiling curve. The reference solutions come from our own Crank-Nicolson solver rather than StaMiSim or Abaqus. The grids are structured and mesh independence is not fully verified. These limitations do not invalidate the benchmark study, but they define the boundary between the current validated framework and full industrial deployment.\nTR: Bu slaytta limitations kısmını net açıklıyorum. Tezde constant material properties ve constant convection coefficient kullanıldı; gerçek endüstriyel A356 quenching ise temperature-dependent properties ve five-regime boiling curve içerir. Reference solutionlar StaMiSim veya Abaqus yerine kendi Crank-Nicolson solverımızdan geldi. Gridler structured ve mesh independence tam olarak doğrulanmadı. Bu sınırlılıklar benchmark çalışmasını geçersiz kılmaz; fakat mevcut validated framework ile full industrial deployment arasındaki sınırı belirler.")
    add_card(s, 0.65, 1.0, 3.9, 5.65, "Physics simplification", "Constant kT, rho cp and h. Real quenching uses temperature-dependent material properties and five boiling regimes.")
    add_card(s, 4.75, 1.0, 3.9, 5.65, "Solver validation", "Reference fields come from the in-house Crank-Nicolson solver, not from StaMiSim or Abaqus. Mesh independence remains future work.", WARM_NOTE)
    add_card(s, 8.85, 1.0, 3.75, 5.65, "NAS and data scope", "NAS budget and benchmark geometries are finite. The real Mortensen mesh is the next deployment target, not yet the main experiment.")
    i += 1

    s = new_slide(prs, i, "Conclusion: Future Work and Closing Remarks", "EN: This final conclusion slide gives the roadmap. The immediate next step is applying Bayesian/TPE to the real Mortensen industrial mesh and FEM snapshots. Then the model should include the full temperature-dependent boiling curve h(T), architecture-specific adaptive-k thresholds, residual-based collocation sampling, and validation against StaMiSim or Abaqus. The closing message is that the results are not only a standalone benchmark; they are the validated foundation for a PINN-based computational accelerator for industrial aluminium quenching.\nTR: Bu final conclusion slaytı roadmap'i veriyorum. En yakın sonraki adım Bayesian/TPE'yi gerçek Mortensen industrial mesh ve FEM snapshotlarına uygulamaktır. Sonra full temperature-dependent boiling curve h(T), architecture-specific adaptive-k thresholds, residual-based collocation sampling ve StaMiSim/Abaqus validation eklenmelidir. Kapanış mesajı şu: bu sonuçlar sadece standalone benchmark değildir; endüstriyel aluminium quenching için PINN-based computational accelerator'ın validated foundation kısmıdır.")
    add_table(s, 0.65, 1.0, 12.0, 5.9, ["Future step", "Why it matters"], [
        ["Apply to Mortensen industrial mesh", "Direct test on the target deployment geometry"],
        ["Use full h(T) boiling curve", "Closer physics to industrial water quenching"],
        ["Architecture-specific adaptive-k", "Better thresholds for each architecture's baseline error"],
        ["Residual-based sampling", "More training effort where physics residual is hardest"],
        ["Validate against StaMiSim/Abaqus", "Confirm that benchmark errors reflect true physics accuracy"],
    ], size=10)
    i += 1

    # References
    s = new_slide(prs, i, "References", "Referans slaytı. Burada özellikle Mortensen makalesini ve NAS-PINN makalesini vurgula. Sorulursa tekrar söyle: problem kaynağı Mortensen, yöntem temeli Wang and Zhong NAS-PINN'dir.")
    add_bullets(s, 0.75, 1.0, 12.0, 5.95, [
        "Mortensen, D., Noorsumar, G., Fjaer, H. G., Babaei, R., & Dronen, P. E. (2026). Mitigating distortions in cast automotive subframes: A finite element simulation approach. The International Journal of Advanced Manufacturing Technology, 142, 5879-5898.",
        "Wang, Y., & Zhong, L. (2024). NAS-PINN: Neural architecture search-guided physics-informed neural network for solving PDEs. Journal of Computational Physics, 496, 112603.",
        "Raissi, M., Perdikaris, P., & Karniadakis, G. E. (2019). Physics-informed neural networks. Journal of Computational Physics, 378, 686-707.",
        "Deb, K. et al. (2002). A fast and elitist multiobjective genetic algorithm: NSGA-II. IEEE Transactions on Evolutionary Computation.",
        "Deb, K., & Jain, H. (2014). An evolutionary many-objective optimization algorithm using reference-point-based nondominated sorting: NSGA-III. IEEE Transactions on Evolutionary Computation.",
        "Akiba, T. et al. (2019). Optuna: A next-generation hyperparameter optimization framework. KDD.",
    ], size=11)
    i += 1

    s = new_slide(prs, i, "Q&A Preparation", "EN: In this slide, I prepare for likely questions. I do not need to read every row during the defense. I can use it as a backup slide if the committee asks about FEM, PINN, NAS-PINN, MAE, L2, adaptive-k, or why PINN-only fails. My goal is to answer simply and connect each answer back to the results.\nTR: Bu slaytta olası sorulara hazırlanıyorum. Savunmada her satırı okumam gerekmez. Committee FEM, PINN, NAS-PINN, MAE, L2, adaptive-k veya PINN-only failure sorarsa backup slide olarak kullanabilirim. Amacım basit cevap vermek ve her cevabı result kısmına bağlamaktır.")
    add_table(s, 0.55, 0.95, 12.3, 5.95, ["Question", "Short answer", "Result link"], [
        ["What is FEM?", "A trusted numerical solver used as reference and anchor.", "It gives the correction snapshots."],
        ["What is PINN?", "A neural network trained with data error and heat-equation physics.", "It predicts between FEM anchors."],
        ["What is NAS-PINN?", "A PINN whose architecture is selected by search, not by hand.", "Bayesian/TPE is the most robust choice."],
        ["MAE vs L2?", "MAE is average Celsius error; L2 is field error at t=30 s.", "Both are needed to judge accuracy."],
        ["Why adaptive-k?", "It changes window length based on error.", "It saves FEM calls safely in simpler cases."],
        ["Why not PINN-only?", "Without FEM anchors, errors accumulate.", "Thermal Fin PINN-only fails clearly."],
        ["Why is PINN useful if training is slow?", "Training is paid once; inference is fast for repeated use.", "Useful when many FEM solves are needed."],
    ], size=7)
    add_card(s, 0.85, 6.35, 11.7, 0.45, "How I use this slide", "I answer in one or two sentences, then point back to the MAE/L2 or heat-map evidence.", COOL_PANEL)
    i += 1

    s = new_slide(prs, i, "Thank You", "EN: This is my final slide. I thank the committee and invite questions. If there is time, I can return to the Q&A backup table or to the heat maps and MAE/L2 result slides.\nTR: Bu benim final slaytım. Committee'ye teşekkür ediyorum ve soruları davet ediyorum. Zaman olursa Q&A backup tablosuna veya heat map ve MAE/L2 result slaytlarına geri dönebilirim.")
    add_text(s, 0.9, 2.15, 11.5, 0.75, "Thank you", 36, True, PRIMARY, PP_ALIGN.CENTER)
    add_text(s, 0.9, 3.05, 11.5, 0.45, "Questions and discussion", 20, False, DARK, PP_ALIGN.CENTER)
    add_card(s, 2.0, 4.15, 9.3, 1.35, "Main takeaway", "FEM-anchored NAS-PINN can reduce FEM calls while keeping useful thermal-field accuracy. FEM is accelerated, not removed.", WARM_NOTE)
    add_text(s, 0.9, 6.35, 11.5, 0.28, "Teşekkürler", 16, True, MUTED, PP_ALIGN.CENTER)

    prs.save(OUT)
    return OUT

if __name__ == "__main__":
    print(make_deck())
