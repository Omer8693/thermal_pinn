"""
rebuild_pptx_v3.py
Rebuild NAS_PINN_Defense_Final_v2.pptx → NAS_PINN_Defense_Final_v3.pptx
applying all 22 requested changes.
"""

import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.opc.packuri import PackURI
from lxml import etree

SRC = "thermal_pinn/reports/thesis_uia/NAS_PINN_Defense_Final_v2.pptx"
DST = "thermal_pinn/reports/thesis_uia/NAS_PINN_Defense_Final_v3.pptx"

# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────

def delete_slide(prs, idx):
    """Delete slide at 0-based index idx."""
    xml_slides = prs.slides._sldIdLst
    prs.part.drop_rel(xml_slides[idx].rId)
    del xml_slides[idx]


def move_slide(prs, old_idx, new_idx):
    """Move slide from old_idx to new_idx (0-based)."""
    xml_slides = prs.slides._sldIdLst
    elem = xml_slides[old_idx]
    xml_slides.remove(elem)
    xml_slides.insert(new_idx, elem)


def set_tf_text(tf, text):
    """Set full text frame to a single string, clearing existing content."""
    paras = tf.paragraphs
    for para in paras[1:]:
        p = para._p
        p.getparent().remove(p)
    p0 = paras[0]
    for run in p0.runs[1:]:
        run._r.getparent().remove(run._r)
    if p0.runs:
        p0.runs[0].text = text
    else:
        r = etree.SubElement(p0._p, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'), attrib={'lang': 'en-US', 'dirty': '0'})
        t = etree.SubElement(r, qn('a:t'))
        t.text = text


def set_notes(slide, text):
    """Set speaker notes on a slide, creating notes slide if needed."""
    try:
        ns = slide.notes_slide
        ntf = ns.notes_text_frame
        set_tf_text(ntf, text)
    except Exception as e:
        pass  # Notes creation can fail on freshly added blank slides; skip silently


def find_text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame]


def remove_picture_shapes(slide):
    """Remove all picture shapes from a slide."""
    pics_to_remove = [s for s in slide.shapes if s.shape_type == 13]
    for pic in pics_to_remove:
        try:
            blip = pic._element.find('.//' + qn('a:blip'))
            if blip is not None:
                rId = blip.get(qn('r:embed'))
                if rId:
                    try:
                        slide.part.drop_rel(rId)
                    except Exception:
                        pass
            pic._element.getparent().remove(pic._element)
        except Exception as e:
            print(f"  WARNING removing picture: {e}")


def slide_title_text(slide):
    """Return first non-empty text shape text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                return t
    return ""


def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=None):
    """Add a new text box to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def rename_slide_parts(prs):
    """
    Rename all slide parts to use sequential numbering: slide1.xml, slide2.xml, …
    This eliminates duplicate-name warnings when saving.
    Uses a two-pass approach to avoid collisions.
    """
    n = len(prs.slides)
    # Pass 1: rename to temp names to avoid any collisions
    for i in range(n):
        prs.slides[i].part.partname = PackURI(f'/ppt/slides/slideTEMP{i+1}.xml')
    # Pass 2: rename to final names
    for i in range(n):
        prs.slides[i].part.partname = PackURI(f'/ppt/slides/slide{i+1}.xml')
    print("  Renamed slide parts to sequential numbering.")


# ──────────────────────────────────────────────
# LOAD
# ──────────────────────────────────────────────
prs = Presentation(SRC)
print(f"Loaded: {len(prs.slides)} slides")

# ──────────────────────────────────────────────
# PHASE 1: DELETIONS (back-to-front to keep indices stable)
# ──────────────────────────────────────────────

slides_to_delete_1based = [48, 47, 46, 45, 44, 43, 42, 41, 24, 20, 19, 18, 15, 13, 12, 11]
slides_to_delete_0based = sorted([x - 1 for x in slides_to_delete_1based], reverse=True)

print("Deleting slides (0-based):", slides_to_delete_0based)
for idx in slides_to_delete_0based:
    title = slide_title_text(prs.slides[idx])
    print(f"  Deleting slide {idx+1}: {title[:60]}")
    delete_slide(prs, idx)

print(f"After deletions: {len(prs.slides)} slides")

# Rename slide parts to eliminate gaps and avoid future collisions
rename_slide_parts(prs)

# Print current slide order
print("Current slide order after deletions:")
for i, slide in enumerate(prs.slides):
    print(f"  {i+1:2d}: {slide_title_text(slide)[:70]}")

# ──────────────────────────────────────────────
# PHASE 2: MODIFY REMAINING SLIDES
# After deletions and rename, order is:
#  1  Title slide
#  2  Presentation Structure
#  3  Intro: What We Do
#  4  A356 aluminium
#  5  From Reference Paper
#  6  Reference Paper Inputs
#  7  Research Question
#  8  Heat Transfer Model
#  9  FEM/PINN/NAS defs
# 10  FEM Process: How It Works  (merge with deleted 11)
# 11  Method Source: NAS-PINN
# 12  NAS-PINN Process
# 13  Architecture Search
# 14  MSWP Framework
# 15  Skip Factor
# 16  Benchmark Domains 2D
# 17  Evaluation Metrics
# 18  2D Adap-k MAE/L2
# 19  2D Adap-k Curves
# 20  Canon2D Setup
# 21  Canon2D MAE
# 22  Canon2D Variability
# 23  Canon2D Heat Maps
# 24  3D MAE/L2
# 25  3D Variability
# 26  3D Heat Maps
# 27  Thermal Fin MAE
# 28  Thermal Fin Variability
# 29  Thermal Fin Fixed-k 2D
# 30  Thermal Fin Adaptive 2D
# 31  Thermal Fin Fixed-k 3D
# 32  Thermal Fin Adaptive 3D
# 33  Q&A
# 34  Thank You
# ──────────────────────────────────────────────

def S(n):
    return prs.slides[n - 1]

# ── CHANGE 1: Slide 1 — Update title ──
print("\n[Change 1] Slide 1: Update title")
slide1 = S(1)
for shape in slide1.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text.strip()
    if 'NAS-Guided FEM-Anchored PINNs' in t and 'Multi-Step' not in t:
        set_tf_text(shape.text_frame,
            "NAS-Guided FEM-Anchored PINNs with Multi-Step Window Prediction for Transient Heat Transfer Simulation")
        print("  Updated main title")
    elif 'Multi-Step Window Prediction for Transient A356' in t:
        set_tf_text(shape.text_frame,
            "Step Optimization in Thermal Simulation Using Selective Finite Element Method Step Skipping")
        print("  Updated subtitle")
    elif 'FEM accelerator inspired' in t:
        set_tf_text(shape.text_frame, "")
        print("  Cleared third subtitle line")

# ── CHANGE 3: Slide 4 — Remove image ──
print("\n[Change 3] Slide 4: Remove image(s)")
remove_picture_shapes(S(4))
print("  Removed picture shapes from slide 4")

# ── CHANGE 4: Slide 10 — Add mesh config summary (slide 11 already deleted) ──
print("\n[Change 4] Slide 10: Add mesh config + update title")
slide10 = S(10)
for shape in slide10.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if 'FEM Process: How It Works' in t:
            set_tf_text(shape.text_frame, "FEM Process and Mesh Configuration")
            print("  Updated slide 10 title")
            break

mesh_text = (
    "Mesh config: Grid: 131×61 (Rectangle, 7991 nodes), 121×121 (Circle, 11285 nodes), 81×81 (L-Shape, 4961 nodes)"
)
tb = slide10.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.0), Inches(0.8))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
run = p.add_run()
run.text = mesh_text
run.font.size = Pt(11)
print("  Added mesh config summary box to slide 10")

# ── CHANGE 7: Slide 11 (was 14: Method Source NAS-PINN) — Rewrite as 4-step ──
print("\n[Change 7] Slide 11: Rewrite as NAS-PINN How It Works")
slide11 = S(11)
for shape in slide11.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text.strip()
    if 'Method Source' in t or ('NAS-PINN' in t and 'Process' not in t and 'How' not in t):
        set_tf_text(shape.text_frame, "NAS-PINN: How It Works — Four Simple Steps")
        print("  Updated title")
        break

step11_updates = {
    'NAS-PINN source': 'Step 1: Choose a search space',
    'Wang and Zhong (2024) provide': 'Layers, neurons, activation functions, Fourier embedding options form the search space.',
    'How this thesis adapts it': 'Step 2: Run NAS optimizer',
    'The thesis keeps the NAS-PINN': 'Bayesian/TPE, NSGA-II, or NSGA-III tests candidate architectures on the heat equation loss.',
}
for shape in slide11.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text.strip()
    for old_key, new_text in step11_updates.items():
        if old_key in t:
            set_tf_text(shape.text_frame, new_text)
            print(f"  Updated: {old_key[:40]}")
            break

add_textbox(slide11, Inches(0.5), Inches(4.8), Inches(6.0), Inches(0.55),
            "Step 3: Train best PINN", font_size=13, bold=True)
add_textbox(slide11, Inches(0.5), Inches(5.4), Inches(11.5), Inches(0.55),
            "Each candidate network is trained with Adam + L-BFGS to minimise the physics residual and anchor loss.",
            font_size=11)
add_textbox(slide11, Inches(0.5), Inches(6.0), Inches(6.0), Inches(0.55),
            "Step 4: Transfer architecture", font_size=13, bold=True)
add_textbox(slide11, Inches(0.5), Inches(6.5), Inches(11.5), Inches(0.55),
            "The winning architecture is reused across all benchmark geometries.  Citation: Wang & Zhong (2024) — NAS-PINN",
            font_size=11)
print("  Added Steps 3 and 4")

# ── CHANGE 8: Slide 12 (NAS-PINN Process) — Update step descriptions ──
print("\n[Change 8] Slide 12: NAS-PINN Process — update step descriptions")
slide12 = S(12)
step12_updates = {
    'Step 1: anchor': 'Step 1: Anchor',
    'Use the FEM field at the start': 'FEM (Finite Element Method) provides the trusted starting state at the beginning of each prediction window.',
    'Step 2: prediction': 'Step 2: Prediction',
    'PINN predicts T(x,t)': 'NAS-PINN (Neural Architecture Search PINN) predicts temperature field T(x,t) inside the window.',
    'Step 3: loss': 'Step 3: Loss',
    'Physics and endpoint terms': 'Physics residual loss + endpoint supervision keep each window prediction anchored to FEM.',
    'Step 4: NAS': 'Step 4: NAS',
    'Architecture search chooses': 'Neural Architecture Search (NAS) selects the best network design for the given geometry.',
}
for shape in slide12.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text.strip()
    for old_key, new_text in step12_updates.items():
        if old_key in t:
            set_tf_text(shape.text_frame, new_text)
            break

# ── CHANGE 12: Slide 14 (MSWP) — Spell out abbreviations ──
print("\n[Change 12] Slide 14: MSWP — spell out abbreviations")
slide14 = S(14)
mswp_updates = {
    'MSWP Framework: Working Logic': 'Multi-Step Window Prediction (MSWP) Framework — Working Logic',
    'Step 1: FEM anchors': 'Step 1: FEM (Finite Element Method) anchors',
    'FEM gives trusted reference states': 'FEM provides trusted reference temperature states at selected time steps.',
    'Step 2: choose k': 'Step 2: Choose skip factor k',
    'The skip factor decides': 'The skip factor k decides how many FEM steps are replaced by PINN (Physics-Informed Neural Network) predictions.',
    'Step 3: PINN predicts': 'Step 3: NAS-PINN (Neural Architecture Search PINN) predicts',
    'NAS-PINN fills the skipped': 'NAS-PINN fills the k skipped thermal windows between FEM anchors.',
    'Step 4: correct': 'Step 4: FEM corrects',
    'The next FEM anchor limits drift': 'The next FEM anchor resets the field, preventing error drift from accumulating across windows.',
}
for shape in slide14.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text.strip()
    for old_key, new_text in mswp_updates.items():
        if old_key in t:
            set_tf_text(shape.text_frame, new_text)
            print(f"  Updated: {old_key[:45]}")
            break

# ── CHANGE 13: Move Slide 15 (Skip Factor) to right after Slide 7 (Research Question) ──
print("\n[Change 13] Move Skip Factor slide to after Research Question")
# Skip Factor is currently at index 14 (slide 15), Research Question at index 6 (slide 7)
move_slide(prs, 14, 7)
print("  Moved Skip Factor from position 15 to position 8")

# Renumber after move:
#  7=Research Question, 8=Skip Factor, 9=Heat Transfer, 10=FEM/PINN/NAS,
#  11=FEM Process+Mesh, 12=NAS-PINN Steps, 13=NAS-PINN Process, 14=Architecture Search,
#  15=MSWP, 16=Benchmark Domains 2D, ...

print("Slide order after move:")
for i, slide in enumerate(prs.slides):
    print(f"  {i+1:2d}: {slide_title_text(slide)[:65]}")

# ── CHANGE 14: Slide 16 (Benchmark Domains 2D) — Update to include 3D/Fin ──
print("\n[Change 14] Slide 16: Benchmark Domains — update")
slide16 = S(16)
for shape in slide16.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text.strip()
    if 'Benchmark Domains: 2D Domain Groups' in t:
        set_tf_text(shape.text_frame, "Benchmark Domains: 2D, 3D, and Thermal Fin")
        print("  Updated title")
    elif 'Level 1: NAS transfer' in t or t == 'Level 1: NAS transfer + adaptive-k':
        set_tf_text(shape.text_frame, "Level 1 — NAS Transfer + Adaptive-k")
    elif 'Domains: Square' in t and 'Purpose' in t:
        set_tf_text(shape.text_frame,
            "Square, Circle, L-Shape, Flower — k_init=1, adaptive. Tests NAS-PINN architecture transfer across diverse 2D shapes.")
    elif 'Level 2' in t:
        set_tf_text(shape.text_frame, "Level 2 — Canonical 2D Fixed-k")
    elif 'Rectangle, Circle, L-Shape' in t and 'Purpose' in t:
        set_tf_text(shape.text_frame,
            "Rectangle, Circle, L-Shape — k=1…5. Compares Bayesian/TPE, NSGA-II, NSGA-III at fixed skip factors.")

add_textbox(slide16, Inches(6.6), Inches(1.4), Inches(6.3), Inches(0.5),
            "Level 3 — Canonical 3D", font_size=12, bold=True)
add_textbox(slide16, Inches(6.6), Inches(1.9), Inches(6.3), Inches(0.9),
            "Rectangular, Cylinder, Stacked Cubes, L-Shape Prism — tests method on 3D volumetric thermal fields.",
            font_size=11)
add_textbox(slide16, Inches(6.6), Inches(3.0), Inches(6.3), Inches(0.5),
            "Level 4 — Thermal Fin", font_size=12, bold=True)
add_textbox(slide16, Inches(6.6), Inches(3.5), Inches(6.3), Inches(0.9),
            "Most demanding — thin-fin geometry creates local gradients; tests FEM anchoring and Fourier features.",
            font_size=11)
print("  Added 3D/Fin descriptions")

# ── CHANGE 16: Slide 19 (FEM-PINN Curves) — Shrink text boxes ──
print("\n[Change 16] Slide 19: FEM-PINN Curves — reduce font size")
slide19 = S(19)
for shape in slide19.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text.strip()
    if t and 'Master Thesis' not in t and 'FEM-PINN' not in t and 'Adaptive-k' not in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(10)

# ──────────────────────────────────────────────
# PHASE 3: ADD NEW SLIDES
# We add at end, then move into position.
# ──────────────────────────────────────────────

print("\n[Phase 3] Adding new slides...")

def add_blank_slide(prs):
    """Add a blank slide at the end using layout 6 (blank)."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

# ── CHANGE 17: New 3D Parameters slide ──
# Will insert before slide 24 (currently 3D MAE/L2).
# After our moves, 3D MAE/L2 is at index 23 (slide 24).
print("[Change 17] Adding 3D Parameters slide")
slide_3d = add_blank_slide(prs)
new_slide_idx = len(prs.slides) - 1

# Populate 3D params slide
add_textbox(slide_3d, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
            "3D Benchmark: Physical and Numerical Parameters",
            font_size=22, bold=True, color=RGBColor(0x1F, 0x38, 0x64))
add_textbox(slide_3d, Inches(0.5), Inches(0.9), Inches(4), Inches(0.3),
            "Master Thesis | 3D Parameters", font_size=9)

params = [
    ("Parameter", "Symbol", "Value"),
    ("Thermal conductivity", "kT", "160 W/(m·K)"),
    ("Volumetric heat capacity", "ρcp", "2.4 × 10⁶ J/(m³·K)"),
    ("Convection coefficient", "h", "4000 W/(m²·K)"),
    ("Initial temperature", "T₀", "540°C"),
    ("Quench bath temperature", "Tw", "20°C"),
    ("PINN anchor spacing", "Δtₐ", "1.5 s"),
    ("FEM time step", "Δtᵢₙₜ", "0.5 s"),
    ("Adam epochs", "—", "1000"),
    ("L-BFGS iterations/window", "—", "300"),
]

rows = len(params)
cols = 3
table_shape = slide_3d.shapes.add_table(rows, cols, Inches(1.0), Inches(1.4), Inches(11.0), Inches(4.6))
table = table_shape.table
table.columns[0].width = Inches(5.5)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(3.5)

for r, row_data in enumerate(params):
    for c, cell_text in enumerate(row_data):
        cell = table.cell(r, c)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        if r == 0:
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', '1F3864')
        elif r % 2 == 0:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', 'E8EEF7')

# Move 3D params slide to before slide 24 (index 23 in current 0-based)
move_slide(prs, len(prs.slides) - 1, 23)
print(f"  3D Parameters slide added and moved to position 24")
print(f"  Slide count: {len(prs.slides)}")

# ── CHANGE 20: New Discussion slide ──
print("\n[Change 20] Adding Discussion slide")
slide_disc = add_blank_slide(prs)

add_textbox(slide_disc, Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.65),
            "Discussion", font_size=26, bold=True, color=RGBColor(0x1F, 0x38, 0x64))
add_textbox(slide_disc, Inches(0.4), Inches(0.72), Inches(4), Inches(0.3),
            "Master Thesis | Discussion", font_size=9)

NAVY = RGBColor(0x1F, 0x38, 0x64)
disc_sections = [
    ("2D Results",
     "83-85% FEM saving across all architectures. Adaptive-k controller reduces errors during steep early-transient gradients. "
     "Bayesian/TPE stays most robust as k increases. Complex boundaries (L-Shape, Circle) raise local errors at large k."),
    ("3D Results",
     "Bayesian/TPE reaches k=5 within the 15°C band on all four 3D geometries. "
     "Geometry sensitivity grows with k — curved surfaces and prism corners increase prediction difficulty."),
    ("Thermal Fin",
     "Fourier features stabilise Bayesian/TPE across all k values. "
     "FEM anchors prevent drift accumulation in the sharp fin regions. "
     "PINN-only errors are 4-18× larger — periodic correction is essential."),
    ("Cross-Domain Guidance + Ablation",
     "FEM anchoring is essential: PINN-only errors are 4-18× larger. "
     "Fourier features lower MAE; endpoint supervision prevents window drift. "
     "Use smaller k for accuracy; larger k only if MAE/L2 budget allows."),
]

positions = [
    (Inches(0.3), Inches(1.1)),
    (Inches(6.6), Inches(1.1)),
    (Inches(0.3), Inches(3.8)),
    (Inches(6.6), Inches(3.8)),
]
box_w = Inches(6.0)

for (left, top), (header, body) in zip(positions, disc_sections):
    add_textbox(slide_disc, left, top, box_w, Inches(0.45),
                header, font_size=13, bold=True, color=NAVY)
    add_textbox(slide_disc, left, top + Inches(0.45), box_w, Inches(2.3),
                body, font_size=11)

# Insert Discussion before Q&A — Q&A is currently at index 32 (slide 33), then Thank You
# After adding 3D params, count is 35. Q&A is slide 34 (index 33), Thank You slide 35 (index 34)
move_slide(prs, len(prs.slides) - 1, 33)
print(f"  Discussion slide added at position 34, slide count={len(prs.slides)}")

# ── CHANGE 21: New Conclusion slide ──
print("\n[Change 21] Adding Conclusion slide")
slide_conc = add_blank_slide(prs)

add_textbox(slide_conc, Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.65),
            "Conclusion", font_size=26, bold=True, color=NAVY)
add_textbox(slide_conc, Inches(0.4), Inches(0.72), Inches(4), Inches(0.3),
            "Master Thesis | Conclusion", font_size=9)

conc_sections = [
    ("What We Did",
     "Built a FEM-anchored NAS-PINN framework with Multi-Step Window Prediction (MSWP). "
     "Validated across 2D (Square, Circle, L-Shape, Flower, Rectangle), "
     "3D (Rectangular, Cylinder, Stacked Cubes, L-Shape Prism), and Thermal Fin domains."),
    ("What We Found",
     "65-85% FEM saving achieved with MAE within engineering limits: ~10°C for 2D, ~15°C for 3D/Fin. "
     "Bayesian/TPE + Fourier features is the most robust combination. "
     "FEM anchoring prevents error accumulation — PINN-only errors are 4-18× larger."),
    ("Limitations and Next Steps",
     "Simplified material model: constant kT, ρcp, h (real quenching uses temperature-dependent properties). "
     "Future: temperature-dependent properties, industrial Mortensen automotive subframe mesh, "
     "StaMiSim solver validation, and real manufacturing boundary conditions."),
]

tops = [Inches(1.1), Inches(3.1), Inches(5.1)]
for top_pos, (header, body) in zip(tops, conc_sections):
    add_textbox(slide_conc, Inches(0.4), top_pos, Inches(12.5), Inches(0.48),
                header, font_size=14, bold=True, color=NAVY)
    add_textbox(slide_conc, Inches(0.4), top_pos + Inches(0.48), Inches(12.5), Inches(1.52),
                body, font_size=12)

# Insert Conclusion after Discussion (index 34), before Q&A (index 35)
move_slide(prs, len(prs.slides) - 1, 34)
print(f"  Conclusion slide added at position 35, slide count={len(prs.slides)}")

print("\nCurrent slide order (after all adds and moves):")
for i, slide in enumerate(prs.slides):
    print(f"  {i+1:2d}: {slide_title_text(slide)[:65]}")

# ──────────────────────────────────────────────
# PHASE 4: FURTHER SLIDE CONTENT CHANGES
# ──────────────────────────────────────────────

# ── CHANGE 18: 3D Heat Maps — Cylinder only ──
print("\n[Change 18] Update 3D Heat Maps slide")
for i, slide in enumerate(prs.slides):
    t = slide_title_text(slide)
    if 'Canonical 3D Heat Maps' in t or ('3D Heat Map' in t and 'Cylinder' not in t):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            st = shape.text_frame.text.strip()
            if 'Canonical 3D Heat Maps' in st:
                set_tf_text(shape.text_frame, "3D Heat Map: Cylinder Domain (Bayesian/TPE, Best Result)")
            elif 'NSGA-II: L-Shape Prism' in st:
                set_tf_text(shape.text_frame, "")
            elif 'NSGA-III: L-Shape Prism' in st:
                set_tf_text(shape.text_frame, "")
            elif 'Bayesian/TPE: Cylinder' in st:
                set_tf_text(shape.text_frame, "Bayesian/TPE: Cylinder (Best 3D Result)")
        print(f"  Updated 3D Heat Maps at slide {i+1}")
        break

# ── CHANGE 19: Thermal Fin heat map slides — reduce text ──
print("\n[Change 19] Reduce text on Fixed-k heat map slides")
for i, slide in enumerate(prs.slides):
    t = slide_title_text(slide)
    if 'Fixed-k Heat Maps (2D)' in t or 'Fixed-k Volumetric' in t:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            st = shape.text_frame.text.strip()
            if (st and 'Master Thesis' not in st and 'Fixed-k' not in st
                    and 'Thermal Fin' not in st and 'Bayesian' not in st
                    and 'NSGA' not in st and 'Interpretation' not in st):
                if len(st) > 80:
                    set_tf_text(shape.text_frame, st[:120])
        print(f"  Reduced text on slide {i+1}: {t[:50]}")

# Merge Thermal Fin Adaptive 2D → combined title (slide 31 after our additions)
print("\n[Change 19b] Update Thermal Fin Adaptive slide to combined title")
for i, slide in enumerate(prs.slides):
    t = slide_title_text(slide)
    if 'Adaptive vs PINN-only Heat Maps (2D)' in t:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            st = shape.text_frame.text.strip()
            if 'Adaptive vs PINN-only Heat Maps (2D)' in st:
                set_tf_text(shape.text_frame,
                    "Thermal Fin: Anchored vs PINN-only (2D and 3D comparison)")
                print(f"  Updated adaptive slide title at {i+1}")
        add_textbox(slide, Inches(0.3), Inches(6.45), Inches(12.3), Inches(0.35),
                    "Note: 3D volumetric comparison follows — FEM-anchored keeps field stable; PINN-only shows drift accumulation.",
                    font_size=10)
        break

# ──────────────────────────────────────────────
# PHASE 5: UPDATE ALL SPEAKER NOTES (EN+TR format)
# ──────────────────────────────────────────────

print("\n[Phase 5] Updating all speaker notes with EN+TR format...")

def format_note(tag, english, turkish):
    return f"{tag}\n{english}\n\nTürkçe: {turkish}"

ALL_NOTES = [
    # (title_fragment, tag, english, turkish)
    ("NAS-Guided FEM-Anchored PINNs",
     "[INTRO]",
     "This thesis builds a FEM-anchored NAS-PINN framework with Multi-Step Window Prediction for transient heat transfer simulation. "
     "FEM is not removed — it provides trusted temperature anchors at selected time steps. "
     "NAS-PINN predicts the thermal field between anchors, reducing FEM calls while keeping error within engineering limits.",
     "Bu tez, geçici ısı transferi simülasyonu için SEY-çıpalı NAS-PINN çerçevesini Çok Adımlı Pencere Tahmini ile oluşturur. "
     "SEY kaldırılmaz; seçili zaman adımlarında güvenilir sıcaklık çıpaları sağlar. "
     "NAS-PINN, çıpalar arasındaki termal alanı tahmin ederek SEY çağrılarını azaltır ve hatayı mühendislik sınırları içinde tutar."),

    ("Presentation Structure",
     "[INTRO]",
     "The talk follows this order: physical problem, theory, method, benchmarks, metrics, results, discussion, and conclusion. "
     "Each section builds on the previous one, so understanding the method is essential before interpreting the results.",
     "Sunum şu sırayı izler: fiziksel problem, teori, yöntem, kıyaslama alanları, metrikler, sonuçlar, tartışma ve sonuç. "
     "Her bölüm bir önceki üzerine inşa edilir, bu nedenle sonuçları yorumlamadan önce yöntemi anlamak çok önemlidir."),

    ("Intro: What We Do",
     "[INTRO]",
     "The physical problem is A356 aluminium quenching: hot cast parts are cooled rapidly in water. "
     "The goal is to predict the temperature field over time — tracking a changing thermal field, not just a single final value. "
     "This makes the simulation both physically important and computationally demanding.",
     "Fiziksel problem A356 alüminyum su verilmesidir: sıcak döküm parçalar suda hızla soğutulur. "
     "Amaç, zaman içindeki sıcaklık alanını tahmin etmektir — tek bir son değer değil, değişen bir termal alanı takip etmek. "
     "Bu durum simülasyonu hem fiziksel açıdan önemli hem de hesaplamalı olarak zorlu kılar."),

    ("A356 aluminium",
     "[INTRO]",
     "A356 is a common aluminium-silicon casting alloy used in automotive and aerospace parts. "
     "During T6 heat treatment, the quenching step controls the cooling history and determines material strength. "
     "This makes repeated quenching simulation both industrially relevant and computationally intensive.",
     "A356, otomotiv ve havacılık parçalarında kullanılan yaygın bir alüminyum-silikon döküm alaşımıdır. "
     "T6 ısıl işleminde su verme adımı, soğuma geçmişini ve malzeme mukavemetini belirler. "
     "Bu, tekrarlanan su verme simülasyonunu hem endüstriyel açıdan önemli hem de hesaplamalı olarak yoğun kılar."),

    ("From the Reference Paper",
     "[INTRO]",
     "Mortensen et al. (2026) use FEM to simulate water quenching of an A356 automotive subframe. "
     "This thesis extends that work: instead of running full FEM at every time step, FEM is used as trusted anchors and NAS-PINN fills the gaps. "
     "The reference paper provides physical parameters and validates the simulation setup.",
     "Mortensen ve ark. (2026), A356 otomotiv alt çerçevesinin su verilmesini SEY ile simüle eder. "
     "Bu tez, söz konusu çalışmayı genişletir: her zaman adımında tam SEY çalıştırmak yerine SEY çıpa olarak kullanılır ve NAS-PINN boşlukları doldurur. "
     "Referans makale fiziksel parametreler sağlar ve simülasyon kurulumunu doğrular."),

    ("Reference Paper Inputs",
     "[INTRO]",
     "The reference work provides thermal conductivity, heat capacity, convection coefficient, and initial conditions. "
     "My thesis takes these parameters as given and asks: how much FEM work can we skip while keeping useful accuracy? "
     "The NAS-PINN architecture and MSWP method are the contributions that answer this question.",
     "Referans çalışma; termal iletkenlik, ısı kapasitesi, konveksiyon katsayısı ve başlangıç koşullarını sağlar. "
     "Bu tez bu parametreleri veri kabul ederek şunu sorar: Kullanışlı doğruluğu korurken ne kadar SEY işini atlayabiliriz? "
     "NAS-PINN mimarisi ve MSWP yöntemi bu soruyu yanıtlayan katkılardır."),

    ("Research Question",
     "[INTRO]",
     "The research question is: how much FEM computation can be skipped while keeping temperature-field error within engineering limits? "
     "Contributions: the FEM-anchored MSWP framework, NAS-guided architecture search, and benchmarking across 2D, 3D, and Thermal Fin. "
     "The answer connects skip factor k, MAE, and FEM saving rate.",
     "Araştırma sorusu: sıcaklık alanı hatasını mühendislik sınırları içinde tutarken ne kadar SEY hesabı atlanabilir? "
     "Katkılar; SEY-çıpalı MSWP çerçevesi, NAS güdümlü mimari arama ve 2D, 3D ile Termal Kanat üzerindeki kıyaslamadır. "
     "Yanıt, atlama faktörü k, Ortalama Mutlak Hata (OMH) ve SEY tasarrufu oranını birbirine bağlar."),

    ("Skip Factor",
     "[METHOD - SKIP FACTOR]",
     "The skip factor k defines how many FEM steps are replaced by PINN prediction. "
     "k=1: every step is anchored by FEM. k=5: four steps are predicted between two FEM anchors. "
     "Fixed-k uses a constant value; adaptive-k adjusts k based on the cooling rate; PINN-only removes all FEM anchors as a stress test.",
     "Atlama faktörü k, kaç SEY adımının PINN tahminiyle değiştirileceğini tanımlar. "
     "k=1: her adım SEY tarafından çıpalanır. k=5: iki SEY çıpası arasında dört adım tahmin edilir. "
     "Sabit-k sabit bir değer kullanır; adaptif-k soğuma hızına göre k'yı ayarlar; yalnızca-PINN ise stres testi olarak tüm SEY çıpalarını kaldırır."),

    ("Heat Transfer Model",
     "[THEORY]",
     "The governing equation is the transient heat equation: rho*cp*(dT/dt) = div(kT * grad(T)) + Q. "
     "FEM solves this on a mesh by discretising in space and time. "
     "PINN embeds this equation inside the loss function so the neural network is physically constrained.",
     "Yöneten denklem, geçici ısı denklemdir: rho*cp*(dT/dt) = div(kT * grad(T)) + Q. "
     "SEY bu denklemi uzay ve zamanda ayrıklaştırarak bir ağ üzerinde çözer. "
     "PINN ise bu denklemi kayıp fonksiyonuna gömer; böylece yapay sinir ağı fiziksel olarak kısıtlanır."),

    ("Finite Element Method",
     "[THEORY]",
     "FEM is the trusted numerical reference solver that discretises the domain and steps forward in time. "
     "PINN is a neural network trained with a physics loss: the heat equation residual must be small. "
     "NAS-PINN adds architecture search so the network structure is selected by testing candidates systematically.",
     "SEY, alanı ayrıklaştıran ve zaman içinde ilerleyen güvenilir sayısal referans çözücüdür. "
     "PINN, fizik kaybıyla eğitilen bir yapay sinir ağıdır: ısı denklemi artığı küçük olmalıdır. "
     "NAS-PINN mimari aramayı ekler; böylece ağ yapısı aday mimarileri sistematik olarak test ederek seçilir."),

    ("FEM Process",
     "[METHOD - FEM PROCESS AND MESH CONFIG]",
     "The FEM workflow has four steps: set up physics, build the mesh, solve the heat equation, save snapshots. "
     "Grid resolution: 131x61 (Rectangle, 7991 nodes), 121x121 (Circle, 11285 nodes), 81x81 (L-Shape, 4961 nodes). "
     "These FEM fields become the baseline for MAE, L2 error, and anchor snapshots. "
     "Runtime note: FEM provides trusted anchors; PINN reduces repeated solve calls.",
     "SEY iş akışının dört adımı vardır: fizik kurulumu, ağ oluşturma, ısı denklemini çözme, anlık görüntüleri kaydetme. "
     "Ağ çözünürlüğü: 131x61 (Dikdörtgen, 7991 düğüm), 121x121 (Daire, 11285 düğüm), 81x81 (L-Şekli, 4961 düğüm). "
     "Bu SEY alanları OMH, L2 hatası ve çıpa anlık görüntüleri için temel oluşturur. "
     "Çalışma süresi notu: SEY güvenilir çıpa sağlar; PINN tekrarlanan çözüm çağrılarını azaltır."),

    ("NAS-PINN: How It Works",
     "[METHOD - NAS-PINN SOURCE]",
     "NAS-PINN searches for the best network architecture instead of fixing it manually. "
     "Step 1: define search space (layers, neurons, activations, Fourier embedding). "
     "Step 2: run optimizer (Bayesian/TPE, NSGA-II, or NSGA-III) to test candidates. "
     "Step 3: train best PINN with Adam+L-BFGS. Step 4: transfer the winner across all geometries. "
     "Citation: Wang and Zhong (2024).",
     "NAS-PINN, en iyi ağ mimarisini elle belirlemek yerine arar. "
     "1. adım: arama uzayını tanımla (katmanlar, nöronlar, aktivasyonlar, Fourier gömme). "
     "2. adım: adayları test etmek için optimizer (Bayesian/TPE, NSGA-II veya NSGA-III) çalıştır. "
     "3. adım: en iyi PINN'i Adam+L-BFGS ile eğit. 4. adım: kazananı tüm geometrilere aktar. "
     "Kaynak: Wang ve Zhong (2024)."),

    ("NAS-PINN Process",
     "[METHOD - NAS-PINN PROCESS]",
     "The NAS-PINN process has four steps per prediction window. "
     "Step 1 (Anchor): FEM provides the trusted state at the window start. "
     "Step 2 (Prediction): NAS-PINN predicts T(x,t) inside the window. "
     "Step 3 (Loss): physics residual and endpoint supervision keep prediction physically consistent. "
     "Step 4 (NAS): architecture search selects the best network design for the geometry.",
     "NAS-PINN süreci, her tahmin penceresi için dört adıma sahiptir. "
     "1. adım (Çıpa): SEY pencere başında güvenilir durumu sağlar. "
     "2. adım (Tahmin): NAS-PINN pencere içinde T(x,t) sıcaklık alanını tahmin eder. "
     "3. adım (Kayıp): fizik artığı ve uç nokta gözetimi tahmini fiziksel olarak tutarlı tutar. "
     "4. adım (NAS): mimari arama geometri için en iyi ağ tasarımını seçer."),

    ("Architecture Search",
     "[METHOD - TRAINING]",
     "Architecture search is performed first, then the selected model is trained. "
     "Adam optimizer handles early broad convergence; L-BFGS refines the solution in each window. "
     "This two-phase training protocol is applied consistently across all benchmark geometries.",
     "Önce mimari arama yapılır, ardından seçilen model eğitilir. "
     "Adam optimizer erken geniş yakınsamayı sağlar; L-BFGS her penceredeki çözümü iyileştirir. "
     "Bu iki aşamalı eğitim protokolü tüm kıyaslama geometrilerinde tutarlı biçimde uygulanır."),

    ("Multi-Step Window Prediction",
     "[METHOD - MSWP]",
     "MSWP divides the time simulation into overlapping windows. "
     "FEM (Finite Element Method) anchors provide the starting state for each window. "
     "NAS-PINN fills the k skipped steps. The next FEM anchor corrects any accumulated drift.",
     "MSWP (Çok Adımlı Pencere Tahmini), zaman simülasyonunu örtüşen pencerelere böler. "
     "SEY (Sonlu Elemanlar Yöntemi) çıpaları her pencere için başlangıç durumunu sağlar. "
     "NAS-PINN k atlanan adımı doldurur. Bir sonraki SEY çıpası birikmiş herhangi bir kayışı düzeltir."),

    ("Benchmark Domains",
     "[BENCHMARKS]",
     "Four benchmark levels test the method at increasing difficulty. "
     "Level 1 (2D NAS transfer): Square, Circle, L-Shape, Flower with adaptive-k. "
     "Level 2 (canonical 2D): Rectangle, Circle, L-Shape with fixed k=1 to 5. "
     "Level 3 (3D): Rectangular, Cylinder, Stacked Cubes, L-Shape Prism. "
     "Level 4 (Thermal Fin): most demanding — thin-fin gradients and Fourier features.",
     "Dört kıyaslama seviyesi yöntemi artan güçlükte test eder. "
     "Seviye 1 (2D NAS transferi): adaptif-k ile Kare, Daire, L-Şekli, Çiçek. "
     "Seviye 2 (kanonik 2D): sabit k=1-5 ile Dikdörtgen, Daire, L-Şekli. "
     "Seviye 3 (3D): Dikdörtgen, Silindir, Yığılmış Küpler, L-Şekli Prizma. "
     "Seviye 4 (Termal Kanat): en zorlu — ince kanat gradyanları ve Fourier özellikleri."),

    ("Evaluation Metrics",
     "[BENCHMARKS - METRICS]",
     "Three metrics are used. MAE (Mean Absolute Error) measures average temperature prediction error in Celsius. "
     "Relative L2 measures full-field error as a normalised norm showing whether the field shape is correct. "
     "FEM Anchor Reduction shows what percentage of FEM calls are replaced by PINN predictions.",
     "Üç metrik kullanılır. OMH (Ortalama Mutlak Hata), Celsius cinsinden ortalama sıcaklık tahmin hatasını ölçer. "
     "Göreceli L2, tam alan hatasını normalleştirilmiş norm olarak ölçer; alan şeklinin doğruluğunu gösterir. "
     "SEY Çıpası Azaltma, SEY çağrılarının yüzde kaçının PINN tahminleriyle değiştirildiğini gösterir."),

    ("2D Adaptive-k: MAE and L2 Summary",
     "[RESULTS - 2D NAS TRANSFER]",
     "The 2D NAS-transfer adaptive-k results show all three architectures achieve 83-85% FEM saving. "
     "MAE remains within 10°C for all tested geometries. "
     "Adaptive-k successfully controls prediction window size during the steep early transient where gradients are largest.",
     "2D NAS-transfer adaptif-k sonuçları, üç mimarinin yüzde 83-85 SEY tasarrufu sağladığını gösterir. "
     "OMH, test edilen tüm geometriler için 10°C içinde kalır. "
     "Adaptif-k, gradyanların en büyük olduğu dik erken geçici dönemde tahmin penceresi boyutunu başarıyla kontrol eder."),

    ("2D Adaptive-k: Figure",
     "[RESULTS - 2D CURVES]",
     "FEM and PINN temperature curves are close throughout the simulation. "
     "Error is highest in early windows because steep temperature gradients make prediction harder. "
     "This explains why adaptive-k reduces the window size early and allows larger k later.",
     "SEY ve PINN sıcaklık eğrileri simülasyon boyunca birbirine yakındır. "
     "Hata, dik sıcaklık gradyanları tahmini zorlaştırdığı için ilk pencerelerde en yüksek değerini alır. "
     "Bu durum, adaptif-k'nın erken dönemde pencere boyutunu neden küçülttüğünü ve sonra daha büyük k'ya neden izin verdiğini açıklar."),

    ("Canonical 2D: Setup",
     "[RESULTS - CANONICAL 2D SETUP]",
     "The canonical 2D study uses fixed-k values from 1 to 5 on Rectangle, Circle, and L-Shape. "
     "The adaptive-k schedule shows window size changes: early windows use smaller k; later windows use larger k. "
     "Understanding this schedule is essential before reading the MAE and L2 tables.",
     "Kanonik 2D çalışması, Dikdörtgen, Daire ve L-Şekli üzerinde 1'den 5'e kadar sabit-k değerleri kullanır. "
     "Adaptif-k takvimi pencere boyutu değişimini gösterir: erken pencereler daha küçük k, sonraki pencereler daha büyük k kullanır. "
     "Bu takvimi anlamak, OMH ve L2 tablolarını okumadan önce çok önemlidir."),

    ("Canonical 2D: MAE and L2",
     "[RESULTS - CANONICAL 2D TABLES]",
     "Canonical 2D MAE and L2 tables compare all three architectures across k values and geometries. "
     "Bayesian/TPE stays strongest as k increases. Circle and L-Shape show higher errors than Rectangle at large k. "
     "This confirms that geometry and skip factor together determine prediction difficulty.",
     "Kanonik 2D OMH ve L2 tabloları, üç mimariyi k değerleri ve geometriler arasında karşılaştırır. "
     "Bayesian/TPE, k arttıkça en güçlü kalmaya devam eder. Daire ve L-Şekli, büyük k'da Dikdörtgen'den yüksek hatalar gösterir. "
     "Bu, geometri ve atlama faktörünün birlikte tahmin güçlüğünü belirlediğini doğrular."),

    ("Canonical 2D: MAE Variability",
     "[RESULTS - CANONICAL 2D TRENDS]",
     "Variability plots show whether MAE is consistent across random seeds. "
     "L2 curves show how field error changes over time windows. "
     "Stable MAE and low L2 both confirm the method produces reliable results across different initialisations.",
     "Değişkenlik grafikleri OMH'nin tohumlar arasında tutarlı olup olmadığını gösterir. "
     "L2 eğrileri alan hatasının zaman pencereleri boyunca nasıl değiştiğini gösterir. "
     "Kararlı OMH ve düşük L2, yöntemin farklı başlatmalar üzerinde güvenilir sonuçlar ürettiğini doğrular."),

    ("Canonical 2D Heat Maps",
     "[RESULTS - 2D HEAT MAPS]",
     "The 2D heat maps show where error appears in the predicted temperature field. "
     "Small average MAE can still hide local errors near boundaries or corners. "
     "Maps confirm the method handles complex geometries but curved boundaries concentrate errors.",
     "2D ısı haritaları, tahmin edilen sıcaklık alanında hatanın nerede göründüğünü gösterir. "
     "Küçük ortalama OMH hâlâ sınırlara yakın yerel hataları gizleyebilir. "
     "Haritalar, yöntemin karmaşık geometrileri ele aldığını ancak eğri sınırların hataları yoğunlaştırdığını doğrular."),

    ("3D Benchmark: Physical",
     "[METHOD - 3D PARAMS]",
     "These parameters define the physical setup for all 3D and Thermal Fin benchmarks. "
     "kT=160 W/(m·K) and ρcp=2.4×10⁶ J/(m³·K) are constants from Mortensen et al. (2026). "
     "Anchor spacing Δtₐ=1.5 s and FEM step Δtᵢₙₜ=0.5 s define the temporal resolution.",
     "Bu parametreler, tüm 3D ve Termal Kanat kıyaslamaları için fiziksel kurulumu tanımlar. "
     "kT=160 W/(m·K) ve ρcp=2,4×10⁶ J/(m³·K), Mortensen ve ark.'dan (2026) alınan sabitlerdir. "
     "Çıpa aralığı Δtₐ=1,5 s ve SEY adımı Δtᵢₙₜ=0,5 s simülasyonun zamansal çözünürlüğünü belirler."),

    ("3D Results: MAE and L2",
     "[RESULTS - 3D TABLES]",
     "The 3D results show MAE and L2 across four geometries: Rectangular, Cylinder, Stacked Cubes, L-Shape Prism. "
     "Bayesian/TPE reaches k=5 on all geometries within the 15°C MAE band. "
     "3D prediction is harder than 2D because the field is a volume and curved boundaries are more complex.",
     "3D sonuçlar, Dikdörtgen, Silindir, Yığılmış Küpler ve L-Şekli Prizma üzerindeki OMH ve L2'yi gösterir. "
     "Bayesian/TPE, tüm geometrilerde k=5'e ulaşırken 15°C OMH bandı içinde kalır. "
     "3D tahmini, alan bir hacim ve eğri sınırlar daha karmaşık olduğundan 2D'den daha zordur."),

    ("3D Results: MAE Variability",
     "[RESULTS - 3D TRENDS]",
     "The 3D variability and L2 curves show that some geometries are more sensitive than others. "
     "NSGA-II and NSGA-III show more variability at large k on complex 3D shapes. "
     "Bayesian/TPE remains the most stable choice for 3D domains.",
     "3D değişkenlik ve L2 eğrileri, bazı geometrilerin diğerlerinden daha hassas olduğunu gösterir. "
     "NSGA-II ve NSGA-III, karmaşık 3D şekillerde büyük k değerlerinde daha fazla değişkenlik gösterir. "
     "Bayesian/TPE, 3D alanlar için en kararlı seçim olmaya devam eder."),

    ("3D Heat Map",
     "[RESULTS - 3D HEAT MAP]",
     "The Cylinder domain with Bayesian/TPE gives the best 3D result. "
     "The volumetric map shows FEM-anchored PINN field closely matches the FEM reference. "
     "This confirms FEM anchoring prevents drift accumulation in 3D.",
     "Bayesian/TPE ile Silindir alanı en iyi 3D sonucu verir. "
     "Hacimsel harita, SEY-çıpalı PINN alanının SEY referansıyla yakından eşleştiğini gösterir. "
     "Bu, SEY çıpalamanın 3D'de kayış birikimini önlediğini doğrular."),

    ("Canonical 3D Heat Maps",
     "[RESULTS - 3D HEAT MAP]",
     "The Cylinder domain with Bayesian/TPE gives the best 3D result. "
     "The volumetric map shows FEM-anchored PINN field closely matches the FEM reference. "
     "This confirms FEM anchoring prevents drift accumulation in 3D.",
     "Bayesian/TPE ile Silindir alanı en iyi 3D sonucu verir. "
     "Hacimsel harita, SEY-çıpalı PINN alanının SEY referansıyla yakından eşleştiğini gösterir. "
     "Bu, SEY çıpalamanın 3D'de kayış birikimini önlediğini doğrular."),

    ("Thermal Fin: MAE and L2",
     "[RESULTS - THERMAL FIN TABLES]",
     "Thermal Fin is the most demanding benchmark due to thin-fin geometry and local gradients. "
     "Bayesian/TPE with Fourier features stays within 15°C MAE at all k values. "
     "NSGA-II and NSGA-III exceed the band at higher k, showing architecture choice matters most here.",
     "Termal Kanat, ince kanat geometrisi ve yerel gradyanlar nedeniyle en zorlu kıyaslamadır. "
     "Fourier özellikleriyle Bayesian/TPE, tüm k değerlerinde 15°C OMH bandı içinde kalır. "
     "NSGA-II ve NSGA-III daha yüksek k'da bandı aşar; bu mimari seçiminin burada en kritik faktör olduğunu gösterir."),

    ("Thermal Fin: MAE Variability",
     "[RESULTS - THERMAL FIN TRENDS]",
     "Variability and L2 curves show higher k is harder for all architectures. "
     "Fourier features make Bayesian/TPE more stable: its L2 curve stays low even at k=5. "
     "NSGA-II and NSGA-III show more spread, confirming Fourier embedding is important for Thermal Fin.",
     "Değişkenlik ve L2 eğrileri, daha yüksek k'nın tüm mimariler için daha zor olduğunu gösterir. "
     "Fourier özellikleri Bayesian/TPE'yi daha kararlı kılar: L2 eğrisi k=5'te bile düşük kalır. "
     "NSGA-II ve NSGA-III daha fazla yayılma gösterir; Fourier gömmenin Termal Kanat için önemli olduğunu doğrular."),

    ("Thermal Fin: Fixed-k Heat Maps (2D)",
     "[RESULTS - THERMAL FIN 2D HEAT MAPS]",
     "The 2D cross-section heat maps show FEM reference, PINN prediction, and absolute error side by side. "
     "FEM anchoring keeps the error pattern local and small near boundaries. "
     "This confirms the method handles the complex thin-fin geometry at the 2D cross-section level.",
     "2D kesit ısı haritaları, SEY referansını, PINN tahminini ve mutlak hatayı yan yana gösterir. "
     "SEY çıpalama, hata örüntüsünü sınırlara yakın yerel ve küçük tutar. "
     "Bu, yöntemin 2D kesit seviyesinde karmaşık ince kanat geometrisini ele alabileceğini doğrular."),

    ("Anchored vs PINN-only",
     "[RESULTS - ADAPTIVE VS PINN-ONLY]",
     "This comparison shows the role of FEM anchors directly. "
     "FEM-anchored: errors are small and local — periodic correction keeps the field stable. "
     "PINN-only: errors accumulate and spread because each imperfect window starts the next. "
     "PINN-only errors are 4-18x larger, proving FEM anchoring is essential.",
     "Bu karşılaştırma, SEY çıpalarının rolünü doğrudan gösterir. "
     "SEY-çıpalı: hatalar küçük ve yerel — periyodik düzeltme alanı kararlı tutar. "
     "Yalnızca-PINN: her kusurlu pencere bir sonrakini başlattığından hatalar birikir ve yayılır. "
     "Yalnızca-PINN hataları 4-18 kat daha büyük; SEY çıpalamanın zorunlu olduğunu kanıtlar."),

    ("Adaptive vs PINN-only Heat Maps",
     "[RESULTS - ADAPTIVE VS PINN-ONLY]",
     "This comparison shows the role of FEM anchors directly. "
     "FEM-anchored: errors are small and local — periodic correction keeps the field stable. "
     "PINN-only: errors accumulate and spread because each imperfect window starts the next. "
     "PINN-only errors are 4-18x larger, proving FEM anchoring is essential.",
     "Bu karşılaştırma, SEY çıpalarının rolünü doğrudan gösterir. "
     "SEY-çıpalı: hatalar küçük ve yerel — periyodik düzeltme alanı kararlı tutar. "
     "Yalnızca-PINN: her kusurlu pencere bir sonrakini başlattığından hatalar birikir ve yayılır. "
     "Yalnızca-PINN hataları 4-18 kat daha büyük; SEY çıpalamanın zorunlu olduğunu kanıtlar."),

    ("Thermal Fin: Fixed-k Volumetric",
     "[RESULTS - THERMAL FIN 3D VIEWS]",
     "The 3D volumetric views confirm sharp fin regions are the hardest to predict. "
     "FEM anchoring keeps the volumetric field stable across windows. "
     "Bayesian/TPE with Fourier shows the best volumetric accuracy.",
     "3D hacimsel görünümler, keskin kanat bölgelerinin tahmin edilmesinin en zor olduğunu doğrular. "
     "SEY çıpalama hacimsel alanı pencereler boyunca kararlı tutar. "
     "Fourier ile Bayesian/TPE en iyi hacimsel doğruluğu gösterir."),

    ("Adaptive vs PINN-only Volumetric",
     "[RESULTS - ADAPTIVE VS PINN-ONLY 3D]",
     "The 3D comparison confirms the 2D finding: FEM-anchored prediction stays stable; PINN-only shows drift. "
     "Periodic FEM correction is essential in 3D because volume errors accumulate more rapidly. "
     "This result supports using FEM anchors at regular intervals even for large k values.",
     "3D karşılaştırma, 2D bulguyu doğrular: SEY-çıpalı tahmin kararlı kalır; yalnızca-PINN kayış gösterir. "
     "Periyodik SEY düzeltmesi 3D'de zorunludur çünkü hacim hataları daha hızlı birikir. "
     "Bu sonuç, büyük k değerlerinde bile düzenli aralıklarla SEY çıpası kullanımını destekler."),

    ("Discussion",
     "[DISCUSSION]",
     "2D: 83-85% FEM saving; adaptive-k keeps errors low during steep transients. "
     "3D: Bayesian/TPE reaches k=5 within 15°C on all geometries; geometry sensitivity grows with k. "
     "Thermal Fin: Fourier features stabilise Bayesian/TPE; FEM anchors prevent drift. "
     "PINN-only errors are 4-18x larger — periodic correction is required. Ablation: Fourier + endpoint supervision are essential choices.",
     "2D: yüzde 83-85 SEY tasarrufu; adaptif-k dik geçici dönemlerde hataları düşük tutar. "
     "3D: Bayesian/TPE tüm geometrilerde 15°C içinde k=5'e ulaşır; geometri hassasiyeti k ile artar. "
     "Termal Kanat: Fourier özellikleri Bayesian/TPE'yi kararlılaştırır; SEY çıpaları kayışı önler. "
     "Yalnızca-PINN hataları 4-18 kat daha büyük — periyodik düzeltme gereklidir. Ablasyon: Fourier + uç nokta gözetimi zorunlu seçimlerdir."),

    ("Conclusion",
     "[CONCLUSION]",
     "Built a FEM-anchored NAS-PINN framework with MSWP; validated across 2D, 3D, and Thermal Fin. "
     "Key result: 65-85% FEM saving; MAE within 10°C (2D) and 15°C (3D/Fin); Bayesian/TPE+Fourier is most robust. "
     "Limitations: constant material model; future: temperature-dependent properties, Mortensen mesh, StaMiSim validation.",
     "SEY-çıpalı NAS-PINN çerçevesi MSWP ile oluşturuldu; 2D, 3D ve Termal Kanat üzerinde doğrulandı. "
     "Temel sonuç: yüzde 65-85 SEY tasarrufu; OMH 10°C (2D) ve 15°C (3D/Kanat) içinde; Bayesian/TPE+Fourier en sağlam. "
     "Sınırlamalar: sabit malzeme modeli; gelecek: sıcaklığa bağlı özellikler, Mortensen ağı, StaMiSim doğrulaması."),

    ("Q&A",
     "[Q AND A]",
     "If asked about accuracy: MAE within 10°C for 2D and 15°C for 3D/Thermal Fin. "
     "If asked about speed: FEM anchor reduction of 65-85%. "
     "If asked about reliability: FEM anchoring, heat-map evidence, and PINN-only comparison show the method is reliable.",
     "Doğruluk sorulursa: OMH 2D için 10°C, 3D/Termal Kanat için 15°C içinde. "
     "Hız sorulursa: yüzde 65-85 SEY çıpası azaltması. "
     "Güvenilirlik sorulursa: SEY çıpalama, ısı haritası kanıtı ve yalnızca-PINN karşılaştırması yöntemin güvenilir olduğunu gösterir."),

    ("Preparation",
     "[Q AND A]",
     "If asked about accuracy: MAE within 10°C for 2D and 15°C for 3D/Thermal Fin. "
     "If asked about speed: FEM anchor reduction of 65-85%. "
     "If asked about reliability: FEM anchoring, heat-map evidence, and PINN-only comparison show the method is reliable.",
     "Doğruluk sorulursa: OMH 2D için 10°C, 3D/Termal Kanat için 15°C içinde. "
     "Hız sorulursa: yüzde 65-85 SEY çıpası azaltması. "
     "Güvenilirlik sorulursa: SEY çıpalama, ısı haritası kanıtı ve yalnızca-PINN karşılaştırması yöntemin güvenilir olduğunu gösterir."),

    ("Thank You",
     "[CLOSING]",
     "FEM remains the trusted reference, and NAS-PINN reduces repeated FEM work by predicting selected windows with controlled error. "
     "Thank you for your attention. I am happy to answer questions.",
     "SEY güvenilir referans olmaya devam eder ve NAS-PINN, seçili pencereleri kontrollü hatayla tahmin ederek tekrarlanan SEY işini azaltır. "
     "Dikkatiniz için teşekkür ederim. Soruları yanıtlamaktan memnuniyet duyarım."),
]

notes_set_count = 0
for i, slide in enumerate(prs.slides):
    title = slide_title_text(slide)
    matched = False
    for (frag, tag, eng, tr) in ALL_NOTES:
        if frag.lower() in title.lower():
            set_notes(slide, format_note(tag, eng, tr))
            print(f"  Notes set for slide {i+1}: {title[:55]}")
            notes_set_count += 1
            matched = True
            break
    if not matched and title:
        set_notes(slide,
            f"[SLIDE {i+1}]\n"
            f"This slide covers: {title}.\n\n"
            f"Türkçe: Bu slayt şunu kapsar: {title}.")

print(f"\nNotes set for {notes_set_count} slides with full EN+TR content.")

# ──────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────
print(f"\nSaving to {DST}...")
prs.save(DST)
print(f"Saved. Final slide count: {len(prs.slides)}")

print("\nFinal slide list:")
for i, slide in enumerate(prs.slides):
    print(f"  {i+1:2d}: {slide_title_text(slide)[:70]}")

# ──────────────────────────────────────────────
# POST-PROCESSING: Fix notes for slides that had broken notes parts
# Some slides in the source PPTX had notes files not registered in Content_Types,
# causing python-pptx to fail to create proper NotesSlidePart objects.
# We fix this by directly editing the saved PPTX zip file.
# ──────────────────────────────────────────────
import zipfile
import shutil
import io
import re

print("\n[Post-processing] Fixing notes content types and text for broken slides...")

def escape_xml(text):
    return (text
        .replace('&', '&amp;')
        .replace('<', '&lt;')
        .replace('>', '&gt;')
        .replace('"', '&quot;')
        .replace("'", '&apos;')
    )

def build_notes_xml_text(text):
    """Build notes XML body text element content."""
    safe_text = escape_xml(text)
    return f"""<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
         xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
         xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm>
      </p:grpSpPr>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Notes Placeholder 1"/>
          <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
          <p:nvPr><p:ph type="body" idx="1"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
        <p:txBody>
          <a:bodyPr/>
          <a:lstStyle/>
          <a:p><a:r><a:t>{safe_text}</a:t></a:r></a:p>
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
</p:notes>"""

# Step 1: Find which slides have notes files missing from Content_Types
with zipfile.ZipFile(DST, 'r') as z:
    ct_xml = z.read('[Content_Types].xml').decode('utf-8')
    all_files = z.namelist()

    # Map: notes_file_shortname -> True/False (in content types)
    notes_files_in_zip = set(
        n.replace('ppt/notesSlides/', '')
        for n in all_files
        if re.match(r'ppt/notesSlides/notesSlide.*\.xml$', n) and '_rels' not in n
    )
    notes_files_in_ct = set(
        m.group(1)
        for m in re.finditer(r'PartName="/ppt/notesSlides/([^"]+\.xml)"', ct_xml)
    )
    missing_from_ct = notes_files_in_zip - notes_files_in_ct
    print(f"  Notes files missing from Content_Types: {sorted(missing_from_ct)}")

    # Map: slide_num -> notes_file it points to
    slide_to_notes = {}
    for i in range(1, 40):
        rels_file = f'ppt/slides/_rels/slide{i}.xml.rels'
        if rels_file in all_files:
            content = z.read(rels_file).decode('utf-8')
            m = re.search(r'notesSlides/([^"]+\.xml)', content)
            if m:
                slide_to_notes[i] = m.group(1)

# Find which slides are pointing to missing notes files
prs_check = Presentation(DST)
slides_titles = {i+1: slide_title_text(prs_check.slides[i]) for i in range(len(prs_check.slides))}

# Build the note text for each slide that needs fixing
notes_replacement = {}  # notes_filename -> new notes XML content
for slide_num, notes_file in slide_to_notes.items():
    if notes_file in missing_from_ct:
        title = slides_titles.get(slide_num, f"Slide {slide_num}")
        note_text = f"[SLIDE {slide_num}]\nThis slide covers: {title}.\n\nTürkçe: Bu slayt şunu kapsar: {title}."
        for (frag, tag, eng, tr) in ALL_NOTES:
            if frag.lower() in title.lower():
                note_text = format_note(tag, eng, tr)
                break
        notes_replacement[notes_file] = build_notes_xml_text(note_text)
        print(f"  Will update notes for slide {slide_num}: {notes_file} -> {title[:40]}")

# Now rewrite the PPTX to fix everything in one pass
DST_TMP = DST + ".tmp"
with zipfile.ZipFile(DST, 'r') as zin:
    with zipfile.ZipFile(DST_TMP, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)

            # Fix Content_Types.xml: add missing notes file entries
            if item.filename == '[Content_Types].xml':
                xml_str = data.decode('utf-8')
                for nf in sorted(missing_from_ct):
                    ct_entry = (f'<Override PartName="/ppt/notesSlides/{nf}" '
                                f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>')
                    if f'notesSlides/{nf}' not in xml_str:
                        xml_str = xml_str.replace('</Types>', ct_entry + '\n</Types>')
                data = xml_str.encode('utf-8')
                print("  Updated Content_Types.xml")

            # Fix the actual notes XML files to contain correct text
            for nf, new_xml in notes_replacement.items():
                if item.filename == f'ppt/notesSlides/{nf}':
                    data = new_xml.encode('utf-8')
                    print(f"  Replaced content of {nf}")
                    break

            zout.writestr(item, data)

os.replace(DST_TMP, DST)
print(f"  Post-processing complete.")

# ──────────────────────────────────────────────
# VERIFY
# ──────────────────────────────────────────────
prs2 = Presentation(DST)
print(f"\nVerification: file opens correctly with {len(prs2.slides)} slides.")

# Check notes
print("Notes check:")
no_notes_final = []
for i, slide in enumerate(prs2.slides):
    try:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        has_tr = 'Türkçe:' in notes
        if not notes:
            no_notes_final.append(i+1)
            print(f"  Slide {i+1}: EMPTY NOTES")
        else:
            print(f"  Slide {i+1}: OK (has_turkish={has_tr})")
    except Exception as e:
        no_notes_final.append(i+1)
        print(f"  Slide {i+1}: ERROR - {e}")

if no_notes_final:
    print(f"\nWARNING: {len(no_notes_final)} slides still missing notes: {no_notes_final}")
else:
    print("\nAll slides have notes.")

import os
size = os.path.getsize(DST)
print(f"\nFinal file size: {size:,} bytes ({size/1024/1024:.1f} MB)")
print("Done.")
